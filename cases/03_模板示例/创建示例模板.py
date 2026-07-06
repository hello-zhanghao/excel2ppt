"""
创建 PPT 模板示例
生成「数据报告模板.pptx」到同目录，演示各种占位符用法

使用流程：
  1. 先运行透视分析生成结果：
     python app\main.py pivot -c cases\03_模板示例\项目配置.xlsx
  2. 用模板填充数据：
     python app\main.py template cases\03_模板示例\数据报告模板.pptx --pivot 透视结果.xlsx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(SCRIPT_DIR, "数据报告模板.pptx")


def _set_run(run, text, size=18, bold=False, color=None):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _set_run(p.runs[0] if p.runs else p.add_run(), text, size, bold, color)
    return box


def _add_chart(slide, chart_type, left, top, width, height, title, categories, series_data):
    """添加图表，标题里带占位符"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)
    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart
    try:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    except Exception:
        pass
    return chart


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 颜色
    NAVY = RGBColor(0x18, 0x2B, 0x49)
    BLUE = RGBColor(0x2E, 0x75, 0xB6)
    GRAY = RGBColor(0x59, 0x59, 0x59)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    # ========== 第1页：标题页 ==========
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.2),
                 "5G 基站数据分析报告", 40, True, NAVY)
    _add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.6),
                 "总基站数：{{按地区汇总.基站数.sum}}  |  总用户数：{{按频段汇总.总用户数.sum}}",
                 20, False, GRAY)
    _add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 "覆盖地区数：{{按地区汇总.行数}}  |  频段数：{{按频段汇总.行数}}",
                 16, False, GRAY)
    # 备注配置
    try:
        slide.notes_slide.notes_text_frame.text = "# 标题页\n数据源=透视结果.xlsx\n"
    except Exception:
        pass

    # ========== 第2页：地区基站汇总 ==========
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "各地区基站分布", 32, True, NAVY)

    # 左侧：文本占位符
    _add_textbox(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                 "华东基站数：{{按地区汇总.基站数.华东}}", 18, False, GRAY)
    _add_textbox(slide, Inches(0.5), Inches(2.1), Inches(5.5), Inches(0.5),
                 "华北基站数：{{按地区汇总.基站数.华北}}", 18, False, GRAY)
    _add_textbox(slide, Inches(0.5), Inches(2.7), Inches(5.5), Inches(0.5),
                 "华南基站数：{{按地区汇总.基站数.华南}}", 18, False, GRAY)
    _add_textbox(slide, Inches(0.5), Inches(3.3), Inches(5.5), Inches(0.5),
                 "华西基站数：{{按地区汇总.基站数.华西}}", 18, False, GRAY)
    _add_textbox(slide, Inches(0.5), Inches(4.0), Inches(5.5), Inches(0.5),
                 "平均基站数：{{按地区汇总.基站数.avg}}", 18, True, BLUE)

    # 右侧：柱状图占位符
    _add_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
               Inches(6.5), Inches(1.5), Inches(6.3), Inches(5),
               "{{图表:按地区汇总}}",
               ["A", "B", "C"], {"占位": (1, 2, 3)})

    # ========== 第3页：频段分析 ==========
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "频段维度统计", 32, True, NAVY)

    # 文本占位符
    _add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
                 "n78 平均用户数：{{按频段汇总.平均用户数.avg}}  |  "
                 "n41 平均用户数：{{按频段汇总.平均用户数.n41}}  |  "
                 "n28 平均用户数：{{按频段汇总.平均用户数.n28}}",
                 16, False, GRAY)

    # 柱状图：基站数
    _add_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
               Inches(0.5), Inches(2.2), Inches(6), Inches(4.5),
               "{{图表:按频段汇总}}",
               ["A", "B"], {"占位": (1, 2)})

    # 饼图：用户占比
    _add_chart(slide, XL_CHART_TYPE.PIE,
               Inches(7), Inches(2.2), Inches(5.8), Inches(4.5),
               "{{图表:地区用户占比}}",
               ["A", "B", "C"], {"占位": (1, 2, 3)})

    # ========== 第4页：速率多维统计 ==========
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "下行速率多维统计", 32, True, NAVY)

    # 文本占位符
    _add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
                 "峰值：{{速率多维统计.峰值.max}} Mbps  |  "
                 "谷值：{{速率多维统计.谷值.min}} Mbps  |  "
                 "均值：{{速率多维统计.均值.avg}} Mbps",
                 18, True, BLUE)

    # 图表占位符
    _add_chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED,
               Inches(0.5), Inches(2.2), Inches(12), Inches(4.8),
               "{{图表:速率多维统计}}",
               ["A", "B", "C"], {"占位": (1, 2, 3)})

    # ========== 第5页：季度汇总 ==========
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                 "季度指标趋势", 32, True, NAVY)

    _add_textbox(slide, Inches(0.5), Inches(1.4), Inches(12), Inches(0.5),
                 "Q1 总用户数：{{季度汇总.总用户数(千人).Q1}}  |  "
                 "Q2 总用户数：{{季度汇总.总用户数(千人).Q2}}",
                 18, False, GRAY)

    # 折线图占位符
    _add_chart(slide, XL_CHART_TYPE.LINE,
               Inches(0.5), Inches(2.2), Inches(12), Inches(4.8),
               "{{图表:季度汇总}}",
               ["A", "B", "C", "D"], {"占位": (1, 2, 3, 4)})

    # ========== 第6页：结尾页 ==========
    slide = prs.slides.add_slide(blank)
    _add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.8),
                 "— 报告结束 —", 32, True, NAVY)
    _add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(0.5),
                 "本报告由 Excel 统一分析工具 自动生成  |  总分析维度：16  |  数据行数：{{按地区汇总.行数}}",
                 16, False, GRAY)

    # 保存
    prs.save(OUTPUT)
    print(f"[OK] 模板已生成: {OUTPUT}")
    print(f"     共 {len(prs.slides)} 页")
    print(f"\n使用方法：")
    print(f"  1. 先运行透视分析：python app\\main.py pivot -c cases\\03_模板示例\\项目配置.xlsx")
    print(f"  2. 再填充模板：python app\\main.py template \"{OUTPUT}\" --pivot 透视结果.xlsx")


if __name__ == "__main__":
    main()
