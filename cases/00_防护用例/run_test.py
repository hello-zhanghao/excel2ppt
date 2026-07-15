"""
防护用例一键测试脚本
用法: python run_test.py

执行流程:
  1. 重建测试数据
  2. 运行透视分析（pivot）
  3. 运行PPT生成（ppt），引用透视结果
  4. 读取并展示输出结果，供人工检查
  5. PPT转图片 + 生成HTML报告（手机可查看）
"""
import os
import sys
import subprocess
import time
import openpyxl
import glob
import base64

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
# SCRIPT_DIR = .../cases/00_防护用例
# PROJECT_DIR = .../excel2ppt  (项目根目录，向上退两级)
PROJECT_DIR = os.path.dirname(os.path.dirname(SCRIPT_DIR))
PYTHON = sys.executable

# 颜色输出
GREEN = "\033[92m"
RED = "\033[91m"
YELLOW = "\033[93m"
CYAN = "\033[96m"
RESET = "\033[0m"
BOLD = "\033[1m"


def print_header(title):
    print(f"\n{CYAN}{'='*60}")
    print(f"  {BOLD}{title}")
    print(f"{'='*60}{RESET}")


def run_cmd(cmd, cwd=PROJECT_DIR):
    """运行命令，返回是否成功"""
    print(f"  {YELLOW}$ {cmd}{RESET}")
    result = subprocess.run(cmd, shell=True, cwd=cwd, capture_output=True, text=True, encoding="utf-8")
    if result.stdout.strip():
        print(result.stdout.rstrip())
    if result.returncode != 0:
        print(f"  {RED}[FAIL] 退出码={result.returncode}{RESET}")
        if result.stderr.strip():
            print(f"  {RED}stderr: {result.stderr.rstrip()}{RESET}")
        return False
    return True


def find_latest_output():
    """查找最新的输出目录"""
    output_dirs = glob.glob(os.path.join(SCRIPT_DIR, "output_*"))
    if not output_dirs:
        return None
    return max(output_dirs, key=os.path.getmtime)


def inspect_excel(excel_path):
    """读取并展示透视分析Excel输出"""
    print_header("Excel 输出内容检查")
    if not excel_path or not os.path.exists(excel_path):
        print(f"  {RED}Excel 文件不存在{RESET}")
        return

    print(f"  文件: {os.path.basename(excel_path)}\n")
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    print(f"  Sheet 列表: {wb.sheetnames}\n")

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        print(f"  {BOLD}--- {sheet_name} ({len(rows)-1} 行数据) ---{RESET}")

        # 找到真正的表头行（跳过区块标题行）
        header_row_idx = 0
        for i, row in enumerate(rows[:3]):
            non_none = [c for c in row if c is not None]
            if len(non_none) >= 2 and any(isinstance(c, str) and c.strip() for c in row):
                header_row_idx = i
                break

        headers = [str(c) if c is not None else "" for c in rows[header_row_idx]]
        # 只保留非空表头
        col_indices = [i for i, h in enumerate(headers) if h.strip()]
        headers = [headers[i] for i in col_indices]

        # 打印表头
        print(f"    列名: {' | '.join(headers)}")

        # 打印数据行
        for row in rows[header_row_idx + 1:]:
            vals = []
            for i in col_indices:
                v = row[i] if i < len(row) else None
                if v is None:
                    vals.append("")
                elif isinstance(v, float):
                    vals.append(f"{v:.4f}" if v != int(v) else str(int(v)))
                else:
                    vals.append(str(v))
            # 跳过全空行
            if any(v.strip() for v in vals):
                print(f"    {' | '.join(vals)}")
        print()
    wb.close()


def inspect_ppt(ppt_path):
    """检查PPT输出"""
    print_header("PPT 输出检查")
    if not ppt_path or not os.path.exists(ppt_path):
        print(f"  {RED}PPT 文件不存在{RESET}")
        return

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    print(f"  文件: {os.path.basename(ppt_path)}")
    print(f"  文件大小: {os.path.getsize(ppt_path) / 1024:.1f} KB\n")

    prs = Presentation(ppt_path)
    print(f"  总页数: {len(prs.slides)}\n")

    for idx, slide in enumerate(prs.slides, 1):
        print(f"  {BOLD}--- 第 {idx} 页 ---{RESET}")

        # 收集文本
        texts = []
        charts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
            if shape.has_chart:
                chart = shape.chart
                charts.append(chart)

        if texts:
            print(f"    文本内容:")
            for t in texts:
                # 截断过长的文本
                t_display = t.replace("\n", " | ") if len(t) < 100 else t[:100] + "..."
                print(f"      - {t_display}")

        if charts:
            print(f"    图表数量: {len(charts)}")
            for ci, chart in enumerate(charts, 1):
                chart_type = str(chart.chart_type) if chart.chart_type else "未知"
                print(f"      图表{ci}: 类型={chart_type}")
                if chart.has_title and chart.chart_title:
                    print(f"               标题={chart.chart_title.text_frame.text}")

                # 打印图表数据
                try:
                    plot = chart.plots[0]
                    categories = list(plot.categories)
                    print(f"               X轴类别: {categories}")

                    for si, series in enumerate(plot.series):
                        vals = list(series.values)
                        name = series.name if hasattr(series, "name") else f"系列{si+1}"
                        # 截断显示
                        vals_str = ", ".join([f"{v:.2f}" if isinstance(v, float) else str(v) for v in vals])
                        if len(vals_str) > 80:
                            vals_str = vals_str[:80] + "..."
                        print(f"               系列[{name}]: {vals_str}")
                except Exception as e:
                    print(f"               (数据读取异常: {e})")

        if not texts and not charts:
            print(f"    (空页)")
        print()


def verify_ppt_features(ppt_path):
    """验证 P0/P1 新功能是否正确实现"""
    print_header("P0/P1 新功能验证")
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    checks = []
    prs = Presentation(ppt_path)
    slides = list(prs.slides)

    # 检查页数（期望24页配置→23页实际，页24跳过）
    checks.append(("总页数=23", len(slides) == 23, f"实际{len(slides)}页"))

    def get_texts(slide):
        return [s.text_frame.text for s in slide.shapes if s.has_text_frame]

    # P1: 目录页（第2页）
    if len(slides) >= 2:
        texts = get_texts(slides[1])
        has_toc = any("目录" in t or "CONTENTS" in t for t in texts)
        checks.append(("P1 目录页", has_toc, f"文本: {texts[:3]}"))
    else:
        checks.append(("P1 目录页", False, "页数不足"))

    # P1: 章节页（第3页，深色背景+居中标题）
    if len(slides) >= 3:
        texts = get_texts(slides[2])
        has_section = any("销售汇总" in t for t in texts)
        checks.append(("P1 章节页", has_section, f"文本: {texts[:2]}"))
    else:
        checks.append(("P1 章节页", False, "页数不足"))

    # P0: 左图右文文字区（第7页，应含多行文字内容）
    if len(slides) >= 7:
        texts = get_texts(slides[6])
        has_text = any("透视分析" in t or "图表类型" in t for t in texts)
        checks.append(("P0 左图右文文字区", has_text, f"文本块数: {len(texts)}"))
    else:
        checks.append(("P0 左图右文文字区", False, "页数不足"))

    # P1: 长标签X轴旋转（第10页，8个城市>6触发旋转）
    if len(slides) >= 10:
        has_chart = False
        for shape in slides[9].shapes:
            if shape.has_chart:
                has_chart = True
                try:
                    cats = list(shape.chart.plots[0].categories)
                    checks.append(("P1 长标签旋转(>6类别)", len(cats) > 6, f"{len(cats)}个类别"))
                except Exception as e:
                    checks.append(("P1 长标签旋转(>6类别)", False, f"异常: {e}"))
                break
        if not has_chart:
            checks.append(("P1 长标签旋转(>6类别)", False, "无图表"))
    else:
        checks.append(("P1 长标签旋转(>6类别)", False, "页数不足"))

    # P1: 结论卡片（第4页，应含✦标记的独立卡片）
    if len(slides) >= 4:
        texts = get_texts(slides[3])
        has_conclusion = any("✦" in t for t in texts)
        # 检查是否有圆角矩形背景框
        has_card_bg = any(
            s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and hasattr(s, "adjustments")
            and len(s.adjustments) > 0
            for s in slides[3].shapes
        )
        checks.append(("P1 结论卡片", has_conclusion, "✦标记" + ("+背景框" if has_card_bg else "无背景框")))
    else:
        checks.append(("P1 结论卡片", False, "页数不足"))

    # P1: 结尾页（页23→0-indexed 22，已在下方"结束页2"检查，此处跳过）

    # 是否生成=否 的页面被跳过（配置24页，实际23页，"跳过测试页"不应出现）
    all_texts = [t for slide in slides for t in get_texts(slide)]
    no_skip_page = not any("跳过测试页" in t or "不应出现" in t for t in all_texts)
    checks.append(("是否生成跳过(页24)", no_skip_page and len(slides) == 23, f"配置24页→实际{len(slides)}页"))

    # P0: 百分比饼图数据为 0~1 小数（第5页饼图，引用占比透视结果）
    # 验证 PPT 拿到的是 0~1 小数，未被 Excel 的 0.0% 格式 ×100 影响
    if len(slides) >= 5:
        pct_ok = False
        pct_detail = "无饼图"
        from pptx.enum.chart import XL_CHART_TYPE
        for shape in slides[4].shapes:
            if shape.has_chart:
                chart = shape.chart
                try:
                    if chart.chart_type == XL_CHART_TYPE.PIE:
                        plot = chart.plots[0]
                        vals = list(plot.series[0].values)
                        # 饼图占比值应为 0~1 小数（如 0.376）
                        pct_ok = all(isinstance(v, (int, float)) and 0 <= v <= 1 for v in vals)
                        pct_detail = f"饼图值={vals}"
                        break
                except Exception as e:
                    pct_detail = f"异常: {e}"
        checks.append(("P0 百分比饼图0~1小数", pct_ok, pct_detail))

    # P0: 主题色验证 — 配置"珊瑚活力"主题，封面左侧块的填充色应为 dark=#2F3C7E
    if len(slides) >= 1:
        theme_ok = False
        for shape in slides[0].shapes:
            try:
                if shape.fill.type is not None:
                    rgb = str(shape.fill.fore_color.rgb)
                    if rgb == "2F3C7E":
                        theme_ok = True
                        break
            except Exception:
                pass
        checks.append(("主题色_珊瑚活力", theme_ok, "dark=#2F3C7E"))
    else:
        checks.append(("主题色_珊瑚活力", False, "页数不足"))

    # === 新增图表类型和布局验证 ===

    def _count_charts(slide):
        return sum(1 for s in slide.shapes if s.has_chart)
    from pptx.enum.chart import XL_CHART_TYPE

    # 2图上下布局（页12→0-indexed 11）：应含area + column 两种图表
    if len(slides) >= 12:
        n = _count_charts(slides[11])
        types = set()
        for s in slides[11].shapes:
            if s.has_chart:
                types.add(s.chart.chart_type)
        has_area = XL_CHART_TYPE.AREA in types
        has_column = XL_CHART_TYPE.COLUMN_CLUSTERED in types
        checks.append(("2图上下(area+column)", n == 2 and has_area and has_column,
                       f"charts={n}, area={has_area}, col={has_column}"))
    else:
        checks.append(("2图上下(area+column)", False, "页数不足"))

    # 4图布局（页13→0-indexed 12）：应含4个图表
    if len(slides) >= 13:
        n = _count_charts(slides[12])
        checks.append(("4图网格", n == 4, f"charts={n}"))
    else:
        checks.append(("4图网格", False, "页数不足"))

    # 上文下图（页14→0-indexed 13）：应有文字区+1个图表
    if len(slides) >= 14:
        n = _count_charts(slides[13])
        has_top_text = any("布局说明" in t for t in get_texts(slides[13]))
        checks.append(("上文下图", n == 1 and has_top_text, f"charts={n}, text={has_top_text}"))
    else:
        checks.append(("上文下图", False, "页数不足"))

    # 散点图（页15→0-indexed 14）
    if len(slides) >= 15:
        has_scatter = any(s.has_chart and s.chart.chart_type == XL_CHART_TYPE.XY_SCATTER
                          for s in slides[14].shapes)
        checks.append(("散点图(scatter)", has_scatter, ""))
    else:
        checks.append(("散点图(scatter)", False, "页数不足"))

    # 环形图（页16→0-indexed 15）
    if len(slides) >= 16:
        has_doughnut = any(s.has_chart and s.chart.chart_type == XL_CHART_TYPE.DOUGHNUT
                           for s in slides[15].shapes)
        checks.append(("环形图(doughnut)", has_doughnut, ""))
    else:
        checks.append(("环形图(doughnut)", False, "页数不足"))

    # 面积图（页17→0-indexed 16）
    if len(slides) >= 17:
        has_area2 = any(s.has_chart and s.chart.chart_type == XL_CHART_TYPE.AREA
                        for s in slides[16].shapes)
        checks.append(("面积图(area)", has_area2, ""))
    else:
        checks.append(("面积图(area)", False, "页数不足"))

    # 组合图（页18→0-indexed 17）
    if len(slides) >= 18:
        has_combo = any("柱+折" in t for t in get_texts(slides[17]))
        checks.append(("组合图(combo)", has_combo, ""))
    else:
        checks.append(("组合图(combo)", False, "页数不足"))

    # 结束页2（页23→0-indexed 22）
    if len(slides) >= 23:
        texts = get_texts(slides[22])
        has_ending2 = any("全特性" in t for t in texts)
        checks.append(("结束页2", has_ending2, f"text={texts[:2]}"))
    else:
        checks.append(("结束页2", False, "页数不足"))

    # 打印结果
    all_ok = True
    for name, ok, detail in checks:
        icon = "✓" if ok else "✗"
        color = GREEN if ok else RED
        print(f"  {color}{icon} {name}: {detail}{RESET}")
        if not ok:
            all_ok = False

    return all_ok


def verify_excel_output(excel_path):
    """验证透视分析 Excel 输出的完整性和数据正确性"""
    print_header("Excel 输出验证")
    checks = []

    if not excel_path or not os.path.exists(excel_path):
        print(f"  {RED}✗ Excel 文件不存在{RESET}")
        return False

    try:
        wb = openpyxl.load_workbook(excel_path, data_only=True)
    except Exception as e:
        print(f"  {RED}✗ 无法打开 Excel: {e}{RESET}")
        return False

    # 1. 检查 Sheet 数量（期望 >=10）
    expected_sheets = [
        "按地区汇总", "地区多维统计", "地区产品交叉", "地区占比",
        "产品计数占比", "销售额分箱", "华东华北汇总", "产品系列汇总",
        "地区单价分析", "block合并测试",
    ]
    actual_sheets = wb.sheetnames
    checks.append(("Sheet数量>=10", len(actual_sheets) >= 10, f"实际{len(actual_sheets)}个"))

    # 2. 检查关键 Sheet 是否存在
    for sname in expected_sheets:
        exists = sname in actual_sheets
        checks.append((f"Sheet[{sname}]存在", exists, ""))

    # 2.5 检查区块名是否使用配置的值（首行=区块名）
    expected_titles = {
        "按地区汇总": "简单分组求和",
        "地区多维统计": "多字段多聚合",
        "地区产品交叉": "交叉表",
        "销售额分箱": "分箱统计",
    }
    for sname, expected_title in expected_titles.items():
        if sname in actual_sheets:
            ws = wb[sname]
            first_cell = ws.cell(row=1, column=1).value
            checks.append((f"区块名[{sname}]", first_cell == expected_title, f"期望'{expected_title}', 实际'{first_cell}'"))

    def _read_sheet_data(sname):
        """读取 Sheet 数据（跳过区块标题行，返回表头+第一个区块的数据行）"""
        ws = wb[sname]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return [], []
        # 第1行是区块标题，第2行是表头
        header = [str(c).strip() if c is not None else "" for c in rows[1]]
        data = []
        for r in rows[2:]:
            if r[0] is None and all(c is None for c in r):
                break
            if r[0] is not None:
                data.append(r)
            else:
                break
        return header, data

    # 3. 按地区汇总 - 验证数据值
    if "按地区汇总" in actual_sheets:
        header, data = _read_sheet_data("按地区汇总")
        checks.append(("按地区汇总表头", "地区" in header and "总销售额" in header, str(header)))
        # 华东=4800, 华北=3300, 华南=4650
        vals = {str(r[0]).strip(): r[1] for r in data if r[0]}
        checks.append(("华东销售额=4800", vals.get("华东") == 4800, str(vals.get("华东"))))
        checks.append(("华北销售额=3300", vals.get("华北") == 3300, str(vals.get("华北"))))
        checks.append(("华南销售额=4650", vals.get("华南") == 4650, str(vals.get("华南"))))

    # 4. 地区多维统计 - 验证多聚合字段
    if "地区多维统计" in actual_sheets:
        header, data = _read_sheet_data("地区多维统计")
        expected_cols = ["地区", "总销售额", "平均销售额", "客户数"]
        has_all_cols = all(c in header for c in expected_cols)
        checks.append(("地区多维统计4列", has_all_cols, str(header)))
        # 华东: 总额4800, 均值1200, 客户数4
        vals = {str(r[0]).strip(): r for r in data if r[0]}
        hd = vals.get("华东")
        if hd:
            checks.append(("华东客户数=4", hd[3] == 4, str(hd[3])))
            checks.append(("华东均值=1200", hd[2] == 1200, str(hd[2])))

    # 5. 地区产品交叉 - 验证交叉表含合计列
    if "地区产品交叉" in actual_sheets:
        header, data = _read_sheet_data("地区产品交叉")
        has_total = "合计" in header
        checks.append(("交叉表含合计列", has_total, str(header)))
        # 华东合计=4800
        vals = {str(r[0]).strip(): r for r in data if r[0]}
        hd_row = vals.get("华东")
        if hd_row and has_total:
            total_idx = header.index("合计")
            checks.append(("华东合计=4800", hd_row[total_idx] == 4800, str(hd_row[total_idx])))

    # 6. 地区占比 - 验证占比值合理（0~1 小数，由 0.0% 格式自动 ×100 显示）
    if "地区占比" in actual_sheets:
        header, data = _read_sheet_data("地区占比")
        pct_vals = [r[1] for r in data if r[0] and r[1] is not None]
        all_valid = all(isinstance(v, (int, float)) and 0 <= v <= 1 for v in pct_vals)
        checks.append(("占比值0~1小数", all_valid, str(pct_vals)))
        # 验证 Excel 单元格 number_format 是百分比格式
        ws_pct = wb["地区占比"]
        # 第3行第2列是第一个占比数据单元格（第1行区块标题，第2行表头）
        pct_cell = ws_pct.cell(row=3, column=2)
        fmt_ok = "%" in str(pct_cell.number_format)
        checks.append(("占比列number_format含%", fmt_ok, f"format={pct_cell.number_format}"))
        # 验证存储值是 0.376 而非 37.6（确认未被 ×100）
        store_ok = pct_cell.value is not None and 0 < float(pct_cell.value) <= 1
        checks.append(("占比存储值未×100", store_ok, f"存储值={pct_cell.value}"))

    # 6.1 产品计数占比 - 验证 count_pct 也走 0~1 小数
    if "产品计数占比" in actual_sheets:
        header, data = _read_sheet_data("产品计数占比")
        pct_vals = [r[1] for r in data if r[0] and r[1] is not None]
        all_valid = all(isinstance(v, (int, float)) and 0 <= v <= 1 for v in pct_vals)
        checks.append(("计数占比0~1小数", all_valid, str(pct_vals)))

    # 7. 销售额分箱 - 验证分箱区间存在
    if "销售额分箱" in actual_sheets:
        header, data = _read_sheet_data("销售额分箱")
        has_bins = len(data) >= 3 and "~" in str(data[0][0])
        checks.append(("分箱区间存在", has_bins, str([r[0] for r in data[:3]])))

    # 8. block合并测试 - 验证多列合并
    if "block合并测试" in actual_sheets:
        header, data = _read_sheet_data("block合并测试")
        has_multi_col = "客户数_A" in header and "客户数_B" in header
        checks.append(("block合并双列", has_multi_col, str(header)))

    # 9. 无行维度汇总 - 横向一行输出（位置1:1对应）
    if "无行维度汇总" in actual_sheets:
        header, data = _read_sheet_data("无行维度汇总")
        # 验证列名：位置1:1 → 总销售额_sum, 平均客户数_avg
        has_col_sum = "总销售额" in header or any("销售额" in h and "求和" in h for h in header)
        has_col_avg = "平均客户数" in header or any("客户数" in h and ("均值" in h or "avg" in h) for h in header)
        # 验证只有1行数据（横向一行）
        is_single_row = len(data) == 1
        checks.append(("无行维度→横向一行", is_single_row, f"data行数={len(data)}"))
        checks.append(("无行维度→列名匹配", has_col_sum and has_col_avg, f"header={header}"))

    # 10. 区块名合并：task1(按地区汇总) + task14(按地区汇总, 同名"简单分组求和") 应合并
    if "按地区汇总" in actual_sheets:
        header, data = _read_sheet_data("按地区汇总")
        # 合并后应有 总销售额 + 总销量 两列（task1 sum销售额 + task14 sum销量，同区块名合并）
        has_merged_cols = "总销售额" in header and "总销量" in header
        checks.append(("区块名不连续合并", has_merged_cols, f"header={header}"))
        # 验证行维度列（地区）头有特殊颜色
        ws_block = wb["按地区汇总"]
        # header在第2行（第1行是区块标题），找到"地区"列的位置
        for ci in range(1, len(header) + 1):
            cell = ws_block.cell(row=2, column=ci)
            if str(cell.value) == "地区":
                dim_fill = cell.fill
                # 行维度列头颜色应区别于普通列头（HEADER_FILL="4472C4" vs DIM_HEADER_FILL="5B9BD5"）
                checks.append(("行维度列头区分色", dim_fill.start_color.rgb != "4472C4",
                               f"地区 fill={dim_fill.start_color.rgb}"))
                break

    # 11. 历史标量引用：任务15生产"总销售额"→任务16公式"销量/总销售额=销量占比"
    if "历史标量" in actual_sheets:
        _, data = _read_sheet_data("历史标量")
        if data and "总销售额" in _read_sheet_data("历史标量")[0]:
            total_sales = float(data[0][0]) if len(data) > 0 else None
        else:
            total_sales = None
        checks.append(("标量生产者_总销售额", total_sales is not None and total_sales > 0,
                       f"total_sales={total_sales}"))
        
        if "按地区汇总" in actual_sheets:
            # 任务16的"标量消费者"区块在第二个区块位置（行8-12），需要跳过第一个区块
            ws_block = wb["按地区汇总"]
            all_rows = list(ws_block.iter_rows(values_only=True))
            # 找到第二个区块标题行（特征：单元格值=="标量消费者"）
            block2_start = None
            for ri, row in enumerate(all_rows):
                if row[0] and str(row[0]).strip() == "标量消费者":
                    block2_start = ri
                    break
            if block2_start is not None and block2_start + 2 < len(all_rows):
                header2 = [str(c).strip() if c is not None else "" for c in all_rows[block2_start + 1]]
                data2 = all_rows[block2_start + 2:]
                # 只取到空白行为止
                data2 = [r for r in data2 if r[0] is not None]
            else:
                header2 = []
                data2 = []
            has_ratio_col = "销量占比" in header2
            checks.append(("标量消费者_销量占比列", has_ratio_col, f"header={header2}"))
            if has_ratio_col and total_sales:
                import pandas as pd
                df2 = pd.DataFrame(data2, columns=header2)
                region_idx = header2.index("地区") if "地区" in header2 else -1
                ratio_idx = header2.index("销量占比") if "销量占比" in header2 else -1
                expected_ratios = {"华北": 270/total_sales, "华东": 390/total_sales, "华南": 370/total_sales}
                ratio_ok = True
                for row in data2:
                    region = str(row[region_idx]) if region_idx >= 0 else ""
                    if region in expected_ratios:
                        actual_ratio = float(row[ratio_idx])
                        if abs(actual_ratio - expected_ratios[region]) > 0.001:
                            ratio_ok = False
                            break
                checks.append(("标量消费者_占比正确", ratio_ok,
                               f"expected={expected_ratios}"))

    # 12. 异名映射多列组合计算：任务21——值映射(金额,数量)与原始字段名(销售额,销量)完全不同
    if "地区均价分析" in actual_sheets:
        ws_price = wb["地区均价分析"]
        first_cell = ws_price.cell(row=1, column=1).value
        checks.append(("异名映射_区块名", first_cell == "多列组合_异名映射",
                       f"期望'多列组合_异名映射', 实际'{first_cell}'"))

        header, data = _read_sheet_data("地区均价分析")
        expected_cols = ["地区", "金额", "数量", "均价(万元/个)"]
        has_all_cols = all(c in header for c in expected_cols)
        checks.append(("异名映射_4列完整", has_all_cols, f"header={header}"))

        if data:
            vals = {str(r[0]).strip(): r for r in data if r[0]}
            # 华东: 金额=4800, 数量=390, 均价=4800/390=12.3077
            hd = vals.get("华东")
            if hd:
                checks.append(("异名映射_华东金额=4800", float(hd[1]) == 4800, str(hd[1])))
                checks.append(("异名映射_华东数量=390", float(hd[2]) == 390, str(hd[2])))
                checks.append(("异名映射_华东均价≈12.31",
                               abs(float(hd[3]) - 12.31) < 0.01, str(hd[3])))

    wb.close()

    # 打印结果
    all_ok = True
    for name, ok, detail in checks:
        icon = "✓" if ok else "✗"
        color = GREEN if ok else RED
        detail_str = f": {detail}" if detail else ""
        print(f"  {color}{icon} {name}{detail_str}{RESET}")
        if not ok:
            all_ok = False

    return all_ok


def verify_theme_config():
    """验证 P0 主题色配置解析（从配置 Excel 到 PptTheme）"""
    print_header("P0 主题色配置解析验证")
    import tempfile
    sys.path.insert(0, PROJECT_DIR)
    from app.src.excel_reader import read_config
    from app.src.ppt_theme import PptTheme

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "PPT配置"
    ws.append(["主色", "#1F4E79"])
    ws.append(["深色", "#1F3864"])
    ws.append(["accent1", "#2E75B6"])
    ws.append(["页脚", "测试报告"])
    ws.append(["字体", "思源黑体"])
    ws.append([])
    ws.append(["页码", "页面类型", "页面标题", "布局"])
    ws.append([1, "封面", "测试", ""])

    f = tempfile.mktemp(suffix=".xlsx")
    wb.save(f)

    checks = []
    try:
        cfg = read_config(f)
        t = PptTheme.from_config(cfg["colors"], cfg["general"])
        checks.append(("主色解析", t.primary == "#1F4E79", t.primary))
        checks.append(("深色解析", t.dark == "#1F3864", t.dark))
        checks.append(("强调色解析", t.accent_colors == ["#2E75B6"], str(t.accent_colors)))
        checks.append(("页脚解析", t.footer_text == "测试报告", t.footer_text))
        checks.append(("字体解析", t.font_name == "思源黑体", t.font_name))
    except Exception as e:
        checks.append(("主题色解析", False, str(e)))
    finally:
        if os.path.exists(f):
            os.remove(f)

    all_ok = True
    for name, ok, detail in checks:
        icon = "✓" if ok else "✗"
        color = GREEN if ok else RED
        print(f"  {color}{icon} {name}: {detail}{RESET}")
        if not ok:
            all_ok = False

    return all_ok


def pptx_to_images(ppt_path, output_dir):
    """用 PowerPoint COM 接口将 PPT 每页导出为 PNG 图片"""
    import win32com.client
    import pythoncom

    os.makedirs(output_dir, exist_ok=True)
    images = []
    abs_ppt = os.path.abspath(ppt_path)

    pythoncom.CoInitialize()
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        # 不可见，避免弹窗
        try:
            ppt_app.Visible = 0
        except Exception:
            pass
        pres = ppt_app.Presentations.Open(abs_ppt, WithWindow=False)

        for i, slide in enumerate(pres.Slides, 1):
            img_path = os.path.join(output_dir, f"slide_{i:02d}.png")
            # 导出为 PNG，分辨率 1280x720
            slide.Export(img_path, "PNG", 1280, 720)
            images.append(img_path)

        pres.Close()
        ppt_app.Quit()
    except Exception as e:
        print(f"  {YELLOW}PPT转图片失败: {e}{RESET}")
        # 退出时确保进程关闭
        try:
            ppt_app.Quit()
        except Exception:
            pass
    finally:
        pythoncom.CoUninitialize()

    return images


def generate_html_report(excel_path, ppt_path, output_dir):
    """使用 html_builder 生成 HTML 报告"""
    print_header("生成 HTML 报告（手机可查看）")
    
    sys.path.insert(0, PROJECT_DIR)
    from app.src.html_builder import generate_html_report as build_html
    from app.src.excel_reader import read_config
    
    config_path = os.path.join(SCRIPT_DIR, "项目配置.xlsx")
    data_path = os.path.join(SCRIPT_DIR, "测试数据.xlsx")
    ppt_pages = []
    try:
        config = read_config(config_path)
        ppt_pages = config.get("pages", [])
        print(f"  PPT 配置: {len(ppt_pages)} 页")
    except Exception as e:
        print(f"  [警告] 读取 PPT 配置失败: {e}")
    
    html_path = build_html(
        excel_path=excel_path,
        ppt_path=ppt_path,
        ppt_config=ppt_pages if ppt_pages else None,
        output_dir=output_dir,
        report_title="防护用例测试报告",
        report_subtitle="透视分析结果 + PPT预览",
    )
    return html_path


def start_preview_server(html_path):
    """使用 html_builder 启动预览服务器"""
    from app.src.html_builder import start_preview_server as start_server
    return start_server(html_path)


def verify_template_mode(pivot_excel_path):
    """验证 PPT 模板替换模式

    测试场景：
    1. 创建带占位符的 PPT 模板（文本占位符 + 图表占位符）
    2. 调用 template 子命令填充数据
    3. 读取生成的 PPT，验证占位符已被替换为透视数据
    4. 验证图表数据源已被替换

    Returns:
        tuple: (success: bool, output_ppt_path: str or None)
    """
    print_header("Step 5.6: 模板替换模式测试 (template)")

    if not pivot_excel_path or not os.path.exists(pivot_excel_path):
        print(f"  {RED}透视结果文件不存在，跳过模板测试{RESET}")
        return False, None

    # 1. 创建带占位符的 PPT 模板
    # 模板放在防护用例目录，输出结果放到透视结果所在的结果文件夹
    template_path = os.path.join(SCRIPT_DIR, "测试模板.pptx")
    # 输出路径放到透视结果所在的结果文件夹（output_xxx）
    pivot_dir = os.path.dirname(pivot_excel_path) if pivot_excel_path else SCRIPT_DIR
    output_path = os.path.join(pivot_dir, "模板填充结果.pptx")

    def _try_remove(path):
        if os.path.exists(path):
            try:
                os.remove(path)
                return True
            except Exception:
                return False
        return True

    # 尝试清理旧文件；若被占用则改用带时间戳的临时名
    use_temp_name = False
    if not _try_remove(template_path) or not _try_remove(output_path):
        use_temp_name = True
        ts = time.strftime("%Y%m%d_%H%M%S")
        template_path = os.path.join(SCRIPT_DIR, f"测试模板_{ts}.pptx")
        output_path = os.path.join(pivot_dir, f"模板填充结果_{ts}.pptx")
        print(f"  {YELLOW}⚠ 检测到旧模板文件被占用，改用临时名: {os.path.basename(template_path)}{RESET}")

    try:
        _create_test_template(template_path)
        print(f"  {GREEN}✓ 创建测试模板: {os.path.basename(template_path)}{RESET}")
    except Exception as e:
        print(f"  {RED}✗ 创建测试模板失败: {e}{RESET}")
        return False, None

    # 2. 调用 template 子命令
    ok = run_cmd(
        f'"{PYTHON}" "{os.path.join(PROJECT_DIR, "app", "main.py")}" template '
        f'"{template_path}" --pivot "{pivot_excel_path}" --image-dir "{SCRIPT_DIR}" -o "{output_path}"'
    )
    if not ok or not os.path.exists(output_path):
        print(f"  {RED}✗ 模板填充执行失败{RESET}")
        return False, None

    print(f"  {GREEN}✓ 模板填充执行成功{RESET}")

    # 3. 读取生成的 PPT 验证
    checks = []
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.chart import XL_CHART_TYPE
        prs = Presentation(output_path)

        # ---------- 页1: 文本占位符 + 图表占位符 ----------
        slide1 = prs.slides[0]
        all_text = ""
        for shape in slide1.shapes:
            if shape.has_text_frame:
                all_text += shape.text_frame.text + "\n"

        # 验证点1: 页1 文本占位符全部替换
        has_unreplaced = "{{" in all_text
        checks.append(("页1 文本占位符已替换", not has_unreplaced))
        if has_unreplaced:
            print(f"  {RED}✗ 页1 存在未替换的占位符:{RESET}")
            import re
            unreplaced = re.findall(r"\{\{[^}]+\}\}", all_text)
            for u in unreplaced[:5]:
                print(f"      {u}")
        else:
            print(f"  {GREEN}✓ 页1 文本占位符全部已替换{RESET}")

        # 验证点2: 页1 文本替换值合理（包含期望的数值）
        # 按地区汇总中总销售额应为 4800（华东）等
        value_ok = "4800" in all_text or "4800.0" in all_text
        checks.append(("页1 文本替换值正确", value_ok))
        if value_ok:
            print(f"  {GREEN}✓ 页1 文本替换值正确（包含期望数据）{RESET}")
        else:
            print(f"  {YELLOW}⚠ 页1 文本替换值未包含期望数据{RESET}")

        # 验证点3: 页1 图表存在且数据已替换
        chart_found = False
        chart_data_replaced = False
        chart_title_preserved = False
        for shape in slide1.shapes:
            if shape.has_chart:
                chart_found = True
                chart = shape.chart
                try:
                    # 验证图表标题保留模板原样（未被占位符污染）
                    if chart.has_title and chart.chart_title.has_text_frame:
                        title = chart.chart_title.text_frame.text
                        if "各地区销售汇总" in title and "{{" not in title:
                            chart_title_preserved = True
                    # 验证图表数据已替换
                    if chart.plots:
                        plot = chart.plots[0]
                        categories = list(plot.categories)
                        cat_str = [str(c) for c in categories]
                        if categories and any("华东" in c for c in cat_str):
                            chart_data_replaced = True
                        elif categories:
                            chart_data_replaced = True  # 至少有数据
                except Exception:
                    pass
                break

        checks.append(("页1 图表数据已替换", chart_found and chart_data_replaced))
        checks.append(("页1 图表标题保留模板原样", chart_title_preserved))
        if chart_found and chart_data_replaced:
            print(f"  {GREEN}✓ 页1 图表数据已替换（分类包含透视数据行值）{RESET}")
        elif chart_found:
            print(f"  {YELLOW}⚠ 页1 图表存在但数据可能未替换{RESET}")
        else:
            print(f"  {RED}✗ 页1 未找到图表{RESET}")
        if chart_title_preserved:
            print(f"  {GREEN}✓ 页1 图表标题保留模板原样{RESET}")
        else:
            print(f"  {RED}✗ 页1 图表标题未保留（被占位符污染）{RESET}")

        # ---------- 页2: 表格整体替换 ----------
        if len(prs.slides) >= 2:
            slide2 = prs.slides[1]
            table_found = False
            table_header_ok = False
            table_data_ok = False
            table_expanded = False
            for shape in slide2.shapes:
                if shape.has_table:
                    table_found = True
                    table = shape.table
                    # 验证表头包含 "地区" 和 "总销售额"
                    header_cells = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
                    table_header_ok = "地区" in header_cells and "总销售额" in header_cells
                    # 验证数据行包含 "华东"
                    all_cells_text = []
                    for r in range(len(table.rows)):
                        for c in range(len(table.columns)):
                            all_cells_text.append(table.cell(r, c).text.strip())
                    table_data_ok = any("华东" in t for t in all_cells_text)
                    # 验证表格已扩展（模板2行2列，数据4行3列，替换后应至少3行3列）
                    table_expanded = len(table.rows) >= 3 and len(table.columns) >= 3
                    break

            checks.append(("页2 表格存在", table_found))
            checks.append(("页2 表头已替换", table_header_ok))
            checks.append(("页2 表格数据已替换", table_data_ok))
            checks.append(("页2 表格行列已扩展", table_expanded))

            if table_found:
                print(f"  {GREEN}✓ 页2 表格存在{RESET}")
            else:
                print(f"  {RED}✗ 页2 未找到表格{RESET}")
            if table_header_ok:
                print(f"  {GREEN}✓ 页2 表头已替换（含地区/总销售额）{RESET}")
            else:
                print(f"  {RED}✗ 页2 表头未正确替换{RESET}")
            if table_data_ok:
                print(f"  {GREEN}✓ 页2 表格数据已替换（含华东）{RESET}")
            else:
                print(f"  {RED}✗ 页2 表格数据未替换{RESET}")
            if table_expanded:
                print(f"  {GREEN}✓ 页2 表格行列已自动扩展{RESET}")
            else:
                print(f"  {YELLOW}⚠ 页2 表格未扩展（可能数据与模板尺寸一致）{RESET}")

        # ---------- 页3: 图片替换（绝对路径 + 透视数据取路径 两种方式） ----------
        if len(prs.slides) >= 3:
            slide3 = prs.slides[2]
            pic_count = 0
            pic_replaced_count = 0
            for shape in slide3.shapes:
                if shape.shape_type == 13:  # PICTURE
                    pic_count += 1
                    # 验证图片名称不再包含 {{图片: 占位符
                    if "{{图片:" not in shape.name:
                        pic_replaced_count += 1

            pic_found = pic_count >= 3  # 应有三张图片（绝对路径 + 透视取路径 + 通配符匹配）
            pic_replaced = pic_replaced_count >= 3
            checks.append(("页3 三张图片存在", pic_found))
            checks.append(("页3 图片全部已替换", pic_replaced))

            if pic_found:
                print(f"  {GREEN}✓ 页3 三张图片存在（绝对路径 + 透视取路径 + 通配符）{RESET}")
            else:
                print(f"  {RED}✗ 页3 图片数量不足（期望3张，实际{pic_count}张）{RESET}")
            if pic_replaced:
                print(f"  {GREEN}✓ 页3 图片全部已替换{RESET}")
            else:
                print(f"  {YELLOW}⚠ 页3 图片可能未替换（{pic_replaced_count}/{pic_count}张已替换）{RESET}")

        # ---------- 页4: 多区块图表 + 默认区块 ----------
        if len(prs.slides) >= 4:
            slide4 = prs.slides[3]
            slide4_text = ""
            for shape in slide4.shapes:
                if shape.has_text_frame:
                    slide4_text += shape.text_frame.text + "\n"

            # 验证默认区块占位符已替换
            default_block_ok = "{{总销售额}}" not in slide4_text and "{{" not in slide4_text
            checks.append(("页4 默认区块占位符已替换", default_block_ok))
            if default_block_ok:
                print(f"  {GREEN}✓ 页4 默认区块占位符已替换（省略前缀生效）{RESET}")
            else:
                print(f"  {RED}✗ 页4 默认区块占位符未替换{RESET}")

            # 验证第二个图表（饼图）存在且数据已替换
            chart2_found = False
            chart2_data_ok = False
            chart2_title_preserved = False
            for shape in slide4.shapes:
                if shape.has_chart:
                    chart2_found = True
                    chart = shape.chart
                    try:
                        # 验证图表标题保留模板原样
                        if chart.has_title and chart.chart_title.has_text_frame:
                            title = chart.chart_title.text_frame.text
                            if "产品销售额占比" in title and "{{" not in title:
                                chart2_title_preserved = True
                        # 验证数据已替换
                        if chart.plots:
                            plot = chart.plots[0]
                            categories = list(plot.categories)
                            cat_str = [str(c) for c in categories]
                            # 产品销售额区块包含 产品A/产品B/产品C
                            if any("产品" in c for c in cat_str):
                                chart2_data_ok = True
                    except Exception:
                        pass
                    break

            checks.append(("页4 第二图表存在", chart2_found))
            checks.append(("页4 第二图表数据已替换", chart2_data_ok))
            checks.append(("页4 第二图表标题保留模板原样", chart2_title_preserved))
            if chart2_found:
                print(f"  {GREEN}✓ 页4 第二图表存在{RESET}")
            else:
                print(f"  {RED}✗ 页4 未找到第二图表{RESET}")
            if chart2_data_ok:
                print(f"  {GREEN}✓ 页4 第二图表数据已替换（分类含产品）{RESET}")
            else:
                print(f"  {YELLOW}⚠ 页4 第二图表数据可能未替换{RESET}")
            if chart2_title_preserved:
                print(f"  {GREEN}✓ 页4 第二图表标题保留模板原样{RESET}")
            else:
                print(f"  {RED}✗ 页4 第二图表标题未保留{RESET}")

        # ---------- 页5: 同 sheet 多区块引用（v2.18.6 核心验证） ----------
        if len(prs.slides) >= 5:
            slide5 = prs.slides[4]
            slide5_text = ""
            for shape in slide5.shapes:
                if shape.has_text_frame:
                    slide5_text += shape.text_frame.text + "\n"

            # 验证5-1: 页5 文本占位符全部替换（不应有 {{ 残留）
            page5_no_placeholder = "{{" not in slide5_text
            checks.append(("页5 文本占位符已替换", page5_no_placeholder))
            if page5_no_placeholder:
                print(f"  {GREEN}✓ 页5 文本占位符已替换{RESET}")
            else:
                print(f"  {RED}✗ 页5 存在未替换的占位符{RESET}")

            # 验证5-2: 第一个区块数据正确（简单分组求和 - 华东总销售额=4800）
            block1_ok = "4800" in slide5_text
            checks.append(("页5 区块1数据正确(4800)", block1_ok))
            if block1_ok:
                print(f"  {GREEN}✓ 页5 区块1数据正确（简单分组求和 华东=4800）{RESET}")
            else:
                print(f"  {RED}✗ 页5 区块1数据不正确{RESET}")

            # 验证5-3: 第二个区块数据正确（标量消费者 - 华东销量占比，应是小数）
            # 标量消费者区块的"销量占比"列值应在 0~1 之间（百分比小数）
            import re as _re
            # 提取"区块2 销量占比(华东):" 后面的数值
            m5 = _re.search(r"区块2 销量占比\(华东\):\s*([\d.]+)", slide5_text)
            block2_ok = False
            if m5:
                try:
                    val = float(m5.group(1))
                    # 占比值应在 0~1 之间（如 0.03）
                    block2_ok = 0 < val < 1
                except Exception:
                    pass
            checks.append(("页5 区块2数据正确(0~1小数)", block2_ok))
            if block2_ok:
                print(f"  {GREEN}✓ 页5 区块2数据正确（标量消费者 华东销量占比={m5.group(1) if m5 else 'N/A'}）{RESET}")
            else:
                print(f"  {RED}✗ 页5 区块2数据不正确（应为0~1小数）{RESET}")

            # 验证5-4: 表格引用第二个区块成功（表头应为 标量消费者 的列）
            table5_found = False
            table5_block2_ok = False
            for shape in slide5.shapes:
                if shape.has_table:
                    table5_found = True
                    table = shape.table
                    header_cells = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
                    # 标量消费者区块表头: 地区/地区销量/销量占比
                    table5_block2_ok = "销量占比" in header_cells or "地区销量" in header_cells
                    break
            checks.append(("页5 表格存在", table5_found))
            checks.append(("页5 表格引用区块2成功", table5_block2_ok))
            if table5_found:
                print(f"  {GREEN}✓ 页5 表格存在{RESET}")
            else:
                print(f"  {RED}✗ 页5 未找到表格{RESET}")
            if table5_block2_ok:
                print(f"  {GREEN}✓ 页5 表格引用区块2成功（表头含销量占比/地区销量）{RESET}")
            else:
                print(f"  {RED}✗ 页5 表格引用区块2失败{RESET}")

            # 验证5-5: 图表引用第一个区块成功（分类应含"华东"等地区）
            chart5_found = False
            chart5_block1_ok = False
            for shape in slide5.shapes:
                if shape.has_chart:
                    chart5_found = True
                    chart = shape.chart
                    try:
                        if chart.plots:
                            plot = chart.plots[0]
                            categories = list(plot.categories)
                            cat_str = [str(c) for c in categories]
                            # 简单分组求和区块行值: 华东/华北/华南
                            if any("华东" in c for c in cat_str):
                                chart5_block1_ok = True
                    except Exception:
                        pass
                    break
            checks.append(("页5 图表存在", chart5_found))
            checks.append(("页5 图表引用区块1成功", chart5_block1_ok))
            if chart5_found:
                print(f"  {GREEN}✓ 页5 图表存在{RESET}")
            else:
                print(f"  {RED}✗ 页5 未找到图表{RESET}")
            if chart5_block1_ok:
                print(f"  {GREEN}✓ 页5 图表引用区块1成功（分类含华东）{RESET}")
            else:
                print(f"  {RED}✗ 页5 图表引用区块1失败{RESET}")

        # ---------- 页6: Sheet名.区块名 精确查找（v2.18.8 核心验证） ----------
        if len(prs.slides) >= 6:
            slide6 = prs.slides[5]
            slide6_text = ""
            for shape in slide6.shapes:
                if shape.has_text_frame:
                    slide6_text += shape.text_frame.text + "\n"

            # 验证6-1: 页6 文本占位符全部替换（不应有 {{ 残留）
            page6_no_placeholder = "{{" not in slide6_text
            checks.append(("页6 文本占位符已替换", page6_no_placeholder))
            if page6_no_placeholder:
                print(f"  {GREEN}✓ 页6 文本占位符已替换{RESET}")
            else:
                print(f"  {RED}✗ 页6 存在未替换的占位符{RESET}")

            # 验证6-2: 精确查找数据正确（按地区汇总.简单分组求和.华东=4800）
            block6_ok = "4800" in slide6_text
            checks.append(("页6 精确查找数据正确(4800)", block6_ok))
            if block6_ok:
                print(f"  {GREEN}✓ 页6 精确查找数据正确（按地区汇总.简单分组求和 华东=4800）{RESET}")
            else:
                print(f"  {RED}✗ 页6 精确查找数据不正确{RESET}")

            # 验证6-3: 表格引用成功（表头应为 简单分组求和 的列：地区/总销售额/总销量）
            table6_found = False
            table6_ok = False
            for shape in slide6.shapes:
                if shape.has_table:
                    table6_found = True
                    table = shape.table
                    header_cells = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
                    # 简单分组求和区块表头: 地区/总销售额/总销量
                    table6_ok = "总销售额" in header_cells or "总销量" in header_cells
                    break
            checks.append(("页6 表格存在", table6_found))
            checks.append(("页6 表格精确查找成功", table6_ok))
            if table6_found:
                print(f"  {GREEN}✓ 页6 表格存在{RESET}")
            else:
                print(f"  {RED}✗ 页6 未找到表格{RESET}")
            if table6_ok:
                print(f"  {GREEN}✓ 页6 表格精确查找成功（表头含总销售额/总销量）{RESET}")
            else:
                print(f"  {RED}✗ 页6 表格精确查找失败{RESET}")

            # 验证6-4: 图表引用成功（分类应含"华东"等地区）
            chart6_found = False
            chart6_ok = False
            for shape in slide6.shapes:
                if shape.has_chart:
                    chart6_found = True
                    chart = shape.chart
                    try:
                        if chart.plots:
                            plot = chart.plots[0]
                            categories = list(plot.categories)
                            cat_str = [str(c) for c in categories]
                            if any("华东" in c for c in cat_str):
                                chart6_ok = True
                    except Exception:
                        pass
                    break
            checks.append(("页6 图表存在", chart6_found))
            checks.append(("页6 图表精确查找成功", chart6_ok))
            if chart6_found:
                print(f"  {GREEN}✓ 页6 图表存在{RESET}")
            else:
                print(f"  {RED}✗ 页6 未找到图表{RESET}")
            if chart6_ok:
                print(f"  {GREEN}✓ 页6 图表精确查找成功（分类含华东）{RESET}")
            else:
                print(f"  {RED}✗ 页6 图表精确查找失败{RESET}")

        # ---------- 页7: 选列替换（v2.18.11 核心验证） ----------
        if len(prs.slides) >= 7:
            slide7 = prs.slides[6]

            # 验证7-1: 表格存在
            table7_found = any(s.has_table for s in slide7.shapes)
            checks.append(("页7 表格存在", table7_found))
            if table7_found:
                print(f"  {GREEN}✓ 页7 表格存在{RESET}")
            else:
                print(f"  {RED}✗ 页7 未找到表格{RESET}")

            # 验证7-2: 表格只含选定列（地区+总销售额，不含总销量）
            table7_cols_ok = False
            if table7_found:
                for s in slide7.shapes:
                    if s.has_table:
                        t = s.table
                        headers7 = [t.cell(0, c).text_frame.text for c in range(len(t.columns))]
                        has_total_sales = any("总销售额" in h for h in headers7)
                        no_total_volume = all("总销量" not in h for h in headers7)
                        table7_cols_ok = has_total_sales and no_total_volume and len(headers7) == 2
                        break
            checks.append(("页7 表格只含选定列", table7_cols_ok))
            if table7_cols_ok:
                print(f"  {GREEN}✓ 页7 表格只含选定列（地区+总销售额，排除总销量）{RESET}")
            else:
                print(f"  {RED}✗ 页7 表格选列失败{RESET}")

            # 验证7-3: 图表存在
            chart7_found = any(s.has_chart for s in slide7.shapes)
            checks.append(("页7 图表存在", chart7_found))
            if chart7_found:
                print(f"  {GREEN}✓ 页7 图表存在{RESET}")
            else:
                print(f"  {RED}✗ 页7 未找到图表{RESET}")

            # 验证7-4: 图表只有1个系列且名称含"总销量"
            chart7_series_ok = False
            if chart7_found:
                for s in slide7.shapes:
                    if s.has_chart:
                        try:
                            plot = s.chart.plots[0]
                            series_list = list(plot.series)
                            if len(series_list) == 1:
                                name = series_list[0].name if hasattr(series_list[0], "name") else ""
                                chart7_series_ok = "总销量" in str(name)
                        except Exception:
                            pass
                        break
            checks.append(("页7 图表只含选定系列", chart7_series_ok))
            if chart7_series_ok:
                print(f"  {GREEN}✓ 页7 图表只含选定系列（总销量，单系列）{RESET}")
            else:
                print(f"  {RED}✗ 页7 图表选列失败{RESET}")

        # ---------- 页8: 计算占位符（v2.18.15 核心验证） ----------
        if len(prs.slides) >= 8:
            slide8 = prs.slides[7]
            slide8_text = ""
            for shape in slide8.shapes:
                if shape.has_text_frame:
                    slide8_text += shape.text_frame.text + "\n"

            # 验证8-1: 页8 计算占位符全部替换（不应有 {{ 残留）
            page8_no_placeholder = "{{" not in slide8_text
            checks.append(("页8 计算占位符已替换", page8_no_placeholder))
            if page8_no_placeholder:
                print(f"  {GREEN}✓ 页8 计算占位符已替换{RESET}")
            else:
                print(f"  {RED}✗ 页8 存在未替换的计算占位符{RESET}")

            # 验证8-2: 华东销售额占比 = 4800/12750 ≈ 37.65%
            calc8a_ok = "37.65%" in slide8_text
            checks.append(("页8 占比计算正确(37.65%)", calc8a_ok))
            if calc8a_ok:
                print(f"  {GREEN}✓ 页8 占比计算正确（华东 4800/12750=37.65%）{RESET}")
            else:
                print(f"  {RED}✗ 页8 占比计算不正确{RESET}")

            # 验证8-3: 华东-华北差额 = 1500
            calc8b_ok = "1500" in slide8_text
            checks.append(("页8 差额计算正确(1500)", calc8b_ok))
            if calc8b_ok:
                print(f"  {GREEN}✓ 页8 差额计算正确（4800-3300=1500）{RESET}")
            else:
                print(f"  {RED}✗ 页8 差额计算不正确{RESET}")

            # 验证8-4: 极差取整 = 1500
            calc8c_ok = "销售额极差: 1500" in slide8_text
            checks.append(("页8 极差计算正确(1500)", calc8c_ok))
            if calc8c_ok:
                print(f"  {GREEN}✓ 页8 极差计算正确（max-min=1500，取整）{RESET}")
            else:
                print(f"  {RED}✗ 页8 极差计算不正确{RESET}")

            # 验证8-5: 默认区块单段列名运算（390/1030≈37.86%）
            calc8d_ok = "37.86%" in slide8_text
            checks.append(("页8 默认区块运算正确(37.86%)", calc8d_ok))
            if calc8d_ok:
                print(f"  {GREEN}✓ 页8 默认区块运算正确（华东销量 390/1030=37.86%）{RESET}")
            else:
                print(f"  {RED}✗ 页8 默认区块运算不正确{RESET}")

            # 验证8-6: 别名-文本表达式（华东销售额=4800，来自备注 别名.华东销售额=按地区汇总.总销售额.华东）
            calc8e_ok = "别名-华东销售额: 4800" in slide8_text
            checks.append(("页8 别名文本表达式(4800)", calc8e_ok))
            if calc8e_ok:
                print(f"  {GREEN}✓ 页8 别名文本表达式正确（华东销售额=4800）{RESET}")
            else:
                print(f"  {RED}✗ 页8 别名文本表达式不正确{RESET}")

            # 验证8-7: 别名-计算表达式（华东/合计=4800/12750≈37.65%）
            calc8f_ok = "别名-利润率: 37.65%" in slide8_text
            checks.append(("页8 别名计算表达式(37.65%)", calc8f_ok))
            if calc8f_ok:
                print(f"  {GREEN}✓ 页8 别名计算表达式正确（4800/12750=37.65%）{RESET}")
            else:
                print(f"  {RED}✗ 页8 别名计算表达式不正确{RESET}")

            # 验证8-8: 文本占位符带格式后缀（华东销售额4800 → .2f 显示 4800.00）
            calc8g_ok = "文本格式-2位小数: 4800.00" in slide8_text
            checks.append(("页8 文本占位符格式后缀(.2f)", calc8g_ok))
            if calc8g_ok:
                print(f"  {GREEN}✓ 页8 文本占位符格式后缀正确（4800 → 4800.00）{RESET}")
            else:
                print(f"  {RED}✗ 页8 文本占位符格式后缀不正确{RESET}")

            # 验证8-9: 文本占位符百分比格式（占比列华东值0.3765 → .2% 显示 37.65%）
            calc8h_ok = "文本格式-百分比: 37.65%" in slide8_text
            checks.append(("页8 文本占位符格式后缀(.2%)", calc8h_ok))
            if calc8h_ok:
                print(f"  {GREEN}✓ 页8 文本占位符百分比格式正确（0.3765 → 37.65%）{RESET}")
            else:
                print(f"  {RED}✗ 页8 文本占位符百分比格式不正确{RESET}")

        # ---------- 页9: 图表标题文本占位符（v2.18.18 验证） ----------
        if len(prs.slides) >= 9:
            slide9 = prs.slides[8]
            # 找到页9的图表，读取标题文本
            chart9_title = ""
            chart9_data_ok = False
            for shape in slide9.shapes:
                if shape.has_chart:
                    try:
                        if shape.chart.has_title and shape.chart.chart_title.has_text_frame:
                            chart9_title = shape.chart.chart_title.text_frame.text
                        # 验证图表数据已替换（3个类别：华东/华北/华南）
                        try:
                            cats = list(shape.chart.plots[0].categories)
                            chart9_data_ok = "华东" in [str(c) for c in cats]
                        except Exception:
                            pass
                    except Exception:
                        pass
                    break

            # 验证9-1: 图表标题文本占位符已替换（华东4800万元）
            title9a_ok = "华东销售额4800万元" in chart9_title
            checks.append(("页9 图表标题文本占位符", title9a_ok))
            if title9a_ok:
                print(f"  {GREEN}✓ 页9 图表标题文本占位符正确（华东销售额4800万元）{RESET}")
            else:
                print(f"  {RED}✗ 页9 图表标题文本占位符不正确: {repr(chart9_title)}{RESET}")

            # 验证9-2: 图表标题计算占位符带格式后缀（占比37.65%）
            title9b_ok = "占比37.65%" in chart9_title
            checks.append(("页9 图表标题计算占位符", title9b_ok))
            if title9b_ok:
                print(f"  {GREEN}✓ 页9 图表标题计算占位符正确（占比37.65%）{RESET}")
            else:
                print(f"  {RED}✗ 页9 图表标题计算占位符不正确: {repr(chart9_title)}{RESET}")

            # 验证9-3: 图表标题无 {{ 残留
            title9c_ok = "{{" not in chart9_title
            checks.append(("页9 图表标题无残留", title9c_ok))
            if title9c_ok:
                print(f"  {GREEN}✓ 页9 图表标题无占位符残留{RESET}")
            else:
                print(f"  {RED}✗ 页9 图表标题有占位符残留: {repr(chart9_title)}{RESET}")

            # 验证9-4: 图表数据源也已替换
            checks.append(("页9 图表数据源替换", chart9_data_ok))
            if chart9_data_ok:
                print(f"  {GREEN}✓ 页9 图表数据源替换成功（华东类别存在）{RESET}")
            else:
                print(f"  {RED}✗ 页9 图表数据源未替换{RESET}")

        # ---------- 页10: 表格数据源备注区声明（v2.18.19 验证） ----------
        if len(prs.slides) >= 10:
            slide10 = prs.slides[9]
            table10_ok = False
            for shape in slide10.shapes:
                if shape.has_table:
                    table = shape.table
                    try:
                        # 验证选列：总销售额,平均销售额 → 表头应有这两列，不应有客户数
                        h0 = table.cell(0, 0).text
                        h1 = table.cell(0, 1).text if len(table.columns) > 1 else ""
                        h2 = table.cell(0, 2).text if len(table.columns) > 2 else ""
                        # 地区多维统计：地区 | 总销售额 | 平均销售额
                        table10_ok = (h0 == "地区" and "总销售额" in (h1, h2) and "平均销售额" in (h1, h2)
                                      and "客户数" not in h0 and "客户数" not in h1 and "客户数" not in h2)
                        if table10_ok:
                            print(f"  {GREEN}✓ 页10 表格备注区声明正确（地区+总销售额+平均销售额，已筛掉客户数）{RESET}")
                        else:
                            print(f"  {RED}✗ 页10 表格备注区声明不正确: 表头={h0}/{h1}/{h2}{RESET}")
                    except Exception as e:
                        print(f"  {RED}✗ 页10 表格读取异常: {e}{RESET}")
                    break
            checks.append(("页10 表格备注区声明", table10_ok))

        # ---------- 页11: 散点图单区块替换 ----------
        if len(prs.slides) >= 11:
            slide11 = prs.slides[10]
            scatter11_ok = False
            scatter11_series = 0
            for shape in slide11.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    try:
                        if chart.chart_type == XL_CHART_TYPE.XY_SCATTER:
                            scatter11_ok = True
                            scatter11_series = len(chart.series)
                    except Exception:
                        pass
                    break
            checks.append(("页11 散点图存在", scatter11_ok))
            checks.append(("页11 散点图系列数>=1", scatter11_ok and scatter11_series >= 1))
            if scatter11_ok:
                print(f"  {GREEN}✓ 页11 散点图数据替换成功（{scatter11_series}个系列）{RESET}")
            else:
                print(f"  {RED}✗ 页11 散点图未找到或替换失败{RESET}")

        # ---------- 页12: 多区块散点图跨Sheet ----------
        if len(prs.slides) >= 12:
            slide12 = prs.slides[11]
            scatter12_ok = False
            scatter12_series = 0
            scatter12_names = []
            for shape in slide12.shapes:
                if shape.has_chart:
                    chart = shape.chart
                    try:
                        if chart.chart_type == XL_CHART_TYPE.XY_SCATTER:
                            scatter12_ok = True
                            scatter12_series = len(chart.series)
                            scatter12_names = [s.name for s in chart.series]
                    except Exception:
                        pass
                    break
            checks.append(("页12 多区块散点图存在", scatter12_ok))
            checks.append(("页12 系列数>=2(多区块)", scatter12_ok and scatter12_series >= 2))
            if scatter12_ok:
                print(f"  {GREEN}✓ 页12 多区块散点图替换成功（{scatter12_series}个系列={scatter12_names}）{RESET}")
            else:
                print(f"  {RED}✗ 页12 多区块散点图未找到或替换失败{RESET}")

        # ---------- 整体文件检查 ----------
        # 验证点N: 输出文件大小合理（>0）
        file_size = os.path.getsize(output_path)
        size_ok = file_size > 1000
        checks.append((f"输出文件大小合理 ({file_size}B)", size_ok))
        if size_ok:
            print(f"  {GREEN}✓ 输出文件大小合理 ({file_size} bytes){RESET}")
        else:
            print(f"  {RED}✗ 输出文件过小 ({file_size} bytes){RESET}")

    except Exception as e:
        import traceback
        print(f"  {RED}✗ 验证过程异常: {e}{RESET}")
        traceback.print_exc()
        return False, output_path

    # 总结
    print(f"\n  {CYAN}模板替换验证项:{RESET}")
    all_ok = True
    for name, ok in checks:
        status = f"{GREEN}✓{RESET}" if ok else f"{RED}✗{RESET}"
        print(f"    {status} {name}")
        if not ok:
            all_ok = False

    return all_ok, output_path


def _create_test_template(template_path):
    """创建带占位符的测试 PPT 模板（12页，覆盖文本/图表/表格/图片四类替换）

    页1 - 文本占位符 + 图表占位符（原基础场景）
    页2 - 表格整体替换 + 表格行数扩展
    页3 - 图片替换（绝对路径 + 透视数据取路径两种方式）
    页4 - 多区块图表 + 备注声明默认区块
    页5 - 同 sheet 内多区块按区块名独立引用
    页6 - Sheet名.区块名 精确查找（区分不同 sheet 中的同名区块）
    页7 - 选列功能（{{图表:区块|列1,列2}}）
    页8 - 计算占位符 + 别名 + 文本占位符格式后缀
    页9 - 图表标题文本占位符 + 备注区声明图表数据源（方案C）
    页10 - 备注区声明表格数据源（方案C）+ 选列
    页11 - 散点图数据替换（单区块 + 多区块跨Sheet）
    页12 - 多区块散点图跨Sheet（备注区 ; 分隔）
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    # ========== 页1: 文本占位符 + 图表占位符 ==========
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    tf.text = "数据报告 - {{简单分组求和.总销售额}}"
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(28)
            run.font.bold = True

    placeholders = [
        ("总销售额", "{{简单分组求和.总销售额}}"),
        ("总销量", "{{简单分组求和.总销量}}"),
        ("行数", "{{简单分组求和.行数}}"),
        ("华东销售额", "{{按地区汇总.总销售额.华东}}"),
        ("华北销售额", "{{按地区汇总.总销售额.华北}}"),
    ]
    y_offset = Inches(1.5)
    for label, placeholder in placeholders:
        box = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(6), Inches(0.5))
        tf = box.text_frame
        tf.text = f"{label}: {placeholder}"
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(16)
        y_offset += Inches(0.6)

    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("占位数据", (1, 2, 3))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7), Inches(1.5), Inches(5.5), Inches(4),
        chart_data
    )
    chart = chart_shape.chart
    # 图表标题保留模板文字（不被占位符污染）
    try:
        chart.has_title = True
        chart.chart_title.text_frame.text = "各地区销售汇总"
    except Exception:
        pass
    # 占位符写在形状名称中（不污染标题）
    chart_shape.name = "{{图表:按地区汇总}}"

    try:
        slide.notes_slide.notes_text_frame.text = "# 模板测试\n数据源=透视结果.xlsx\n"
    except Exception:
        pass

    # ========== 页2: 表格整体替换（模板表格只有2行2列，数据4行3列，验证行列扩展） ==========
    slide2 = prs.slides.add_slide(blank_layout)

    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title2.text_frame.text = "表格替换测试页"
    for para in title2.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    rows, cols = 2, 2  # 模板表只有 2x2，透视数据 4行3列，会触发自动扩展
    table_shape = slide2.shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3))
    table = table_shape.table
    # 填充占位内容
    table.cell(0, 0).text = "占位表头1"
    table.cell(0, 1).text = "占位表头2"
    table.cell(1, 0).text = "占位数据1"
    table.cell(1, 1).text = "占位数据2"
    # 给表格形状打标记（替代文字）
    try:
        table_shape._element.attrib
        # python-pptx 没有 set alt text 的 API，通过设置 name 实现
        table_shape.name = "{{表格:按地区汇总}}"
    except Exception:
        pass

    # 页2备注
    try:
        slide2.notes_slide.notes_text_frame.text = "# 表格替换测试\n"
    except Exception:
        pass

    # ========== 页3: 图片替换（两种方式：绝对路径 + 透视数据取路径） ==========
    slide3 = prs.slides.add_slide(blank_layout)

    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title3.text_frame.text = "图片替换测试页"
    for para in title3.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # 生成一张测试图片（避免外部依赖）
    test_image_path = _generate_test_image()

    # 方式1: 绝对路径 — 占位符直接写完整路径
    if test_image_path:
        pic_shape = slide3.shapes.add_picture(
            test_image_path,
            Inches(0.5), Inches(1.5), Inches(4), Inches(3.5)
        )
        pic_shape.name = "{{图片:" + test_image_path + "}}"

    # 方式2: 从透视数据取路径 — 占位符写 区块名.列名.行值
    # 透视结果"地区图片"区块含"图片路径"列，华东行值为 _test_image.png
    if test_image_path:
        pic_shape2 = slide3.shapes.add_picture(
            test_image_path,
            Inches(5.5), Inches(1.5), Inches(4), Inches(3.5)
        )
        pic_shape2.name = "{{图片:地区图片.图片路径.华东}}"

    # 方式3: 通配符匹配 — 文件名带时间戳时模糊查找
    # _test_image*.png 匹配 _test_image.png（验证通配符功能）
    if test_image_path:
        pic_shape3 = slide3.shapes.add_picture(
            test_image_path,
            Inches(10.5), Inches(1.5), Inches(3), Inches(3.5)
        )
        pic_shape3.name = "{{图片:_test_image*.png}}"

    # 页3备注
    try:
        slide3.notes_slide.notes_text_frame.text = "# 图片替换测试\n"
    except Exception:
        pass

    # ========== 页4: 多区块图表 + 备注声明默认区块（省略前缀占位符） ==========
    slide4 = prs.slides.add_slide(blank_layout)

    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title4.text_frame.text = "多区块图表 + 默认区块测试"
    for para in title4.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # 使用备注声明本页默认区块 = 按地区汇总
    try:
        slide4.notes_slide.notes_text_frame.text = "区块=按地区汇总\n"
    except Exception:
        pass

    # 省略前缀的占位符（依赖备注声明的默认区块）
    default_block_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
    default_block_box.text_frame.text = "默认区块总销售额: {{总销售额}}"

    # 第二个图表（饼图），引用不同区块
    chart_data2 = CategoryChartData()
    chart_data2.categories = ["X", "Y", "Z"]
    chart_data2.add_series("占位", (10, 20, 30))
    chart2_shape = slide4.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(7), Inches(1.5), Inches(5.5), Inches(4),
        chart_data2
    )
    chart2 = chart2_shape.chart
    # 图表标题保留模板文字
    try:
        chart2.has_title = True
        chart2.chart_title.text_frame.text = "产品销售额占比"
    except Exception:
        pass
    # 占位符写在形状名称中
    chart2_shape.name = "{{图表:产品销售额}}"

    # ========== 页5: 同 sheet 内多区块引用（核心验证 v2.18.6） ==========
    # 透视结果 sheet "按地区汇总" 内有两个区块：
    #   - "简单分组求和"（3行3列：地区/总销售额/总销量）
    #   - "标量消费者"（3行3列：地区/地区销量/销量占比）
    # 之前只读取第一个区块，现在两个都能独立引用
    slide5 = prs.slides.add_slide(blank_layout)

    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title5.text_frame.text = "同 Sheet 多区块引用测试"
    for para in title5.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # 第一个区块：简单分组求和（第一个区块）
    block1_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(6), Inches(0.5))
    block1_box.text_frame.text = "区块1 总销售额(华东): {{简单分组求和.总销售额.华东}}"

    # 第二个区块：标量消费者（同 sheet 内第二个区块，之前访问不到）
    block2_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6), Inches(0.5))
    block2_box.text_frame.text = "区块2 销量占比(华东): {{标量消费者.销量占比.华东}}"

    # 表格引用第二个区块（验证表格也能取到正确区块）
    table_shape5 = slide5.shapes.add_table(2, 2, Inches(0.5), Inches(2.5), Inches(8), Inches(3))
    table5 = table_shape5.table
    table5.cell(0, 0).text = "占位表头1"
    table5.cell(0, 1).text = "占位表头2"
    table5.cell(1, 0).text = "占位数据1"
    table5.cell(1, 1).text = "占位数据2"
    table_shape5.name = "{{表格:标量消费者}}"

    # 图表引用第一个区块
    chart_data5 = CategoryChartData()
    chart_data5.categories = ["A", "B", "C"]
    chart_data5.add_series("占位", (1, 2, 3))
    chart5_shape = slide5.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7), Inches(1.2), Inches(5.5), Inches(4),
        chart_data5
    )
    chart5 = chart5_shape.chart
    try:
        chart5.has_title = True
        chart5.chart_title.text_frame.text = "区块1 图表"
    except Exception:
        pass
    chart5_shape.name = "{{图表:简单分组求和}}"

    try:
        slide5.notes_slide.notes_text_frame.text = "# 多区块引用测试\n"
    except Exception:
        pass

    # ========== 页6: Sheet名.区块名 精确查找（v2.18.8 核心验证） ==========
    # 透视结果中 "按地区汇总" sheet 内有 "简单分组求和" 区块
    # 当不同 sheet 出现同名区块时，单写区块名无法区分，必须用 "Sheet名.区块名" 精确指定
    slide6 = prs.slides.add_slide(blank_layout)

    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title6.text_frame.text = "Sheet名.区块名 精确查找测试"
    for para in title6.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # 文本占位符：Sheet名.区块名.列名.行值
    text6_box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(0.5))
    text6_box.text_frame.text = "精确查找 华东总销售额: {{按地区汇总.简单分组求和.总销售额.华东}}"

    # 表格占位符：{{表格:Sheet名.区块名}}
    table_shape6 = slide6.shapes.add_table(2, 2, Inches(0.5), Inches(2.0), Inches(8), Inches(3))
    table6 = table_shape6.table
    table6.cell(0, 0).text = "占位表头1"
    table6.cell(0, 1).text = "占位表头2"
    table6.cell(1, 0).text = "占位数据1"
    table6.cell(1, 1).text = "占位数据2"
    table_shape6.name = "{{表格:按地区汇总.简单分组求和}}"

    # 图表占位符：{{图表:Sheet名.区块名}}
    chart_data6 = CategoryChartData()
    chart_data6.categories = ["A", "B", "C"]
    chart_data6.add_series("占位", (1, 2, 3))
    chart6_shape = slide6.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(9), Inches(1.2), Inches(4), Inches(4),
        chart_data6
    )
    chart6 = chart6_shape.chart
    try:
        chart6.has_title = True
        chart6.chart_title.text_frame.text = "精确查找图表"
    except Exception:
        pass
    chart6_shape.name = "{{图表:按地区汇总.简单分组求和}}"

    try:
        slide6.notes_slide.notes_text_frame.text = "# Sheet名.区块名 精确查找测试\n"
    except Exception:
        pass

    # ========== 页7: 选列替换（v2.18.11 核心验证） ==========
    # 区块 "按地区汇总" 有 地区/总销售额/总销量 三列
    # 用 {{表格:区块名|列1}} 只填指定列，第一列（地区）自动保留
    # 用 {{图表:区块名|列1}} 只把指定列作为系列
    slide7 = prs.slides.add_slide(blank_layout)

    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title7.text_frame.text = "选列替换测试"
    for para in title7.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # 表格：只选"总销售额"列，结果应是 2列（地区+总销售额），不含总销量
    table_shape7 = slide7.shapes.add_table(2, 2, Inches(0.5), Inches(2.0), Inches(8), Inches(3))
    table7 = table_shape7.table
    table7.cell(0, 0).text = "占位表头1"
    table7.cell(0, 1).text = "占位表头2"
    table7.cell(1, 0).text = "占位数据1"
    table7.cell(1, 1).text = "占位数据2"
    table_shape7.name = "{{表格:按地区汇总|总销售额}}"

    # 图表：只选"总销量"列，结果应只有1个系列（总销量）
    chart_data7 = CategoryChartData()
    chart_data7.categories = ["A", "B", "C"]
    chart_data7.add_series("占位", (1, 2, 3))
    chart7_shape = slide7.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(9), Inches(1.2), Inches(4), Inches(4),
        chart_data7
    )
    chart7 = chart7_shape.chart
    try:
        chart7.has_title = True
        chart7.chart_title.text_frame.text = "选列图表"
    except Exception:
        pass
    chart7_shape.name = "{{图表:按地区汇总|总销量}}"

    try:
        slide7.notes_slide.notes_text_frame.text = "# 选列替换测试\n"
    except Exception:
        pass

    # ========== 页8: 计算占位符（v2.18.15 核心验证） ==========
    # 区块 "按地区汇总" 数据：华东 4800/390, 华北 3300/270, 华南 4650/370
    # 用 {{计算:表达式}} 做两数运算（占比/差额/极差）
    # 用 {{计算:表达式|格式}} 控制输出格式（.2%/.0f）
    slide8 = prs.slides.add_slide(blank_layout)

    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title8.text_frame.text = "计算占位符测试"
    for para in title8.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    # 文本1：行值相除算占比（4800/12750≈37.65%）
    text8a = slide8.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(0.5))
    text8a.text_frame.text = "华东销售额占比: {{计算:按地区汇总.总销售额.华东 / 按地区汇总.总销售额.sum | .2%}}"

    # 文本2：两行值相减算差额（4800-3300=1500）
    text8b = slide8.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12), Inches(0.5))
    text8b.text_frame.text = "华东-华北差额: {{计算:按地区汇总.总销售额.华东 - 按地区汇总.总销售额.华北}}"

    # 文本3：聚合极差取整（max-min=1500）
    text8c = slide8.shapes.add_textbox(Inches(0.5), Inches(2.4), Inches(12), Inches(0.5))
    text8c.text_frame.text = "销售额极差: {{计算:按地区汇总.总销售额.max - 按地区汇总.总销售额.min | .0f}}"

    # 文本4：使用 default_block（备注区声明 区块=按地区汇总），单段列名
    text8d = slide8.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12), Inches(0.5))
    text8d.text_frame.text = "华东销量占比: {{计算:总销量.华东 / 总销量.sum | .2%}}"

    # 文本5：使用别名（备注区声明 别名.华东销售额=按地区汇总.总销售额.华东）
    text8e = slide8.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(12), Inches(0.5))
    text8e.text_frame.text = "别名-华东销售额: {{华东销售额}}"

    # 文本6：使用别名引用计算表达式（备注区声明 别名.利润率=计算:利润/销售额|.2%）
    text8f = slide8.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(0.5))
    text8f.text_frame.text = "别名-利润率: {{利润率}}"

    # 文本7：文本占位符带格式后缀（华东销售额4800 → .2f 显示 4800.00）
    text8g = slide8.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.5))
    text8g.text_frame.text = "文本格式-2位小数: {{按地区汇总.总销售额.华东|.2f}}"

    # 文本8：文本占位符带百分比格式（地区占比区块.销售额占比列.华东=0.3765 → .2% 显示 37.65%）
    text8h = slide8.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12), Inches(0.5))
    text8h.text_frame.text = "文本格式-百分比: {{地区占比.销售额占比.华东|.2%}}"

    try:
        slide8.notes_slide.notes_text_frame.text = "# 计算占位符测试\n区块=按地区汇总\n别名.华东销售额=按地区汇总.总销售额.华东\n别名.利润率=计算:按地区汇总.总销售额.华东 / 按地区汇总.总销售额.sum | .2%\n"
    except Exception:
        pass

    # ========== 页9: 图表标题文本占位符（v2.18.18 验证） ==========
    slide9 = prs.slides.add_slide(blank_layout)

    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title9.text_frame.text = "图表标题占位符测试"

    # 图表标题内嵌文本占位符 + 计算占位符 + 格式后缀
    chart_data9 = CategoryChartData()
    chart_data9.categories = ["A", "B", "C"]
    chart_data9.add_series("占位", (1, 2, 3))
    chart9_shape = slide9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5), Inches(1.2), Inches(12), Inches(5.5),
        chart_data9
    )
    chart9 = chart9_shape.chart
    try:
        chart9.has_title = True
        # 标题里混用：文本占位符 + 计算占位符（带格式后缀） + 普通文字
        chart9.chart_title.text_frame.text = "华东销售额{{按地区汇总.总销售额.华东}}万元 占比{{计算:按地区汇总.总销售额.华东 / 按地区汇总.总销售额.sum|.2%}}"
    except Exception:
        pass
    # 数据源占位符通过备注区的方案C声明（不用形状名称写 {{图表:...}}）
    chart9_shape.name = "图表1"

    try:
        slide9.notes_slide.notes_text_frame.text = "# 图表标题占位符测试\n区块=按地区汇总\n图表1=按地区汇总\n"
    except Exception:
        pass

    # ========== 页10: 表格数据源备注区声明（v2.18.19 验证） ==========
    slide10 = prs.slides.add_slide(blank_layout)

    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title10.text_frame.text = "表格备注区声明测试"
    for para in title10.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    rows10, cols10 = 2, 2
    table10_shape = slide10.shapes.add_table(rows10, cols10, Inches(1), Inches(1.5), Inches(8), Inches(3))
    table10 = table10_shape.table
    table10.cell(0, 0).text = "占位表头1"
    table10.cell(0, 1).text = "占位表头2"
    table10.cell(1, 0).text = "占位数据1"
    table10.cell(1, 1).text = "占位数据2"
    # 方案C：形状名用简短标识，备注区声明映射（含选列）
    table10_shape.name = "表格1"

    try:
        slide10.notes_slide.notes_text_frame.text = "# 表格备注区声明测试\n表格1=地区多维统计|总销售额,平均销售额\n"
    except Exception:
        pass

    # ========== 页11: 散点图单区块替换 ==========
    slide11 = prs.slides.add_slide(blank_layout)

    title11 = slide11.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title11.text_frame.text = "散点图测试（城市销售额 vs 销量）"
    for para in title11.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    from pptx.chart.data import XyChartData
    scatter_data = XyChartData()
    scatter_data.add_series("占位").add_data_point(0, 0)
    scatter_shape = slide11.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(1), Inches(1.2), Inches(11), Inches(5.5),
        scatter_data
    )
    scatter_shape.name = "散点图1"

    try:
        slide11.notes_slide.notes_text_frame.text = "# 散点图单区块替换\n散点图1=城市散点数据|销售额,销量\n"
    except Exception:
        pass

    # ========== 页12: 多区块散点图跨Sheet ==========
    slide12 = prs.slides.add_slide(blank_layout)

    title12 = slide12.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    title12.text_frame.text = "多区块散点图（跨Sheet合并）"
    for para in title12.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(24)
            run.font.bold = True

    scatter2_data = XyChartData()
    scatter2_data.add_series("占位").add_data_point(0, 0)
    scatter2_shape = slide12.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(1), Inches(1.2), Inches(11), Inches(5.5),
        scatter2_data
    )
    scatter2_shape.name = "散点图2"

    try:
        slide12.notes_slide.notes_text_frame.text = "# 多区块散点图跨Sheet\n散点图2=城市散点数据|销售额,销量 ; 地区多维统计|总销售额,平均销售额\n"
    except Exception:
        pass

    prs.save(template_path)


def _generate_test_image():
    """生成一张简单的测试图片（PNG），返回路径"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        img = Image.new('RGB', (400, 300), color=(73, 109, 137))
        draw = ImageDraw.Draw(img)
        # 简单画几个矩形和文字
        draw.rectangle([50, 50, 350, 100], fill=(255, 255, 255))
        draw.rectangle([50, 150, 350, 200], fill=(255, 200, 100))
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except Exception:
            font = ImageFont.load_default()
        draw.text((100, 60), "PLACEHOLDER IMAGE", fill=(0, 0, 0), font=font)
        draw.text((100, 160), "FOR TEMPLATE TEST", fill=(0, 0, 0), font=font)

        path = os.path.join(SCRIPT_DIR, "_test_image.png")
        img.save(path)
        return path
    except Exception as e:
        print(f"  {YELLOW}⚠ 生成测试图片失败: {e}{RESET}")
        return None


def main():
    print_header("防护用例测试 - 开始")

    # Step 0: 确认数据文件存在
    config_path = os.path.join(SCRIPT_DIR, "项目配置.xlsx")
    data_path = os.path.join(SCRIPT_DIR, "测试数据.xlsx")
    if not os.path.exists(config_path) or not os.path.exists(data_path):
        print(f"  {YELLOW}数据/配置文件不存在，先运行 创建测试数据.py{RESET}")
        run_cmd(f'"{PYTHON}" "{os.path.join(SCRIPT_DIR, "创建测试数据.py")}"')

    # Step 1: 透视分析
    print_header("Step 1: 透视分析 (pivot)")
    pivot_out = os.path.join(SCRIPT_DIR, "项目配置_分析.xlsx")
    ok = run_cmd(
        f'"{PYTHON}" "{os.path.join(PROJECT_DIR, "app", "main.py")}" pivot '
        f'-c "{config_path}" --data-dir "{SCRIPT_DIR}" -o "{pivot_out}"'
    )
    if not ok:
        print(f"\n{RED}{BOLD}透视分析失败！请检查上方错误信息。{RESET}")
        sys.exit(1)

    # Step 2: 透视结果（直接使用上一步指定的输出路径）
    pivot_excel = pivot_out
    if not pivot_excel or not os.path.exists(pivot_excel):
        print(f"  {RED}透视结果文件不存在: {pivot_excel}{RESET}")
        sys.exit(1)

    # Step 3: 生成 PPT
    print_header("Step 2: PPT 生成 (ppt)")
    ppt_out = os.path.join(SCRIPT_DIR, "项目配置_报告.pptx")
    ok = run_cmd(
        f'"{PYTHON}" "{os.path.join(PROJECT_DIR, "app", "main.py")}" ppt '
        f'-c "{config_path}" --data-dir "{SCRIPT_DIR}" --pivot-file "{pivot_excel}" -o "{ppt_out}"'
    )
    if not ok:
        print(f"\n{RED}{BOLD}PPT 生成失败！请检查上方错误信息。{RESET}")
        sys.exit(1)

    # Step 4: PPT 结果（直接使用上一步指定的输出路径）
    ppt_path = ppt_out
    if not ppt_path or not os.path.exists(ppt_path):
        print(f"  {RED}PPT 文件不存在: {ppt_path}{RESET}")
        sys.exit(1)

    # Step 5: 检查输出
    inspect_excel(pivot_excel)
    inspect_ppt(ppt_path)

    # Step 5.5: 自动化验证（Excel + PPT功能 + 主题色）
    excel_ok = verify_excel_output(pivot_excel)
    features_ok = verify_ppt_features(ppt_path)
    theme_ok = verify_theme_config()

    # Step 5.6: 模板替换模式测试
    template_ok, template_ppt_path = verify_template_mode(pivot_excel)

    # Step 6: 生成HTML报告（手机可查看）
    ppt_dir = os.path.dirname(ppt_path) if ppt_path else SCRIPT_DIR
    html_path = generate_html_report(pivot_excel, ppt_path, ppt_dir)

    # 总结
    print_header("测试总结")
    print(f"  Excel: {pivot_excel}" if pivot_excel else "  Excel: N/A")
    print(f"  PPT:   {ppt_path}" if ppt_path else "  PPT:   N/A")
    print(f"  HTML:  {html_path}" if html_path else "  HTML:  N/A")

    pivot_ok = pivot_excel is not None and os.path.exists(pivot_excel)
    ppt_ok = ppt_path is not None and os.path.exists(ppt_path)

    print()
    if pivot_ok and ppt_ok and excel_ok and features_ok and theme_ok and template_ok:
        print(f"  {GREEN}{BOLD}✓ 全部通过 - 请人工检查输出格式{RESET}")
    else:
        print(f"  {RED}{BOLD}✗ 存在失败项{RESET}")
        if not pivot_ok:
            print(f"    - Excel 生成失败")
        if not ppt_ok:
            print(f"    - PPT 生成失败")
        if not excel_ok:
            print(f"    - Excel 输出验证未通过")
        if not features_ok:
            print(f"    - P0/P1 功能验证未通过")
        if not theme_ok:
            print(f"    - 主题色配置解析未通过")
        if not template_ok:
            print(f"    - 模板替换模式验证未通过")

    # 启动本地服务器供预览
    if html_path and os.path.exists(html_path):
        print(f"\n  {CYAN}启动预览服务器...{RESET}")
        start_preview_server(html_path)

    print(f"\n  {CYAN}提示: 请打开输出文件检查格式是否正常{RESET}")
    print(f"  {CYAN}  Excel: {pivot_excel}{RESET}" if pivot_excel else "")
    print(f"  {CYAN}  PPT:   {ppt_path}{RESET}" if ppt_path else "")


if __name__ == "__main__":
    main()
