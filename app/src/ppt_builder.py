import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

from src.ppt_theme import PptTheme, LAYOUT_POSITIONS


def build_ppt(config, chart_map, output_path):
    """根据页面定义逐页生成PPT（使用PPT原生图表）"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    pages = config.get("pages", [])
    colors = config.get("colors", {})
    general = config.get("general", {})
    # 美化主题：从 config 的 colors/general 构建，支持主题名
    theme_name = general.get("主题") or general.get("theme") or config.get("theme_name")
    theme = PptTheme.from_config(colors, general, theme_name=theme_name)

    if not pages:
        pages = _generate_auto_pages(config, chart_map)

    total_pages = len(pages)
    # auto 颜色计数器：单系列图填 auto 时，按图表序号从调色板循环取色
    color_counter = [0]

    # 预收集目录条目（所有内容页标题，排除封面/目录/结尾）
    toc_items = []
    for p in pages:
        pt = str(p.get("页面类型", "content")).strip().lower()
        if pt in ("cover", "封面", "toc", "目录", "ending", "结尾", "致谢"):
            continue
        title = str(p.get("页面标题", ""))
        if "|" in title:
            title = title.split("|", 1)[0].strip()
        if title:
            toc_items.append(title)

    for idx, page_def in enumerate(pages):
        page_type = str(page_def.get("页面类型", "content")).strip().lower()

        if page_type in ("cover", "封面"):
            title = str(page_def.get("页面标题", ""))
            subtitle = ""
            if "|" in title:
                parts = title.split("|", 1)
                title = parts[0].strip()
                subtitle = parts[1].strip()
            theme.add_cover_slide(prs, title, subtitle)
        elif page_type in ("toc", "目录"):
            title = str(page_def.get("页面标题", "目录"))
            if "|" in title:
                title = title.split("|", 1)[0].strip()
            theme.add_toc_slide(prs, title, toc_items)
        elif page_type in ("section", "章节"):
            title = str(page_def.get("页面标题", ""))
            subtitle = ""
            if "|" in title:
                parts = title.split("|", 1)
                title = parts[0].strip()
                subtitle = parts[1].strip()
            theme.add_section_slide(prs, title, subtitle)
        elif page_type in ("ending", "结尾", "致谢"):
            title = str(page_def.get("页面标题", "谢谢"))
            subtitle = ""
            if "|" in title:
                parts = title.split("|", 1)
                title = parts[0].strip()
                subtitle = parts[1].strip()
            theme.add_ending_slide(prs, title, subtitle)
        else:
            _add_content_slide_from_def(prs, page_def, chart_map, theme, idx + 1, total_pages, color_counter)

    prs.save(output_path)
    return output_path


def _generate_auto_pages(config, chart_map):
    """没有「PPT页面」Sheet时，自动按每页图表数分页（兼容旧模式）"""
    general = config.get("general", {})
    charts_per_slide = int(general.get("charts_per_slide") or general.get("每页图表数") or 4)
    chart_list = list(chart_map)
    pages = []
    for i in range(0, len(chart_list), charts_per_slide):
        page_charts = chart_list[i:i + charts_per_slide]
        page = {"页面类型": "content", "布局": f"{min(charts_per_slide, len(page_charts))}图"}
        for j, ch in enumerate(page_charts):
            page[f"图表{j + 1}"] = ch.get("图表标题", "")
        pages.append(page)
    return pages


def _add_content_slide_from_def(prs, page_def, chart_map, theme, page_num=0, total_pages=0, color_counter=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    theme.set_slide_bg(slide, "FFFFFF")

    layout = str(page_def.get("布局", "4图")).strip()
    title = str(page_def.get("页面标题", ""))
    # 支持"主标题|副标题"语法
    subtitle = ""
    if "|" in title:
        parts = title.split("|", 1)
        title = parts[0].strip()
        subtitle = parts[1].strip()

    if title:
        theme.add_page_title(slide, title, subtitle)

    positions = theme.layout_positions.get(layout, theme.layout_positions["4图"])

    active_charts = []
    if "charts" in page_def and page_def["charts"]:
        active_charts = page_def["charts"]
    else:
        for i in range(1, 5):
            chart_def = None
            inline_title = page_def.get(f"图表{i}标题", "")
            if inline_title and str(inline_title).strip():
                chart_def = {
                    "图表标题": str(inline_title).strip(),
                    "图表类型": str(page_def.get(f"图表{i}类型", "column")).strip(),
                    "数据Sheet": str(page_def.get(f"图表{i}Sheet", "Sheet1")).strip(),
                    "X轴范围": str(page_def.get(f"图表{i}X范围", "")).strip(),
                    "Y轴范围": str(page_def.get(f"图表{i}Y范围", "")).strip(),
                    "颜色": str(page_def.get(f"图表{i}颜色", "")).strip(),
                    "_categories": page_def.get(f"_图表{i}_categories"),
                    "_values": page_def.get(f"_图表{i}_values"),
                }
            else:
                chart_name = page_def.get(f"图表{i}", "")
                if chart_name and str(chart_name).strip():
                    chart_def = next(
                        (c for c in chart_map if c.get("图表标题") == str(chart_name).strip()),
                        None,
                    )
            if chart_def:
                active_charts.append(chart_def)

    # A4：给每个图表区加浅灰圆角矩形背景
    chart_rects = []
    for idx, chart_def in enumerate(active_charts):
        if idx >= len(positions):
            break
        left, top, width, height = positions[idx]
        theme.add_chart_bg(slide, left, top, width, height)
        chart_rects.append((left, top, width, height))

    for idx, chart_def in enumerate(active_charts):
        if idx >= len(chart_rects):
            break
        left, top, width, height = chart_rects[idx]
        _add_native_chart(slide, chart_def, theme, left, top, width, height, color_counter)

    # 文字区：左图右文/上文下图布局时渲染文字内容
    text_content = str(page_def.get("文字内容", "")).strip()
    if text_content:
        if layout == "左图右文":
            theme.add_right_text(slide, text_content)
        elif layout == "上文下图":
            theme.add_top_text(slide, text_content)

    theme.add_footer(slide, page_num, total_pages)


def _add_native_chart(slide, chart_info, theme, left, top, width, height, color_counter=None):
    chart_type = str(chart_info.get("图表类型", "column")).strip().lower()
    title = str(chart_info.get("图表标题", ""))
    categories = chart_info.get("_categories", [])
    values = chart_info.get("_values", [])
    color = str(chart_info.get("颜色", "")).strip()
    is_hierarchical = bool(chart_info.get("_is_hierarchical"))

    if chart_info.get("_is_map"):
        _add_map_slide(slide, chart_info, left, top, width, height)
        return

    # auto 或空：按图表序号从调色板循环取色，让每页图表颜色自动错开
    if not color or color.lower() == "auto":
        if color_counter is None:
            color_counter = [0]
        color = theme.palette[color_counter[0] % len(theme.palette)]
        color_counter[0] += 1
    if not categories:
        return

    if _is_combo_chart_type(chart_type):
        _add_combo_chart(slide, chart_info, chart_type, title, categories, values, color, theme, left, top, width, height)
        return

    is_scatter = chart_type == "scatter"
    xl_type = _get_chart_type(chart_type)

    if is_hierarchical and isinstance(values, dict):
        _add_hierarchical_chart(slide, chart_type, title, categories, values, color, theme, xl_type, chart_info, left, top, width, height)
        return

    series_list = []
    if isinstance(values, dict):
        for s_name, s_vals in values.items():
            series_list.append((s_name, s_vals))
    elif isinstance(values, list):
        series_list.append((title or "数据", values))

    if not series_list:
        return

    # 判断是否为百分比数据（系列名含占比/pct等关键词，且值为 0~1 小数）
    # 百分比列：柱/线图把值 ×100 让柱高合理，数据标签用 0.0% 格式
    is_pct_chart = _is_pct_series(series_list) and chart_type not in ("pie", "doughnut", "scatter")
    plot_values = series_list
    if is_pct_chart:
        plot_values = [(s_name, [v * 100 if isinstance(v, (int, float)) and v <= 1 else v for v in s_vals]) for s_name, s_vals in series_list]

    if is_scatter:
        chart_data = XyChartData()
        for s_name, s_vals in plot_values:
            series = chart_data.add_series(s_name)
            for i in range(min(len(categories), len(s_vals))):
                try:
                    x_val = float(categories[i])
                    y_val = float(s_vals[i])
                    series.add_data_point(x_val, y_val)
                except (ValueError, TypeError):
                    continue
    else:
        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]
        for s_name, s_vals in plot_values:
            chart_data.add_series(s_name, s_vals)

    chart_shape = slide.shapes.add_chart(
        xl_type, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for p in chart.chart_title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = theme.hex_to_rgb(theme.dark)

    theme.apply_series_color(chart, color, chart_type)
    theme.style_chart(chart, chart_type, is_pct=is_pct_chart)
    theme.render_conclusion(slide, chart_info, left, top, width, height)


# 百分比列名兜底关键词（与 excel_writer 保持一致，"率"已移除避免误伤）
_PCT_KEYWORDS = ["占比", "pct", "百分比", "比例"]


def _is_pct_name(name):
    if not name:
        return False
    name = str(name)
    return any(kw in name for kw in _PCT_KEYWORDS)


def _is_pct_series(series_list):
    """判断系列是否为百分比数据：系列名含占比关键词，且值域在 0~1 之间"""
    for s_name, s_vals in series_list:
        if not _is_pct_name(s_name):
            return False
    # 所有系列名都含占比关键词，再验证值域
    all_vals = []
    for _, s_vals in series_list:
        for v in s_vals:
            if isinstance(v, (int, float)):
                all_vals.append(float(v))
    if not all_vals:
        return False
    return all(0 <= v <= 1 for v in all_vals)


def _is_combo_chart_type(chart_type):
    return "," in chart_type or chart_type == "combo"


def _parse_combo_types(chart_type):
    parts = [t.strip() for t in chart_type.split(",") if t.strip()]
    return [t for t in parts if t in ("column", "line", "bar")] or ["column", "line"]


def _add_combo_chart(slide, chart_info, chart_type, title, categories, values, color, theme, left, top, width, height):
    series_list = []
    if isinstance(values, dict):
        for s_name, s_vals in values.items():
            series_list.append((s_name, s_vals))
    elif isinstance(values, list):
        series_list.append((title or "数据", values))

    if not series_list:
        return

    type_map = _parse_combo_types(chart_type)
    while len(type_map) < len(series_list):
        type_map.append("column")

    chart_data = CategoryChartData()
    # 检测层级X轴：categories为 [(parent, [(child, {...}), ...]), ...]
    if categories and isinstance(categories[0], (tuple, list)) and len(categories[0]) == 2 \
            and isinstance(categories[0][1], list):
        cats = chart_data.categories
        for parent, children in categories:
            cat = cats.add_category(str(parent))
            for child_label, _child_data in children:
                cat.add_sub_category(str(child_label))
    else:
        chart_data.categories = [str(c) for c in categories]
    for s_name, s_vals in series_list:
        chart_data.add_series(s_name, s_vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    _restructure_as_combo_chart(chart, type_map)

    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for p in chart.chart_title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = theme.hex_to_rgb(theme.dark)

    theme.apply_combo_series_color(chart, chart_type, type_map)
    theme.style_chart(chart, chart_type)

    theme.render_conclusion(slide, chart_info, left, top, width, height)


def _restructure_as_combo_chart(chart, type_map):
    chart_cs = chart._chartSpace
    chart_el = chart_cs.find(qn('c:chart'))
    if chart_el is None:
        return
    plot_area = chart_el.find(qn('c:plotArea'))
    if plot_area is None:
        return
    bar_chart = plot_area.find(qn('c:barChart'))
    if bar_chart is None:
        return

    sers = bar_chart.findall(qn('c:ser'))
    if not sers:
        return

    line_sers = []
    bar_sers = []
    for i, ser in enumerate(sers):
        if i < len(type_map) and type_map[i] == "line":
            line_sers.append(ser)
        else:
            bar_sers.append(ser)

    if not line_sers:
        return

    for ser in line_sers:
        bar_chart.remove(ser)

    if not bar_sers:
        axid_elements = bar_chart.findall(qn('c:axId'))
        for axid in axid_elements:
            bar_chart.remove(axid)

    # 清除所有旧的 axId，重新分配
    for axid in bar_chart.findall(qn('c:axId')):
        bar_chart.remove(axid)

    cat_ax = plot_area.find(qn('c:catAx'))
    left_val_ax = plot_area.find(qn('c:valAx'))
    cat_axid = cat_ax.find(qn('c:axId')).get('val') if cat_ax is not None else None
    left_val_axid = left_val_ax.find(qn('c:axId')).get('val') if left_val_ax is not None else None

    bar_axids = bar_chart.findall(qn('c:axId'))
    bar_refs_right = any(a.get('val') == (str(int(left_val_axid) + 1) if left_val_axid else None) for a in bar_axids)

    line_chart = etree.SubElement(plot_area, qn('c:lineChart'))
    grouping = etree.SubElement(line_chart, qn('c:grouping'))
    grouping.set('val', 'standard')
    for ser in line_sers:
        line_chart.append(ser)

    if cat_axid and left_val_axid and left_val_ax is not None and not bar_refs_right:
        new_val_axid = str(int(left_val_axid) + 1)
        axid_cat = etree.SubElement(line_chart, qn('c:axId'))
        axid_cat.set('val', cat_axid)
        axid_val = etree.SubElement(line_chart, qn('c:axId'))
        axid_val.set('val', new_val_axid)

        right_val_ax = etree.SubElement(plot_area, qn('c:valAx'))
        etree.SubElement(right_val_ax, qn('c:axId')).set('val', new_val_axid)
        scaling = etree.SubElement(right_val_ax, qn('c:scaling'))
        etree.SubElement(scaling, qn('c:orientation')).set('val', 'minMax')
        etree.SubElement(right_val_ax, qn('c:delete')).set('val', '0')
        etree.SubElement(right_val_ax, qn('c:axPos')).set('val', 'r')
        etree.SubElement(right_val_ax, qn('c:majorTickMark')).set('val', 'out')
        etree.SubElement(right_val_ax, qn('c:minorTickMark')).set('val', 'none')
        etree.SubElement(right_val_ax, qn('c:tickLblPos')).set('val', 'nextTo')
        etree.SubElement(right_val_ax, qn('c:crossAx')).set('val', cat_axid)
        etree.SubElement(right_val_ax, qn('c:crosses')).set('val', 'max')
        etree.SubElement(right_val_ax, qn('c:auto')).set('val', '1')

        baxid_cat = etree.SubElement(bar_chart, qn('c:axId'))
        baxid_cat.set('val', cat_axid)
        baxid_val = etree.SubElement(bar_chart, qn('c:axId'))
        baxid_val.set('val', left_val_axid)
    elif bar_sers:
        for axid in bar_axids:
            new_axid = etree.SubElement(line_chart, qn('c:axId'))
            new_axid.set('val', axid.get('val'))


def _add_map_slide(slide, chart_info, left, top, width, height):
    from src.map_builder import build_scatter_map, build_heatmap, save_map_image

    geo_df = chart_info.get("_geo_df")
    if geo_df is None or geo_df.empty:
        return

    chart_type = str(chart_info.get("图表类型", "map")).strip().lower()
    title = str(chart_info.get("图表标题", ""))
    x_range = str(chart_info.get("X轴范围", "")).strip()
    y_range = str(chart_info.get("Y轴范围", "")).strip()
    color_hex = str(chart_info.get("颜色", "")).strip() or None

    x_cols = [c.strip() for c in x_range.split(",") if c.strip()] if x_range else []
    y_cols = [c.strip() for c in y_range.split(",") if c.strip()] if y_range else []

    default_lat = [c for c in ["纬度", "lat", "latitude", "Lat"] if c in geo_df.columns]
    default_lon = [c for c in ["经度", "lon", "longitude", "lng", "Lon"] if c in geo_df.columns]

    lon_col = x_cols[0] if x_cols and x_cols[0] in geo_df.columns else (default_lon[0] if default_lon else None)
    lat_col = y_cols[0] if y_cols and y_cols[0] in geo_df.columns else (default_lat[0] if default_lat else None)

    metric_col = y_cols[-1] if len(y_cols) >= 2 and y_cols[-1] in geo_df.columns else None
    if not metric_col:
        num_cols = geo_df.select_dtypes(include=["number"]).columns.tolist()
        exclude = {lat_col, lon_col, "site_id", "cell_id", "站点ID", "小区ID"}
        metric_cols = [c for c in num_cols if c not in exclude and c not in default_lat and c not in default_lon]
        if metric_cols:
            metric_col = metric_cols[0]

    if not lat_col or not lon_col or not metric_col:
        return

    if chart_type == "heatmap":
        fig = build_heatmap(geo_df, lat_col, lon_col, metric_col, title)
    else:
        fig = build_scatter_map(geo_df, lat_col, lon_col, metric_col, title, color_hex)

    if fig is None:
        return

    # 直接用 BytesIO 流插入图片，避免临时文件残留
    png_stream = save_map_image(fig)
    slide.shapes.add_picture(png_stream, left, top, width, height)


def _get_chart_type(chart_type):
    type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "area": XL_CHART_TYPE.AREA,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    return type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)


def _add_hierarchical_chart(slide, chart_type, title, categories, values, color, theme, xl_type, chart_info, left, top, width, height):
    chart_data = CategoryChartData()
    cats = chart_data.categories

    for parent, children in categories:
        cat = cats.add_category(str(parent))
        for child_label, _child_data in children:
            cat.add_sub_category(str(child_label))

    if isinstance(values, dict):
        for series_name, series_vals in values.items():
            clean = [v if v is not None else 0 for v in series_vals]
            chart_data.add_series(str(series_name), clean)
    else:
        chart_data.add_series(title or "数据", values)

    chart_shape = slide.shapes.add_chart(
        xl_type, left, top, width, height, chart_data
    )
    chart = chart_shape.chart

    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for p in chart.chart_title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.color.rgb = theme.hex_to_rgb(theme.dark)

    theme.apply_series_color(chart, color, chart_type)
    theme.style_chart(chart, chart_type)
    theme.render_conclusion(slide, chart_info, left, top, width, height)


# ==================== PPT 配置校验 ====================

VALID_CHART_TYPES = {"bar", "column", "line", "pie", "scatter", "area", "doughnut", "combo", "map", "heatmap"}
PIVOT_DATA_SOURCE_KEYWORDS = ("{pivot}", "pivot", "透视结果", "透视分析")


def validate_ppt_config(config, config_dir, pivot_data_file=None):
    """
    校验 PPT 配置，返回 (results, all_ok)。
    results: list[dict] 每项含 level(error/warning/info)/page/chart/column/message
    all_ok: True 表示无错误
    """
    results = []
    pages = config.get("pages", [])

    if not pages:
        results.append({
            "level": "warning",
            "page": "-",
            "chart": "-",
            "column": "-",
            "message": "未检测到 PPT 页面配置（pages 为空）",
        })
        return results, True

    # 收集数据文件 Sheet 列名缓存，避免重复 IO
    sheet_columns_cache = {}

    def _get_sheet_columns(file_path, sheet_name, skip_title_row=False):
        key = (file_path, sheet_name, skip_title_row)
        if key in sheet_columns_cache:
            return sheet_columns_cache[key]
        try:
            import openpyxl
            ext = os.path.splitext(str(file_path))[1].lower()
            if ext == ".csv":
                import pandas as pd
                df = pd.read_csv(file_path, encoding="utf-8-sig", nrows=0)
                cols = list(df.columns)
            else:
                wb = openpyxl.load_workbook(file_path, read_only=True)
                if sheet_name and sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                else:
                    ws = wb[wb.sheetnames[0]]
                min_row = 2 if skip_title_row else 1
                rows_iter = ws.iter_rows(min_row=min_row, max_row=min_row, values_only=True)
                row = next(rows_iter)
                cols = [str(c).strip() if c is not None else "" for c in row]
                wb.close()
            sheet_columns_cache[key] = cols
            return cols
        except Exception as e:
            sheet_columns_cache[key] = None
            return None

    for page_idx, page in enumerate(pages, 1):
        page_num = page.get("页码", page_idx)
        page_type = str(page.get("页面类型", "内容")).strip().lower()
        layout = str(page.get("布局", "")).strip()
        charts = page.get("charts", [])

        # 封面页不校验图表配置
        if page_type in ("cover", "封面"):
            continue

        # 校验布局名
        if layout and layout not in LAYOUT_POSITIONS:
            results.append({
                "level": "error",
                "page": str(page_num),
                "chart": "-",
                "column": "布局",
                "message": f"布局 '{layout}' 不支持，可用: {', '.join(LAYOUT_POSITIONS.keys())}",
            })

        # 校验图表数量与布局槽位
        if layout and layout in LAYOUT_POSITIONS and charts:
            slot_count = len(LAYOUT_POSITIONS[layout])
            if len(charts) > slot_count:
                results.append({
                    "level": "warning",
                    "page": str(page_num),
                    "chart": "-",
                    "column": "布局",
                    "message": f"图表数 {len(charts)} 超过布局 '{layout}' 的槽位 {slot_count}，多出的将被忽略",
                })

        for chart in charts:
            chart_title = str(chart.get("图表标题", "")).strip()
            chart_type = str(chart.get("图表类型", "column")).strip().lower()
            data_sheet = str(chart.get("数据Sheet", "")).strip()
            data_source = str(chart.get("数据源", "")).strip()
            x_range = str(chart.get("X轴范围", chart.get("X轴", ""))).strip()
            y_range = str(chart.get("Y轴范围", chart.get("Y轴", ""))).strip()
            chart_id = chart_title or f"第{page_idx}页图表"

            # 校验图表标题
            if not chart_title:
                results.append({
                    "level": "warning",
                    "page": str(page_num),
                    "chart": chart_id,
                    "column": "图表标题",
                    "message": "图表标题为空",
                })

            # 校验图表类型
            if chart_type:
                # 组合图支持逗号分隔
                types = [t.strip() for t in chart_type.split(",")]
                for t in types:
                    if t and t not in VALID_CHART_TYPES:
                        results.append({
                            "level": "error",
                            "page": str(page_num),
                            "chart": chart_id,
                            "column": "图表类型",
                            "message": f"图表类型 '{t}' 不支持，可用: {', '.join(sorted(VALID_CHART_TYPES))}",
                        })

            # 跳过地图类型的列名校验（地图列名逻辑较灵活）
            if chart_type in ("map", "heatmap"):
                continue

            # 校验 X/Y 轴必填
            if not x_range:
                results.append({
                    "level": "error",
                    "page": str(page_num),
                    "chart": chart_id,
                    "column": "X轴",
                    "message": "X轴未指定",
                })
            if not y_range:
                results.append({
                    "level": "error",
                    "page": str(page_num),
                    "chart": chart_id,
                    "column": "Y轴",
                    "message": "Y轴未指定",
                })
            if not x_range or not y_range:
                continue

            # 解析数据文件路径
            is_pivot_ref = data_source.lower() in PIVOT_DATA_SOURCE_KEYWORDS
            if is_pivot_ref:
                if not pivot_data_file or not os.path.exists(str(pivot_data_file)):
                    results.append({
                        "level": "warning",
                        "page": str(page_num),
                        "chart": chart_id,
                        "column": "数据源",
                        "message": f"引用透视结果但未提供 --pivot-file，图表将被跳过",
                    })
                file_path = pivot_data_file
            elif data_source:
                from src.excel_reader import find_data_file as _find
                file_path = _find(data_source, config_dir)
                if not file_path:
                    results.append({
                        "level": "warning",
                        "page": str(page_num),
                        "chart": chart_id,
                        "column": "数据源",
                        "message": f"数据源 '{data_source}' 未找到匹配文件",
                    })
            else:
                # 无数据源时无法校验列名，跳过
                continue

            if not file_path or not os.path.exists(str(file_path)):
                continue

            # 校验 Sheet 存在性与列名匹配
            # 透视结果文件首行是区块标题，需跳过取第2行表头
            is_pivot_result = is_pivot_ref or "分析" in os.path.basename(str(file_path))
            cols = _get_sheet_columns(file_path, data_sheet, skip_title_row=is_pivot_result)
            if cols is None:
                results.append({
                    "level": "warning",
                    "page": str(page_num),
                    "chart": chart_id,
                    "column": "数据Sheet",
                    "message": f"无法读取 Sheet '{data_sheet}' 的列名（文件可能被占用）",
                })
                continue

            # 校验 Sheet 名（Excel 才校验，CSV 无 Sheet）
            ext = os.path.splitext(str(file_path))[1].lower()
            if ext != ".csv" and data_sheet:
                import openpyxl
                try:
                    wb = openpyxl.load_workbook(file_path, read_only=True)
                    if data_sheet not in wb.sheetnames:
                        results.append({
                            "level": "error",
                            "page": str(page_num),
                            "chart": chart_id,
                            "column": "数据Sheet",
                            "message": f"Sheet '{data_sheet}' 不存在，可用: {', '.join(wb.sheetnames)}",
                        })
                    wb.close()
                except Exception:
                    pass

            # 校验 X 轴列名（透视结果文件已跳过区块标题行，列名匹配准确）
            x_cols = [c.strip() for c in x_range.split(",") if c.strip()]
            for xc in x_cols:
                if xc not in cols:
                    results.append({
                        "level": "error",
                        "page": str(page_num),
                        "chart": chart_id,
                        "column": "X轴",
                        "message": f"X轴列名 '{xc}' 在 Sheet '{data_sheet}' 首行中未找到，可用列: {', '.join(c for c in cols[:10] if c)}{'...' if len(cols) > 10 else ''}",
                    })

            # 校验 Y 轴列名
            y_cols = [c.strip() for c in y_range.split(",") if c.strip()]
            for yc in y_cols:
                if yc not in cols:
                    results.append({
                        "level": "error",
                        "page": str(page_num),
                        "chart": chart_id,
                        "column": "Y轴",
                        "message": f"Y轴列名 '{yc}' 在 Sheet '{data_sheet}' 首行中未找到，可用列: {', '.join(c for c in cols[:10] if c)}{'...' if len(cols) > 10 else ''}",
                    })

    all_ok = all(r["level"] != "error" for r in results)
    return results, all_ok


def print_ppt_validation_results(results):
    """格式化打印 PPT 配置校验结果，返回 all_ok"""
    if not results:
        print("[校验] 全部通过 ✓")
        return True

    error_count = sum(1 for r in results if r["level"] == "error")
    warning_count = sum(1 for r in results if r["level"] == "warning")
    info_count = sum(1 for r in results if r["level"] == "info")

    print(f"\n{'='*60}")
    print(f"  PPT 配置校验结果")
    print(f"{'='*60}")
    print(f"  错误: {error_count}  |  警告: {warning_count}  |  提示: {info_count}")
    print(f"{'='*60}")

    for r in results:
        icon = {"error": "✗", "warning": "⚠", "info": "ℹ"}.get(r["level"], " ")
        col_info = f" [{r['column']}]" if r.get("column") and r["column"] != "-" else ""
        chart_info = f" [{r['chart']}]" if r.get("chart") and r["chart"] != "-" else ""
        print(f"  {icon} [页{r['page']}{chart_info}{col_info}] {r['message']}")

    print(f"{'='*60}")

    if error_count > 0:
        print(f"  ❌ 发现 {error_count} 个错误，请修正后再执行")
    elif warning_count > 0:
        print(f"  ⚠  发现 {warning_count} 个警告，可继续执行但可能影响结果")
    print()

    return error_count == 0
