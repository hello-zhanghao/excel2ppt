"""PPT 美化主题模块。

将 PPT 的配色、布局、字体及页面元素（封面 / 标题 / 页脚 / 图表背景 /
图表样式 / 结论卡片）渲染逻辑集中封装为 PptTheme，便于后续通过更换
主题配置或子类化实现多模板切换。

使用方式：
    from src.ppt_theme import PptTheme

    theme = PptTheme()                  # 默认主题
    theme = PptTheme(theme_config)      # 自定义主题（配色/字体/页脚等）
    theme = PptTheme.from_config(colors, general)  # 从 PPT 配置构建

    theme.add_cover_slide(prs, title, subtitle)
    theme.style_chart(chart, "column")

扩展新模板的两种方式：
    1. 数据换肤：传入不同 theme_config（颜色/字体/页脚/布局）即可
    2. 结构换肤：继承 PptTheme 覆盖 add_cover_slide / add_page_title 等方法
"""
from time import strftime
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ==================== 默认主题常量 ====================

HW_RED = "#C8102E"
HW_DARK = "#182B49"
HW_GRAY = "#8C8C8C"
HW_LIGHT = "#F7F8FA"
HW_BORDER = "#E5E7EB"
FONT_NAME = "Microsoft YaHei"

# 商务配色：主色+辅色，避免紫黄混搭
HUAWEI_PALETTE = [
    "#C8102E",  # 主红
    "#182B49",  # 深蓝
    "#5B9BD5",  # 浅蓝
    "#ED7D31",  # 橙
    "#70AD47",  # 绿
    "#A5A5A5",  # 灰
    "#FFC000",  # 金黄
    "#4472C4",  # 商务蓝
]

DEFAULT_FOOTER = "数据分析报告"

# ==================== 预设主题库（基于 PPT Generator Skill 配色方案） ====================

THEME_PRESETS = {
    "默认": {
        "primary": "#C8102E", "dark": "#182B49", "gray": "#8C8C8C",
        "light": "#F7F8FA", "border": "#E5E7EB",
        "palette": ["#C8102E", "#182B49", "#5B9BD5", "#ED7D31", "#70AD47", "#A5A5A5", "#FFC000", "#4472C4"],
        "name": "华为商务红",
        "conclusion_bg": "#FFF8E7", "conclusion_border": "#E0C060",
    },
    "午夜商务": {
        "primary": "#1E2761", "dark": "#0D1333", "gray": "#8890A0",
        "light": "#F2F4F8", "border": "#D0D5DD",
        "palette": ["#1E2761", "#CADCFC", "#3B82F6", "#10B981", "#F59E0B", "#94A3B8", "#6366F1", "#EC4899"],
        "name": "午夜商务蓝",
        "conclusion_bg": "#EEF2FF", "conclusion_border": "#A5B4FC",
    },
    "科技深空": {
        "primary": "#58A6FF", "dark": "#0D1117", "gray": "#6E7681",
        "light": "#161B22", "border": "#30363D",
        "palette": ["#58A6FF", "#3FB950", "#D29922", "#F778BA", "#79C0FF", "#7EE787", "#A5D6FF", "#FFA198"],
        "name": "GitHub 深色科技",
        "conclusion_bg": "#1A2332", "conclusion_border": "#58A6FF",
    },
    "珊瑚活力": {
        "primary": "#F96167", "dark": "#2F3C7E", "gray": "#8890A0",
        "light": "#FFF5F5", "border": "#FECDD3",
        "palette": ["#F96167", "#2F3C7E", "#F9E795", "#3B82F6", "#10B981", "#F59E0B", "#D946EF", "#6366F1"],
        "name": "珊瑚活力",
        "conclusion_bg": "#FFF0F0", "conclusion_border": "#FCA5A5",
    },
    "暖陶简约": {
        "primary": "#B85042", "dark": "#3D2C2A", "gray": "#978E8C",
        "light": "#FCFAF7", "border": "#E7D8CF",
        "palette": ["#B85042", "#A7BEAE", "#E7E8D1", "#D4A574", "#6B9080", "#A4B465", "#CB997E", "#8C6B5A"],
        "name": "暖陶简约",
        "conclusion_bg": "#FDF2EE", "conclusion_border": "#D4A574",
    },
    "海洋渐变": {
        "primary": "#065A82", "dark": "#1C2957", "gray": "#7B8DA0",
        "light": "#F0F6FA", "border": "#C4D6E4",
        "palette": ["#065A82", "#1C7293", "#21295C", "#3B82F6", "#10B981", "#06B6D4", "#0EA5E9", "#6366F1"],
        "name": "海洋渐变",
        "conclusion_bg": "#E8F4FA", "conclusion_border": "#7DD3FC",
    },
    "炭灰极简": {
        "primary": "#212121", "dark": "#36454F", "gray": "#9CA3AF",
        "light": "#F2F2F2", "border": "#D6D6D6",
        "palette": ["#36454F", "#212121", "#6B7280", "#4B5563", "#9CA3AF", "#D1D5DB", "#1F2937", "#78716C"],
        "name": "炭灰极简",
        "conclusion_bg": "#F5F5F5", "conclusion_border": "#AAAAAA",
    },
    "青绿信任": {
        "primary": "#028090", "dark": "#1B3A3A", "gray": "#7B9898",
        "light": "#F2FAFA", "border": "#C4E4E4",
        "palette": ["#028090", "#00A896", "#02C39A", "#05668D", "#02C39A", "#10B981", "#14B8A6", "#0D9488"],
        "name": "青绿信任",
        "conclusion_bg": "#E6F7F5", "conclusion_border": "#5EEAD4",
    },
    "莓果奶油": {
        "primary": "#6D2E46", "dark": "#3D1C2A", "gray": "#A08090",
        "light": "#FDF8FA", "border": "#EFD5DF",
        "palette": ["#6D2E46", "#A26769", "#ECE2D0", "#D4A5A5", "#C6878F", "#8E4D5C", "#BF7B87", "#E8D5D8"],
        "name": "莓果奶油",
        "conclusion_bg": "#FDF0F4", "conclusion_border": "#F9A8D4",
    },
    "鼠尾草静": {
        "primary": "#50808E", "dark": "#2D4A4F", "gray": "#8BA0A5",
        "light": "#F5F9F8", "border": "#D0E0DE",
        "palette": ["#84B59F", "#69A297", "#50808E", "#3B6B6B", "#A7C4B5", "#6B9080", "#B7D3C7", "#52796F"],
        "name": "鼠尾草静",
        "conclusion_bg": "#EEF6F2", "conclusion_border": "#86EFAC",
    },
    "樱桃大胆": {
        "primary": "#990011", "dark": "#1A0002", "gray": "#9C8880",
        "light": "#FFFBF5", "border": "#F2DFD7",
        "palette": ["#990011", "#FCF6F5", "#2F3C7E", "#D32F2F", "#C2185B", "#7B1FA2", "#E64A19", "#455A64"],
        "name": "樱桃大胆",
        "conclusion_bg": "#FFF0F0", "conclusion_border": "#FCA5A5",
    },
}

# 各布局对应的图表坐标（左, 上, 宽, 高），单位 Inches
LAYOUT_POSITIONS = {
    "1图": [
        (Inches(1.2), Inches(1.3), Inches(10.9), Inches(5.6)),
    ],
    "2图上下": [
        (Inches(1.2), Inches(1.0), Inches(10.9), Inches(2.8)),
        (Inches(1.2), Inches(4.2), Inches(10.9), Inches(2.8)),
    ],
    "2图左右": [
        (Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5)),
        (Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5)),
    ],
    "4图": [
        (Inches(0.6), Inches(1.15), Inches(5.8), Inches(2.75)),
        (Inches(6.9), Inches(1.15), Inches(5.8), Inches(2.75)),
        (Inches(0.6), Inches(4.2), Inches(5.8), Inches(2.75)),
        (Inches(6.9), Inches(4.2), Inches(5.8), Inches(2.75)),
    ],
    "左图右文": [
        (Inches(0.6), Inches(1.3), Inches(6.3), Inches(5.5)),
    ],
    "上文下图": [
        (Inches(1.2), Inches(2.9), Inches(10.9), Inches(4.1)),
    ],
}


def get_theme_preset(name):
    """根据主题名查找预设配置（模糊匹配，找不到返回 None）。"""
    name = str(name).strip()
    if name in THEME_PRESETS:
        return dict(THEME_PRESETS[name])
    lower = name.lower()
    for k, v in THEME_PRESETS.items():
        if k.lower() == lower or v.get("name", "").lower() == lower:
            return dict(v)
    return None


def list_theme_names():
    """返回所有预设主题名列表。"""
    return list(THEME_PRESETS.keys())


class PptTheme:
    """PPT 美化主题。

    封装配色方案、布局位置、字体及页面元素渲染逻辑。
    通过传入不同 theme_config 可切换配色 / 字体 / 页脚 / 布局；
    子类化可覆盖具体渲染方法实现全新模板风格。

    Attributes:
        primary: 主强调色（标题装饰、默认柱色）
        dark: 深色（标题文字、坐标轴文字）
        gray: 灰色（副标题、页脚、刻度）
        light: 浅色（图表区背景）
        border: 边框色（图表区边框、页脚分隔线）
        palette: 多系列图表调色板
        accent_colors: 额外强调色（优先于 palette 用于饼图等点状着色）
        font_name: 正文字体
        footer_text: 页脚左侧文本
        layout_positions: 各布局的图表坐标定义
    """

    def __init__(self, theme_config=None, theme_name=None):
        cfg = theme_config or {}
        preset = {}
        if theme_name:
            preset = get_theme_preset(theme_name) or {}
        self.primary = cfg.get("primary") or preset.get("primary") or HW_RED
        self.dark = cfg.get("dark") or preset.get("dark") or HW_DARK
        self.gray = cfg.get("gray") or preset.get("gray") or HW_GRAY
        self.light = cfg.get("light") or preset.get("light") or HW_LIGHT
        self.border = cfg.get("border") or preset.get("border") or HW_BORDER
        self.palette = cfg.get("palette") or preset.get("palette") or list(HUAWEI_PALETTE)
        self.accent_colors = cfg.get("accent_colors") or []
        self.font_name = cfg.get("font_name") or FONT_NAME
        self.footer_text = cfg.get("footer_text") or DEFAULT_FOOTER
        self.layout_positions = cfg.get("layout_positions") or LAYOUT_POSITIONS
        self.conclusion_bg = cfg.get("conclusion_bg") or preset.get("conclusion_bg") or "#FFF8E7"
        self.conclusion_border = cfg.get("conclusion_border") or preset.get("conclusion_border") or "#E0C060"

    @classmethod
    def from_config(cls, colors=None, general=None, theme_name=None):
        """从 PPT 配置的 colors / general 字典构建主题。
        
        theme_name 优先级低于 colors/general 中的显式配置，即：
        若配置了 theme_name="科技深空" 但 colors 里也有 primary="#xxx"，
        则 colors 中的显式配置会覆盖预设值。
        """
        colors = colors or {}
        general = general or {}
        accent_keys = sorted(k for k in colors if str(k).startswith("accent"))
        accent_colors = [colors[k] for k in accent_keys if colors.get(k)]

        def pick(*keys):
            for k in keys:
                v = colors.get(k) or general.get(k)
                if v:
                    return v
            return None

        return cls({
            "primary": pick("primary", "主色"),
            "dark": pick("dark", "深色"),
            "gray": pick("gray", "灰色"),
            "light": pick("light", "浅色"),
            "border": pick("border", "边框"),
            "palette": pick("palette", "调色板"),
            "accent_colors": accent_colors,
            "font_name": pick("font_name", "字体"),
            "footer_text": pick("页脚", "footer"),
        }, theme_name=theme_name)

    # ==================== 工具方法 ====================

    def hex_to_rgb(self, hex_color):
        hex_color = str(hex_color).lstrip("#")
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        return RGBColor(0x18, 0x2B, 0x49)

    def get_color_list(self):
        """返回图表配色列表：强调色优先，后接默认调色板"""
        return list(self.accent_colors) + list(self.palette)

    @staticmethod
    def _hex_lighten(hex_color, factor=0.45):
        """将 hex 颜色混入白色，返回同格式 hex 字符串（浅色版）"""
        c = str(hex_color).lstrip("#")
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        lr = int(r + (255 - r) * factor)
        lg = int(g + (255 - g) * factor)
        lb = int(b + (255 - b) * factor)
        return f"{lr:02X}{lg:02X}{lb:02X}"

    @staticmethod
    def _set_spPr_gradient(series_spPr, top_hex, bottom_hex):
        """在 spPr XML 元素上设置双色垂直渐变填充"""
        from pptx.oxml.ns import qn
        from lxml import etree
        for bad in ("a:solidFill", "a:gradFill", "a:noFill", "a:pattFill"):
            for el in series_spPr.findall(qn(bad)):
                series_spPr.remove(el)

        gradFill = etree.SubElement(series_spPr, qn("a:gradFill"))
        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))

        for pos, color in [(0, top_hex), (100000, bottom_hex)]:
            gs = etree.SubElement(gsLst, qn("a:gs"))
            gs.set("pos", str(pos))
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            srgb.set("val", color)

        lin = etree.SubElement(gradFill, qn("a:lin"))
        lin.set("ang", "5400000")
        lin.set("scaled", "1")

    def set_font_name(self, paragraph, font_name=None):
        fname = font_name or self.font_name
        for run in paragraph.runs:
            run.font.name = fname
        try:
            pPr = paragraph._pPr
            if pPr is None:
                pPr = paragraph._p.get_or_add_pPr()
            defRPr = pPr.find(qn('a:defRPr'))
            if defRPr is None:
                defRPr = pPr.makeelement(qn('a:defRPr'), {})
                pPr.insert(0, defRPr)
            defRPr.set('latin', fname)
            defRPr.set('ea', fname)
        except Exception:
            pass

    # ==================== 页面骨架 ====================

    def set_slide_bg(self, slide, hex_color=None):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.hex_to_rgb(hex_color or "FFFFFF")

    def add_cover_slide(self, prs, title, subtitle):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        self.set_slide_bg(slide, "FFFFFF")

        # 左侧深色色块（占宽 1/3，作背景）
        left_block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), prs.slide_height
        )
        left_block.fill.solid()
        left_block.fill.fore_color.rgb = self.hex_to_rgb(self.dark)
        left_block.line.fill.background()

        # 深色块右侧红色竖条（强调色）
        red_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0), Inches(0.06), prs.slide_height
        )
        red_accent.fill.solid()
        red_accent.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        red_accent.line.fill.background()

        # 左侧块顶部小红色装饰
        top_deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.6), Inches(0.8), Inches(0.06)
        )
        top_deco.fill.solid()
        top_deco.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        top_deco.line.fill.background()

        # 左侧块上的项目标识
        tag_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.8), Inches(3.5), Inches(0.5)
        )
        ttf = tag_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = "DATA REPORT"
        tp.font.size = Pt(11)
        tp.font.bold = True
        tp.font.color.rgb = self.hex_to_rgb("#FFFFFF")
        tp.alignment = PP_ALIGN.LEFT
        self.set_font_name(tp)

        # 主标题（右侧白色区域，左对齐）
        txBox = slide.shapes.add_textbox(
            Inches(5.0), Inches(2.7), Inches(7.8), Inches(1.5)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.hex_to_rgb(self.dark)
        p.alignment = PP_ALIGN.LEFT
        self.set_font_name(p)

        # 副标题
        if subtitle:
            txBox2 = slide.shapes.add_textbox(
                Inches(5.0), Inches(4.0), Inches(7.8), Inches(0.8)
            )
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(16)
            p2.font.color.rgb = self.hex_to_rgb(self.gray)
            p2.alignment = PP_ALIGN.LEFT
            self.set_font_name(p2)

        # 右下日期
        date_str = strftime("%Y.%m.%d")
        date_box = slide.shapes.add_textbox(
            Inches(5.0), Inches(6.5), Inches(5), Inches(0.4)
        )
        dtf = date_box.text_frame
        dp = dtf.paragraphs[0]
        dp.text = date_str
        dp.font.size = Pt(10)
        dp.font.color.rgb = self.hex_to_rgb(self.gray)
        dp.alignment = PP_ALIGN.LEFT
        self.set_font_name(dp)

        # 左侧块底部小标识
        foot_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(6.7), Inches(3.5), Inches(0.4)
        )
        ftf = foot_box.text_frame
        fp = ftf.paragraphs[0]
        fp.text = self.footer_text
        fp.font.size = Pt(9)
        fp.font.color.rgb = self.hex_to_rgb("#A0A8B0")
        fp.alignment = PP_ALIGN.LEFT
        self.set_font_name(fp)

        # 右侧装饰几何元素
        for i, (cx, cy, r, color_hex) in enumerate([
            (Inches(10.5), Inches(5.8), Inches(0.15), self.primary),
            (Inches(11.2), Inches(5.5), Inches(0.10), self.gray),
            (Inches(11.8), Inches(5.9), Inches(0.12), self.dark),
            (Inches(10.8), Inches(6.2), Inches(0.08), self.gray),
        ]):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, r, r)
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.hex_to_rgb(color_hex)
            circle.line.fill.background()

        # 底部装饰条
        bottom_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(6.9), Inches(7.5), Inches(0.03)
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = self.hex_to_rgb(self.border)
        bottom_bar.line.fill.background()

    def add_page_title(self, slide, title, subtitle):
        # 页面标题：左侧强调色竖条 + 标题 + 副标题
        deco_block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.28), Inches(0.10), Inches(0.5)
        )
        deco_block.fill.solid()
        deco_block.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        deco_block.line.fill.background()

        # 主标题
        left = Inches(0.85)
        top = Inches(0.2)
        width = Inches(11.85)
        height = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.hex_to_rgb(self.dark)
        p.alignment = PP_ALIGN.LEFT
        self.set_font_name(p)

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.85), Inches(0.72), Inches(11.85), Inches(0.3)
            )
            stf = sub_box.text_frame
            stf.word_wrap = True
            sp = stf.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(11)
            sp.font.color.rgb = self.hex_to_rgb(self.gray)
            sp.alignment = PP_ALIGN.LEFT
            self.set_font_name(sp)

        # 红色短线（标题下方）
        line_left = Inches(0.6)
        line_top = Inches(0.78 if not subtitle else 1.05)
        line_width = Inches(1.6)
        line_height = Inches(0.04)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, line_left, line_top, line_width, line_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        shape.line.fill.background()

    def add_footer(self, slide, page_num, total_pages):
        """页脚：左侧项目名 + 右侧页码 + 顶部分隔线"""
        # 分隔细线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.hex_to_rgb(self.border)
        line.line.fill.background()

        # 左下项目名
        left_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(7.1), Inches(6), Inches(0.3)
        )
        ltf = left_box.text_frame
        lp = ltf.paragraphs[0]
        lp.text = self.footer_text
        lp.font.size = Pt(8)
        lp.font.color.rgb = self.hex_to_rgb(self.gray)
        lp.alignment = PP_ALIGN.LEFT
        self.set_font_name(lp)

        # 右下页码
        if total_pages:
            page_text = f"{page_num} / {total_pages}"
            right_box = slide.shapes.add_textbox(
                Inches(11.0), Inches(7.1), Inches(1.7), Inches(0.3)
            )
            rtf = right_box.text_frame
            rp = rtf.paragraphs[0]
            rp.text = page_text
            rp.font.size = Pt(8)
            rp.font.color.rgb = self.hex_to_rgb(self.gray)
            rp.alignment = PP_ALIGN.RIGHT
            self.set_font_name(rp)

    def add_chart_bg(self, slide, left, top, width, height):
        """图表区背景：浅灰圆角矩形，让图表与白底页面视觉分离"""
        pad = Inches(0.1)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left - pad, top - pad,
            width + 2 * pad, height + 2 * pad
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.hex_to_rgb(self.light)
        bg.line.color.rgb = self.hex_to_rgb(self.border)
        bg.line.width = Pt(0.5)
        try:
            bg.adjustments[0] = 0.04
        except Exception:
            pass
        # 移到最底层
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def add_right_text(self, slide, text):
        left = Inches(7.2)
        top = Inches(1.2)
        width = Inches(5.6)
        height = Inches(5.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP

        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = self.hex_to_rgb("#333333")
            p.space_after = Pt(8)

    def add_top_text(self, slide, text):
        left = Inches(1.0)
        top = Inches(0.9)
        width = Inches(11.3)
        height = Inches(1.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = text.split("\n")
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = self.hex_to_rgb("#333333")
            p.space_after = Pt(4)

    def add_toc_slide(self, prs, title, toc_items):
        """目录页：左侧深色块 + 右侧目录条目列表。

        Args:
            toc_items: list[str] 目录条目标题
        """
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self.set_slide_bg(slide, "FFFFFF")

        # 左侧深色色块
        left_block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), prs.slide_height
        )
        left_block.fill.solid()
        left_block.fill.fore_color.rgb = self.hex_to_rgb(self.dark)
        left_block.line.fill.background()

        # 红色竖条
        red_accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(0), Inches(0.06), prs.slide_height
        )
        red_accent.fill.solid()
        red_accent.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        red_accent.line.fill.background()

        # 左侧"目录"标识
        tag_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.5), Inches(3.5), Inches(1.0)
        )
        ttf = tag_box.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.text = title or "目录"
        tp.font.size = Pt(36)
        tp.font.bold = True
        tp.font.color.rgb = self.hex_to_rgb("#FFFFFF")
        tp.alignment = PP_ALIGN.LEFT
        self.set_font_name(tp)

        sub_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(3.5), Inches(3.5), Inches(0.5)
        )
        stf = sub_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = "CONTENTS"
        sp.font.size = Pt(12)
        sp.font.color.rgb = self.hex_to_rgb("#A0A8B0")
        sp.alignment = PP_ALIGN.LEFT
        self.set_font_name(sp)

        # 右侧目录条目
        if toc_items:
            txBox = slide.shapes.add_textbox(
                Inches(5.2), Inches(1.5), Inches(7.3), Inches(5.0)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(toc_items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{i + 1:02d}  {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.hex_to_rgb(self.dark)
                p.space_after = Pt(14)
                self.set_font_name(p)

    def add_section_slide(self, prs, title, subtitle=""):
        """章节分隔页：深色背景 + 装饰元素 + 居中标题"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self.set_slide_bg(slide, self.dark)

        # 左上角小型强调色装饰方块
        for (dx, dy, dw) in [(1.0, 1.0, 0.06), (1.25, 1.0, 0.06)]:
            deco = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(dx), Inches(dy), Inches(dw), Inches(0.8)
            )
            deco.fill.solid()
            deco.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
            deco.line.fill.background()

        # 中央装饰短线
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.06)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        top_bar.line.fill.background()

        # 主标题（居中）
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.1), Inches(11.3), Inches(1.2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.hex_to_rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        self.set_font_name(p)

        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.4), Inches(11.3), Inches(0.6)
            )
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(14)
            sp.font.color.rgb = self.hex_to_rgb("#A0A8B0")
            sp.alignment = PP_ALIGN.CENTER
            self.set_font_name(sp)

        # 右下角微小组装饰圆点
        for (dx, dy, r) in [(12.0, 6.8, 0.08), (12.3, 6.6, 0.05)]:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(r), Inches(r))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
            dot.line.fill.background()

    def add_ending_slide(self, prs, title="谢谢", subtitle=""):
        """结尾页：深色背景 + 居中致谢文字 + 装饰元素"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self.set_slide_bg(slide, self.dark)

        # 顶部装饰元素：左上角和右上角各一组强调色圆点
        for (dx, dy, r) in [(1.0, 1.0, 0.06), (1.3, 1.0, 0.04), (11.6, 1.0, 0.06), (11.9, 1.0, 0.04)]:
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(dx), Inches(dy), Inches(r), Inches(r))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
            dot.line.fill.background()

        # 中央装饰短线
        deco = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(2.8), Inches(0.9), Inches(0.06)
        )
        deco.fill.solid()
        deco.fill.fore_color.rgb = self.hex_to_rgb(self.primary)
        deco.line.fill.background()

        # 致谢标题
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(3.1), Inches(11.3), Inches(1.2)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self.hex_to_rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        self.set_font_name(p)

        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.5), Inches(11.3), Inches(0.6)
            )
            stf = sub_box.text_frame
            sp = stf.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(14)
            sp.font.color.rgb = self.hex_to_rgb("#A0A8B0")
            sp.alignment = PP_ALIGN.CENTER
            self.set_font_name(sp)

        # 底部页脚
        foot_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(6.7), Inches(12), Inches(0.4)
        )
        ftf = foot_box.text_frame
        fp = ftf.paragraphs[0]
        fp.text = self.footer_text
        fp.font.size = Pt(9)
        fp.font.color.rgb = self.hex_to_rgb("#A0A8B0")
        fp.alignment = PP_ALIGN.CENTER
        self.set_font_name(fp)

    # ==================== 图表样式 ====================

    def apply_series_color(self, chart, color_hex, chart_type):
        rgb = self.hex_to_rgb(color_hex)
        num_series = len(chart.series)

        if chart_type in ("pie", "doughnut"):
            pie_colors = self.get_color_list()
            plot = chart.plots[0]
            series = plot.series[0]
            for idx, point in enumerate(series.points):
                point_color = pie_colors[idx % len(pie_colors)]
                light = self._hex_lighten(point_color, 0.60)
                mid = self._hex_lighten(point_color, 0.25)
                self._set_chart_point_gradient(point, light, mid)
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_category_name = True
            data_labels.show_value = False
            data_labels.font.size = Pt(8)
            data_labels.font.color.rgb = self.hex_to_rgb(self.dark)
        elif chart_type in ("line", "scatter"):
            palette = self.get_color_list()
            for idx, series in enumerate(chart.series):
                s_color = self.hex_to_rgb(palette[idx % len(palette)]) if num_series > 1 else rgb
                if chart_type == "line":
                    series.format.line.color.rgb = s_color
                    series.format.line.width = Pt(2)
                    series.marker.style = XL_MARKER_STYLE.CIRCLE
                    series.marker.size = 6
                    series.marker.format.fill.solid()
                    series.marker.format.fill.fore_color.rgb = s_color
                    series.marker.format.line.color.rgb = s_color
            # scatter — no fill needed
        else:
            palette = self.get_color_list()
            for idx, series in enumerate(chart.series):
                raw_color = palette[idx % len(palette)] if num_series > 1 else color_hex
                light = self._hex_lighten(raw_color, 0.60)
                mid = self._hex_lighten(raw_color, 0.25)
                self._set_chart_series_gradient(series, light, mid)

    def _set_chart_series_gradient(self, series, top_hex, bottom_hex):
        """在图表系列上设置从上到下的双色渐变填充"""
        from pptx.oxml.ns import qn
        from lxml import etree
        ser_el = series._element
        spPr = ser_el.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(ser_el, qn("c:spPr"))
        PptTheme._set_spPr_gradient(spPr, top_hex, bottom_hex)

    @staticmethod
    def _set_chart_point_gradient(point, top_hex, bottom_hex):
        """在饼图/环形图的单个扇区上设置渐变填充"""
        from pptx.oxml.ns import qn
        pt_el = point._element
        spPr = pt_el.find(qn("c:spPr"))
        if spPr is None:
            spPr = pt_el.makeelement(qn("c:spPr"), {})
            pt_el.insert(0, spPr)
        PptTheme._set_spPr_gradient(spPr, top_hex, bottom_hex)

    def apply_combo_series_color(self, chart, chart_type, type_map):
        palette = self.get_color_list()
        try:
            chart_cs = chart._chartSpace
            chart_el = chart_cs.find(qn('c:chart'))
            if chart_el is None:
                return
            plot_area = chart_el.find(qn('c:plotArea'))
            if plot_area is None:
                return

            bar_chart = plot_area.find(qn('c:barChart'))
            line_chart = plot_area.find(qn('c:lineChart'))

            sers = []
            if bar_chart is not None:
                sers.extend(bar_chart.findall(qn('c:ser')))
            if line_chart is not None:
                sers.extend(line_chart.findall(qn('c:ser')))

            for idx, ser in enumerate(sers):
                c = self.hex_to_rgb(palette[idx % len(palette)])
                is_line = idx < len(type_map) and type_map[idx] == "line"

                if is_line:
                    spPr = ser.find(qn('c:spPr'))
                    if spPr is None:
                        spPr = etree.SubElement(ser, qn('c:spPr'))
                    ln = spPr.find(qn('a:ln'))
                    if ln is None:
                        ln = etree.SubElement(spPr, qn('a:ln'))
                    ln.set('w', '25400')
                    solidFill = ln.find(qn('a:solidFill'))
                    if solidFill is None:
                        solidFill = etree.SubElement(ln, qn('a:solidFill'))
                    srgb = solidFill.find(qn('a:srgbClr'))
                    if srgb is None:
                        for child in list(solidFill):
                            solidFill.remove(child)
                        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
                    srgb.set('val', '{:02X}{:02X}{:02X}'.format(c[0], c[1], c[2]))

                    marker = ser.find(qn('c:marker'))
                    if marker is None:
                        marker = etree.SubElement(ser, qn('c:marker'))
                    symbol = marker.find(qn('c:symbol'))
                    if symbol is None:
                        symbol = etree.SubElement(marker, qn('c:symbol'))
                    symbol.set('val', 'circle')
                else:
                    raw_color = palette[idx % len(palette)]
                    spPr = ser.find(qn('c:spPr'))
                    if spPr is None:
                        spPr = etree.SubElement(ser, qn('c:spPr'))
                    light = self._hex_lighten(raw_color, 0.60)
                    mid = self._hex_lighten(raw_color, 0.25)
                    self.__class__._set_spPr_gradient(spPr, light, mid)
                    ln = spPr.find(qn('a:ln'))
                    if ln is not None:
                        try:
                            spPr.remove(ln)
                        except Exception:
                            pass
        except Exception as e:
            print(f"[警告] 组合图颜色应用失败: {e}")

    def style_chart(self, chart, chart_type, is_pct=False):
        try:
            num_series = len(chart.series)
            chart.has_legend = (chart_type == "pie" or num_series > 1)
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.size = Pt(9)
                chart.legend.font.color.rgb = self.hex_to_rgb(self.dark)
        except Exception as e:
            print(f"[警告] 图例样式设置失败: {e}")

        try:
            chart_cs = chart._chartSpace
            chart_el = chart_cs.find(qn('c:chart'))
            if chart_el is not None:
                plot_area = chart_el.find(qn('c:plotArea'))
                if plot_area is not None:
                    for ax_tag in ['c:valAx', 'c:catAx']:
                        for axis in plot_area.findall(qn(ax_tag)):
                            major_grid = axis.find(qn('c:majorGridlines'))
                            if major_grid is None:
                                major_grid = etree.SubElement(axis, qn('c:majorGridlines'))
                            spPr = major_grid.find(qn('c:spPr'))
                            if spPr is None:
                                spPr = etree.SubElement(major_grid, qn('c:spPr'))
                            ln = spPr.find(qn('a:ln'))
                            if ln is None:
                                ln = etree.SubElement(spPr, qn('a:ln'))
                            sf = ln.find(qn('a:solidFill'))
                            if sf is None:
                                for child in list(ln):
                                    ln.remove(child)
                                sf = etree.SubElement(ln, qn('a:solidFill'))
                            sc = sf.find(qn('a:srgbClr'))
                            if sc is None:
                                for child in list(sf):
                                    sf.remove(child)
                                sc = etree.SubElement(sf, qn('a:srgbClr'))
                            sc.set('val', 'D9D9D9')
                            ln.set('w', '9525')
        except Exception as e:
            print(f"[警告] 网格线样式设置失败: {e}")

        try:
            if chart_type not in ("pie", "doughnut", "scatter"):
                value_axis = chart.value_axis
                value_axis.tick_labels.font.size = Pt(9)
                value_axis.tick_labels.font.color.rgb = self.hex_to_rgb(self.gray)
                value_axis.has_major_gridlines = True
                # 百分比图：值轴刻度也显示 % 符号（数据已 ×100，格式 0.0"%"）
                if is_pct:
                    value_axis.tick_labels.number_format = '0.0"%"'
                if chart.has_legend:
                    value_axis.visible = True
                cat_axis = chart.category_axis
                cat_axis.tick_labels.font.size = Pt(9)
                cat_axis.tick_labels.font.color.rgb = self.hex_to_rgb(self.dark)
                # 根据标签长度/数量自动决定水平或旋转
                try:
                    raw_cats = list(chart.plots[0].categories)
                    flat_cats = []
                    for c in raw_cats:
                        if hasattr(c, "label"):
                            flat_cats.append(str(c.label))
                        else:
                            flat_cats.append(str(c))
                    self.apply_axis_label_rotation(cat_axis, flat_cats)
                except Exception:
                    self._force_axis_text_horizontal(cat_axis)
        except Exception as e:
            print(f"[警告] 坐标轴样式设置失败: {e}")

        try:
            plot = chart.plots[0]
            # scatter 图不支持数据标签，跳过
            if chart_type != "scatter":
                plot.has_data_labels = True
                data_labels = plot.data_labels
                data_labels.font.size = Pt(9)
                data_labels.font.color.rgb = self.hex_to_rgb(self.dark)
                try:
                    data_labels.show_legend_key = False
                    if chart_type in ("pie", "doughnut"):
                        data_labels.number_format = '0.0%'
                    elif is_pct:
                        # 百分比柱/线图：值已 ×100，用 0.0"%" 显示带%号
                        data_labels.show_value = True
                        data_labels.number_format = '0.0"%"'
                        if chart_type in ("line", "area"):
                            data_labels.position = XL_LABEL_POSITION.CENTER
                        else:
                            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                    else:
                        data_labels.show_value = True
                        data_labels.number_format = '#,##0.##'
                        if chart_type in ("line", "area"):
                            data_labels.position = XL_LABEL_POSITION.CENTER
                        else:
                            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
                except Exception:
                    pass
        except Exception as e:
            print(f"[警告] 数据标签样式设置失败: {e}")

    def _force_axis_text_horizontal(self, axis):
        try:
            txPr = axis.tick_labels._txPr
            bodyPr = txPr.find(qn('a:bodyPr'))
            if bodyPr is None:
                bodyPr = txPr.makeelement(qn('a:bodyPr'), {})
                txPr.append(bodyPr)
            bodyPr.set('rot', '0')
            if 'vert' in bodyPr.attrib:
                del bodyPr.attrib['vert']
        except Exception:
            pass

    def apply_axis_label_rotation(self, axis, categories):
        """根据类别标签长度和数量自动决定水平或旋转 -45°。

        规则：类别数 > 6 或任一标签长度 > 4 时旋转，避免重叠。
        """
        try:
            if not categories:
                self._force_axis_text_horizontal(axis)
                return
            max_len = max(len(str(c)) for c in categories)
            count = len(categories)
            if count > 6 or max_len > 4:
                self._rotate_axis_text(axis, -45)
            else:
                self._force_axis_text_horizontal(axis)
        except Exception:
            self._force_axis_text_horizontal(axis)

    def _rotate_axis_text(self, axis, degrees):
        """将坐标轴标签旋转指定角度（degrees 为负数表示逆时针）"""
        try:
            txPr = axis.tick_labels._txPr
            bodyPr = txPr.find(qn('a:bodyPr'))
            if bodyPr is None:
                bodyPr = txPr.makeelement(qn('a:bodyPr'), {})
                txPr.append(bodyPr)
            # rot 单位为 1/60000 度
            bodyPr.set('rot', str(int(degrees * 60000)))
            if 'vert' in bodyPr.attrib:
                del bodyPr.attrib['vert']
        except Exception:
            pass

    # ==================== 结论卡片 ====================

    def render_conclusion(self, slide, chart_info, left, chart_top, width, chart_height):
        """渲染结论卡片：带浅色背景框，位置自适应（图表下方优先，空间不足则上方）。

        Args:
            chart_top: 图表区域的顶部 y 坐标
            chart_height: 图表区域的高度
        """
        tmpl = str(chart_info.get("结论模板", "")).strip()
        if not tmpl:
            return
        categories = chart_info.get("_categories", [])
        values = chart_info.get("_values", {})

        flat_cats, flat_vals = self._flatten_chart_data(categories, values)
        if not flat_vals:
            return

        ctx = self._compute_stats(flat_cats, flat_vals)

        # 多系列时补充每个系列的命名统计
        if isinstance(values, dict) and len(values) > 1:
            for s_name in values:
                svals = [v for v in values[s_name] if v is not None]
                if svals:
                    ctx[f"max_{s_name}"] = max(svals)
                    ctx[f"min_{s_name}"] = min(svals)
                    ctx[f"avg_{s_name}"] = round(sum(svals) / len(svals), 1)
                    ctx[f"total_{s_name}"] = sum(svals)

        try:
            text = tmpl
            for key, val in ctx.items():
                if isinstance(val, float):
                    if val == int(val):
                        text = text.replace("{" + key + "}", str(int(val)))
                    else:
                        text = text.replace("{" + key + "}", f"{val:.2f}")
                else:
                    text = text.replace("{" + key + "}", str(val))
        except Exception:
            text = tmpl

        if text == tmpl:
            return

        # 自适应位置：下方优先，空间不足放上方
        card_height = Inches(0.4)
        below_top = chart_top + chart_height + Inches(0.05)
        if below_top + card_height > Inches(7.0):
            conc_top = chart_top - card_height - Inches(0.05)
        else:
            conc_top = below_top

        # 圆角矩形背景框
        bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, conc_top, width, card_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.hex_to_rgb(self.conclusion_bg)
        bg.line.color.rgb = self.hex_to_rgb(self.conclusion_border)
        bg.line.width = Pt(0.5)
        try:
            bg.adjustments[0] = 0.15
        except Exception:
            pass

        tb = slide.shapes.add_textbox(
            left, conc_top, width, card_height
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "✦ " + text
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self.hex_to_rgb(self.dark)
        p.alignment = PP_ALIGN.CENTER
        self.set_font_name(p)

    def _flatten_chart_data(self, categories, values):
        """
        展平图表数据用于统计。
        多系列时合并所有系列的 (类别, 值)，而非只取第一个系列，
        以确保 _compute_stats 的全局 max/min/avg 反映完整数据。
        """
        # 先展开单系列的类别标签
        if isinstance(categories, list) and categories and isinstance(categories[0], (tuple, list)):
            base_cats = [str(child_label) for parent, children in categories for child_label, _ in children]
        else:
            base_cats = [str(c) for c in categories]

        flat_cats = []
        flat_vals = []

        if isinstance(values, dict):
            # 多系列：每个系列都用相同的 base_cats，合并所有 (cat, val)
            for s_vals in values.values():
                for i, v in enumerate(s_vals):
                    if v is not None:
                        flat_vals.append(v)
                        flat_cats.append(base_cats[i] if i < len(base_cats) else f"项{i+1}")
        elif isinstance(values, list):
            for i, v in enumerate(values):
                if v is not None:
                    flat_vals.append(v)
                    flat_cats.append(base_cats[i] if i < len(base_cats) else f"项{i+1}")

        return flat_cats, flat_vals

    def _compute_stats(self, categories, values):
        if not values:
            return {}
        # 防御非数值数据：过滤掉无法转 float 的值，避免 sort/sum 抛 TypeError
        pairs = []
        for cat, v in zip(categories, values):
            try:
                pairs.append((cat, float(v)))
            except (TypeError, ValueError):
                continue
        if not pairs:
            return {}
        pairs.sort(key=lambda x: x[1], reverse=True)
        total = sum(p[1] for p in pairs)
        avg = round(total / len(pairs), 1)
        ctx = {
            "max_val": pairs[0][1],
            "max_cat": pairs[0][0],
            "min_val": pairs[-1][1],
            "min_cat": pairs[-1][0],
            "avg": avg,
            "total": total,
            "count": len(pairs),
        }
        if len(pairs) >= 3:
            ctx["top3_cats"] = ", ".join(p[0] for p in pairs[:3])
            ctx["top3_vals"] = ", ".join(str(p[1]) for p in pairs[:3])
            ctx["bottom3_cats"] = ", ".join(p[0] for p in pairs[-3:])
            ctx["bottom3_vals"] = ", ".join(str(p[1]) for p in pairs[-3:])
        return ctx
