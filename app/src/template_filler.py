"""
PPT 模板填充器 — 基于已有 PPT 模板和透视结果进行数据替换

用法：
    from src.template_filler import fill_template
    fill_template("模板.pptx", "透视结果.xlsx", "输出.pptx")

占位符语法：
    区块查找支持三种形式（图表/表格/图片占位符的"区块名"位置同样适用）：
        区块名                       跨所有 sheet 查找（兼容旧用法）
        Sheet名.区块名               精确指定 sheet 内的区块
        Sheet名                      该 sheet 第一个区块（兼容旧用法）

    文本占位符（在文本框中）：
        {{区块名.列名}}              取该列合计值（数值列）或首行值
        {{区块名.列名.行值}}         精确取某行某列的值
        {{区块名.行数}}              取该区块的数据行数
        {{标量.标量名}}              取无行维度标量
        {{Sheet名.区块名.列名}}      精确指定 sheet 内区块（同名区块不冲突）
        {{Sheet名.区块名.列名.行值}} 精确指定 sheet 内区块的某行

    文本占位符可选格式后缀（仅取到数值时生效）：
        {{区块名.列名|.2f}}          保留2位小数
        {{区块名.列名|.0f}}          取整
        {{区块名.列名|.2%}}          百分比（值×100 加 %）
        {{区块名.列名|int}}          整数
        格式与计算占位符共用，详见 _KNOWN_FMTS

    图表占位符（在图表形状名称/替代文字中）：
        {{图表:区块名}}              用该区块数据替换图表数据源
        {{图表:Sheet名.区块名}}      精确指定 sheet 内区块
        {{图表:区块名|列1,列2}}      只使用指定列（第一列自动作为X轴类别）
        {{图表:区块名|xy}}           散点图独立 xy 对模式：列按 (X1,Y1,X2,Y2,...) 配对
        {{图表:区块名|xy|列1,列2}}   xy 对模式 + 指定列（列数必须为偶数）

    图表标题（chart title）：
        标题文本框支持所有文本占位符语法（{{区块.列}}、{{计算:...}}、别名、格式后缀等）
        示例：标题写 "华东{{按地区汇总.总销售额.华东}}万元销售情况"
              替换后 "华东4800万元销售情况"

    图片占位符（在图片替代文字/名称中，或在文本框中）：
        {{图片:文件路径}}            替换为指定图片（绝对路径或相对模板目录）
        {{图片:@output/相对路径}}    从输出目录查找（解决带时间戳的动态输出目录）
        {{图片:路径*.png}}           支持通配符 * 和 ?（文件名/目录名带时间戳时模糊匹配）
        {{图片:区块名.列名.行值}}    取透视数据中的图片路径

    表格占位符（在表格替代文字/名称中）：
        {{表格:区块名}}              用该区块数据整表替换（保留模板表格样式）
        {{表格:Sheet名.区块名}}      精确指定 sheet 内区块
        {{表格:区块名|列1,列2}}      只填指定列（第一列自动保留）

    聚合后缀（可选）：
        {{区块名.列名.sum}}          求和（默认）
        {{区块名.列名.avg}}          平均
        {{区块名.列名.max}}          最大
        {{区块名.列名.min}}          最小

    计算占位符（在文本框中，标识符沿用上述文本占位符语法）：
        {{计算:表达式}}               默认智能格式化（整数显示整数，小数2位）
        {{计算:表达式|.2f}}           保留2位小数
        {{计算:表达式|.0f}}           整数
        {{计算:表达式|.1%}}           百分比（值×100 后加 %）
        示例：
            {{计算:(本月-上月)/上月*100}}            环比提升%
            {{计算:利润/销售额|.2%}}                 利润率
            {{计算:A.销售额.max - A.销售额.min}}     极差

PPT 备注配置（每页备注区可写）：
    数据源=透视结果.xlsx
    区块=按地区汇总                # 声明本页默认区块，占位符可省略前缀
    别名.别名名=表达式              # 别名声明，正文用 {{别名名}} 引用
    形状名=区块表达式               # 图表/表格数据源映射（方案C）
                                    形状名需与 PPT 中的形状名称一致
                                    如：图表1=按地区汇总|销售额,利润
    别名.别名名=表达式             # 声明别名，正文用 {{别名名}} 引用，避免长表达式堆在正文
                                  # 表达式可以是文本占位符表达式或 计算: 表达式
    示例：
        别名.华东销售额=按地区汇总.总销售额.华东
        别名.利润率=计算:利润/销售额|.2%
        别名.环比=计算:(本月-上月)/上月*100|.2f
    正文写：华东销售额: {{华东销售额}}，利润率: {{利润率}}
"""
import os
import re
from typing import Dict, Optional, Tuple, List

import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from src.safe_math import evaluate_numeric_expression


# 占位符正则：{{...}}
_PLACEHOLDER_RE = re.compile(r"\{\{([^{}]+)\}\}")
# 图表占位符正则：{{图表:xxx}}
_CHART_PLACEHOLDER_RE = re.compile(r"\{\{图表[:：]([^{}]+)\}\}")
# 图片占位符正则：{{图片:xxx}}
_IMAGE_PLACEHOLDER_RE = re.compile(r"\{\{图片[:：]([^{}]+)\}\}")
# 表格占位符正则：{{表格:xxx}}
_TABLE_PLACEHOLDER_RE = re.compile(r"\{\{表格[:：]([^{}]+)\}\}")

# 缺失标注样式（黄底高亮 [缺失:...]）
_MISSING_HIGHLIGHT = RGBColor(0xFF, 0xFF, 0x00)   # 黄色背景
_MISSING_FONTCOLOR = RGBColor(0xC0, 0x00, 0x00)  # 深红字体


def _mark_missing_text(run, expr: str):
    """将一个 run 标注为缺失：文本=[缺失:expr]，红字加粗。
    文本框黄底由调用方按需设置（避免过度影响混合内容框）。
    """
    run.text = f"[缺失:{expr}]"
    try:
        run.font.color.rgb = _MISSING_FONTCOLOR
        run.font.bold = True
    except Exception:
        pass


def _mark_missing_cell(cell, block_name: str):
    """将表格单元格标注为缺失：文本=[缺失:block_name]，黄底红字加粗。"""
    cell.text = f"[缺失:{block_name}]"
    try:
        # 黄色填充背景
        cell.fill.solid()
        cell.fill.fore_color.rgb = _MISSING_HIGHLIGHT
        # 红字加粗
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = _MISSING_FONTCOLOR
                run.font.bold = True
    except Exception:
        pass


def load_pivot_results(pivot_data_path: str) -> Dict[str, pd.DataFrame]:
    """加载透视结果 xlsx，返回 {区块名或sheet名: DataFrame}。

    一个 sheet 中可能包含多个区块（以区块标题行分隔），每个区块独立解析。
    - 区块标题行：只有第1列有值（区块名），其余列为空
    - 区块标题行的下一行是表头，再下面是数据行

    索引优先级：
    1. 区块名（精确匹配，跨所有 sheet 查找）
    2. sheet 名（指向该 sheet 的第一个区块，兼容旧用法）
    3. "Sheet名.区块名"（精确指定 sheet 内的区块，新增）

    每个 DataFrame 的 attrs 会携带百分比列信息：
        df.attrs["pct_columns"] = set[str]   # 数字格式含 % 的列名集合
    下游 _write_chart_data 优先使用此元信息识别百分比列，避免靠列名猜测。
    """
    if not os.path.exists(pivot_data_path):
        raise FileNotFoundError(f"透视结果文件不存在: {pivot_data_path}")

    wb = openpyxl.load_workbook(pivot_data_path, data_only=True)
    result = {}
    # 记录每个 sheet 的第一个区块名，用于 sheet 名别名
    sheet_first_block = {}
    # 记录 (sheet_name, block_name) → DataFrame，用于 "Sheet名.区块名" 查找
    sheet_block_map = {}

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        # 扫描所有行，按区块标题行分段
        # 区块标题行：只有第1列有值，其余列为空
        block_starts = []  # [(title_row_idx, block_name), ...]
        for i, row in enumerate(rows):
            non_empty = [c for c in row if c is not None and str(c).strip()]
            if len(non_empty) == 1:
                # 只有第1列有值 → 区块标题行
                block_name = str(rows[i][0]).strip()
                block_starts.append((i, block_name))

        if not block_starts:
            # 没有区块标题行：整个 sheet 当作一个区块，sheet 名作为区块名
            block_starts.append((0, sheet_name))

        first_block_in_sheet = True
        for idx, (title_row_idx, block_name) in enumerate(block_starts):
            # 区块标题行存在时，表头在下一行；否则表头在当前行
            if block_name == sheet_name and not _is_title_row(rows, title_row_idx):
                header_row_idx = title_row_idx
            else:
                header_row_idx = title_row_idx + 1

            # 数据行范围：表头行+1 到 下一个区块标题行-1（或 sheet 末尾）
            data_end = block_starts[idx + 1][0] if idx + 1 < len(block_starts) else len(rows)
            data_rows = rows[header_row_idx + 1:data_end]
            # 过滤空行
            data_rows = [r for r in data_rows if any(c is not None and str(c).strip() for c in r)]
            if not data_rows:
                continue
            if header_row_idx >= len(rows):
                continue

            headers = [str(c).strip() if c is not None else f"col_{i}" for i, c in enumerate(rows[header_row_idx])]
            df = pd.DataFrame(data_rows, columns=headers)

            # 读取数据行的单元格数字格式，识别百分比列（格式串含 %）
            # 表头单元格是文本格式无意义，必须读数据行
            # 取第一个数据行的格式作为整列格式（Excel 列格式通常一致）
            pct_cols = set()
            try:
                data_row_idx_1based = header_row_idx + 2  # 表头+1 = 第一个数据行（1-based）
                if data_row_idx_1based <= ws.max_row:
                    headers_row = list(ws[header_row_idx + 1])
                    data_cells = list(ws[data_row_idx_1based])
                    for h_cell, d_cell in zip(headers_row, data_cells):
                        if h_cell.value is None:
                            continue
                        col_name = str(h_cell.value).strip()
                        fmt = str(d_cell.number_format or "")
                        if "%" in fmt:
                            pct_cols.add(col_name)
            except Exception:
                pass
            df.attrs["pct_columns"] = pct_cols

            # 用区块名作为 key（后出现的同名区块覆盖先出现的，与 pivot 输出语义一致）
            result[block_name] = df

            # 记录 (sheet_name, block_name) → DataFrame
            sheet_block_map[f"{sheet_name}.{block_name}"] = df

            # 记录每个 sheet 的第一个区块名
            if first_block_in_sheet:
                sheet_first_block[sheet_name] = block_name
                first_block_in_sheet = False

    # 用 sheet 名作为别名，指向该 sheet 的第一个区块（兼容旧用法）
    for sheet_name, first_block in sheet_first_block.items():
        if sheet_name not in result:
            result[sheet_name] = result.get(first_block)

    # 合并 "Sheet名.区块名" 形式的 key
    result.update(sheet_block_map)

    wb.close()
    return result


def _is_title_row(rows, row_idx):
    """判断指定行是否是区块标题行（只有第1列有值）"""
    if row_idx >= len(rows):
        return False
    row = rows[row_idx]
    non_empty = [c for c in row if c is not None and str(c).strip()]
    return len(non_empty) == 1


def _parse_value(value, default=""):
    """将单元格值转换为字符串，数值保留合理精度"""
    if value is None:
        return default
    if isinstance(value, float):
        if value == int(value):
            return str(int(value))
        return f"{value:.2f}"
    return str(value)


def _aggregate_column(df: pd.DataFrame, col: str, agg: str = "sum") -> Optional[float]:
    """对指定列做聚合"""
    if col not in df.columns:
        return None
    series = df[col]
    # 尝试转为数值
    numeric_series = pd.to_numeric(series, errors="coerce")
    if numeric_series.isna().all():
        return None
    agg = agg.lower().strip()
    if agg == "sum":
        return float(numeric_series.sum())
    elif agg in ("avg", "mean"):
        return float(numeric_series.mean())
    elif agg == "max":
        return float(numeric_series.max())
    elif agg == "min":
        return float(numeric_series.min())
    elif agg == "count":
        return int(numeric_series.count())
    return float(numeric_series.sum())


def _lookup_block(pivot_data: Dict[str, pd.DataFrame], block_expr: str) -> Optional[pd.DataFrame]:
    """按区块表达式查找 DataFrame，支持三种形式：
    1. "区块名"          → 跨 sheet 查找（兼容旧用法）
    2. "Sheet名.区块名"  → 精确指定 sheet 内的区块
    3. "Sheet名"         → 该 sheet 第一个区块（兼容旧用法）
    """
    if not block_expr:
        return None
    # 精确匹配（包括 "Sheet名.区块名" 形式）
    if block_expr in pivot_data:
        return pivot_data[block_expr]
    return None


def _parse_block_and_cols(expr: str) -> Tuple[str, Optional[List[str]]]:
    """解析占位符参数：'区块名|列1,列2,列3' → ('区块名', ['列1','列2','列3'])

    支持的格式：
        区块名                → (区块名, None)  全量列
        区块名|列1,列2        → (区块名, [列1,列2])  指定列
        Sheet名.区块名|列1    → (Sheet名.区块名, [列1])

    列名会去除首尾空白。无 | 时返回 (expr, None) 保持兼容。
    """
    expr = expr.strip()
    if "|" not in expr:
        return expr, None
    block_part, cols_part = expr.split("|", 1)
    block_part = block_part.strip()
    cols = [c.strip() for c in cols_part.split(",") if c.strip()]
    if not cols:
        return block_part, None
    return block_part, cols


def _extract_xy_pair_flag(expr: str) -> Tuple[str, bool]:
    """从图表表达式中提取 |xy 标志，返回 (清理后表达式, xy_pair_flag)

    用法：{{图表:区块名|xy}} 或 {{图表:区块名|xy|列1,列2}}
    xy 标志表示散点图按 (X1,Y1,X2,Y2,...) 列顺序配对，每对形成一个独立系列。
    """
    expr = expr.strip()
    # 匹配 |xy 或 ｜xy（全角竖线），可能在末尾或中间
    for sep in ("|xy", "｜xy", "|XY", "｜XY"):
        if sep in expr:
            cleaned = expr.replace(sep, "").replace("||", "|").strip(" |｜")
            return cleaned, True
    return expr, False


def _extract_transpose_flag(expr: str) -> Tuple[str, bool]:
    """从图表表达式中提取 |t 标志，返回 (清理后表达式, transpose_flag)

    用法：{{图表:区块名|t}} 或 {{图表:区块名|列1,列2|t}}
    |t 标志表示行列转置：行维度值作为系列名，列名作为 X 轴类别。
    对应 PowerPoint 中"切换行列"按钮的功能。
    """
    expr = expr.strip()
    for sep in ("|t", "｜t", "|T", "｜T"):
        if sep in expr:
            cleaned = expr.replace(sep, "").replace("||", "|").strip(" |｜")
            return cleaned, True
    return expr, False


def _filter_df_columns(df: pd.DataFrame, cols: Optional[List[str]], min_keep_cols: int = 1) -> pd.DataFrame:
    """按指定列筛选 DataFrame，至少保留前 min_keep_cols 列 + 指定列。

    无 cols 或 cols 为 None 时返回原 df。
    指定列不存在时跳过。
    min_keep_cols 用于多级分类 X 轴场景：模板有 N 级分类时，前 N 列为分类层级列，
    必须保留即使用户只写了数据列；用户显式写了分类列时自动去重不重复添加。
    保留原 df 的 attrs（含 pct_columns 等元信息）。
    """
    if cols is None or df.empty:
        return df
    # 至少保留前 min_keep_cols 列（分类层级列），用于多级分类 X 轴
    n = max(1, min_keep_cols)
    n = min(n, len(df.columns))
    keep = list(df.columns[:n])
    keep_set = set(keep)
    for c in cols:
        if c in df.columns and c not in keep_set:
            keep.append(c)
            keep_set.add(c)
    if not keep:
        return df
    filtered = df[keep]
    # df[keep] 切片会丢失 attrs，手动保留百分比列等元信息
    try:
        filtered.attrs = dict(df.attrs)
    except Exception:
        pass
    return filtered


def _smart_format_number(value) -> str:
    """智能格式化数字：整数显示整数，小数保留2位"""
    if isinstance(value, bool):
        return str(value)
    if isinstance(value, (int, float)):
        if isinstance(value, float) and value == int(value):
            return str(int(value))
        if isinstance(value, int):
            return str(value)
        return f"{value:.2f}"
    return str(value)


# 已知的格式化字符串集合（计算占位符和文本占位符共用）
# 同时支持半角 | 和全角 ｜ 作为格式分隔符
_KNOWN_FMTS = {
    # 小数位格式
    ".0f", ".1f", ".2f", ".3f", ".4f", ".5f", ".6f",
    # 整数格式
    "int", "d",
    # 百分比格式（自动 ×100 加 %）
    ".0%", ".1%", ".2%", ".3%", ".4%",
    # 科学计数法
    ".0e", ".1e", ".2e", ".3e", ".0E", ".1E", ".2E", ".3E",
    # 千分位
    ",.0f", ",.1f", ",.2f", ",.3f", ",.0%", ",.1%", ",.2%", ",.3%",
}


def _split_expr_and_fmt(expr: str) -> Tuple[str, Optional[str]]:
    """从表达式中分离格式后缀，返回 (表达式, 格式或None)

    仅当最后一个 | 或 ｜ 后的内容命中 _KNOWN_FMTS 时才视为格式串，
    否则保持原样（| 视为表达式的一部分）。
    同时支持半角 | (U+007C) 和全角 ｜ (U+FF5C)。
    """
    last_pipe_idx = max(expr.rfind("|"), expr.rfind("｜"))
    if last_pipe_idx == -1:
        return expr, None
    fmt_candidate = expr[last_pipe_idx + 1:].strip()
    if fmt_candidate in _KNOWN_FMTS:
        return expr[:last_pipe_idx].strip(), fmt_candidate
    return expr, None


def _format_calc_result(value, fmt: str) -> str:
    """按格式化字符串格式化计算结果

    支持的格式：
        .0f ~ .6f               固定小数位（.0f/.1f/.2f/.3f/.4f/.5f/.6f）
        int / d                 整数
        .0% ~ .4%               百分比（值×100 后加 %）
        .0e ~ .3e / .0E ~ .3E   科学计数法
        ,.0f ~ ,.3f             千分位小数
        ,.0% ~ ,.3%             千分位百分比
    """
    fmt = fmt.strip()
    pct_formats = {".0%", ".1%", ".2%", ".3%", ".4%"}
    if fmt in pct_formats:
        # Python format spec ".2%" 本身会自动乘 100 并加 %
        try:
            return format(float(value), fmt)
        except Exception:
            pass
    if fmt in ("int", "d"):
        return str(int(value))
    try:
        return format(value, fmt)
    except Exception:
        return _smart_format_number(value)


def _resolve_calc_expr(expr: str, pivot_data: Dict[str, pd.DataFrame],
                       default_block: Optional[str]) -> str:
    """解析 {{计算:表达式}} 占位符，返回计算结果字符串

    语法：
        {{计算:表达式}}              默认智能格式化（整数显示整数，小数2位）
        {{计算:表达式|.2f}}          保留2位小数
        {{计算:表达式|.0f}}          整数
        {{计算:表达式|.1%}}          百分比（值×100 后加 %，保留1位小数）

    表达式中的标识符沿用现有文本占位符语法（至少含2段以 . 分隔）：
        区块名.列名                  数值列求和 / 非数值列首行
        区块名.列名.行值             精确取某行某列
        区块名.列名.聚合             显式聚合（sum/avg/max/min/count）
        Sheet名.区块名.列名          精确指定 sheet 内区块
        标量.标量名                  取标量值
        列名                         使用 default_block（仅当 default_block 存在）

    支持的运算符：+ - * / () 以及小数和整数常量。
    找不到值的标识符按 0 处理；除零等异常返回空字符串。

    示例：
        {{计算:(本月-上月)/上月*100}}               默认格式
        {{计算:(本月.华东-上月.华东)/上月.华东*100|.2f}}
        {{计算:利润/销售额|.2%}}                    百分比格式
        {{计算:A.销售额.max - A.销售额.min|.0f}}    极差取整
    """
    # 分离表达式与格式（复用模块级 _split_expr_and_fmt）
    expr, fmt = _split_expr_and_fmt(expr)

    # 标识符正则：1段或多段以 . 分隔，第一段必须含字母/汉字/下划线（排除 3.14 这类纯数字）
    # 单段如 "本月"（依赖 default_block），多段如 "区块.列"、"Sheet.区块.列.行值"
    # 字符类含空格以支持含空格的字段名（如 "总 销售额"），匹配后 rstrip 去除运算符前的空格
    # 每段后可选跟 (...) 形式的括号内容，支持含括号的字段名（如 "销售额(万元)"、"率(%)"）
    # 括号内不支持嵌套，仅匹配单层 ()
    # 字符类含 % 以支持字段名直接含百分号（如 "增长%"、"利润率%"）
    ident_pattern = re.compile(
        r'[\u4e00-\u9fa5a-zA-Z_][\u4e00-\u9fa5\w %]*'
        r'(?:\.[\u4e00-\u9fa5\w %]+(?:\([^)]*\))?)*'
    )

    missing_any = False

    def _replace_ident(m):
        nonlocal missing_any
        ident = m.group(0).rstrip()  # 去除末尾空格（运算符前的空格被贪婪匹配吞入）
        value = _resolve_text_placeholder(ident, pivot_data, default_block)
        if not value:
            missing_any = True
            return "(0)"
        try:
            return f"({float(value)})"
        except (ValueError, TypeError):
            missing_any = True
            return "(0)"

    safe_expr = ident_pattern.sub(_replace_ident, expr)

    try:
        result = evaluate_numeric_expression(safe_expr)
    except Exception:
        return ""

    if fmt:
        return _format_calc_result(result, fmt)
    return _smart_format_number(result)


def _resolve_text_placeholder(expr: str, pivot_data: Dict[str, pd.DataFrame],
                              default_block: Optional[str] = None,
                              raw: bool = False) -> str:
    """解析文本占位符表达式，返回替换值字符串

    支持的格式：
        区块名.列名
        区块名.列名.行值
        区块名.行数
        区块名.列名.聚合(sum/avg/max/min/count)
        标量.标量名
        Sheet名.区块名.列名              （精确指定 sheet 内区块，新增）
        Sheet名.区块名.列名.行值         （新增）
        Sheet名.区块名.列名.聚合         （新增）
        列名                             （使用 default_block）
        列名.行值                        （使用 default_block）

    查找优先级：先尝试 "Sheet名.区块名" 双段精确匹配，
    再回退到 "区块名" 单段跨 sheet 匹配。

    raw=True 时返回原始 Python 值（int/float/str），不经过 _parse_value 截断，
    用于带格式后缀的文本占位符（避免精度丢失，如 0.3765 被截成 "0.38"）。
    找不到值时仍返回空字符串 ""。
    """
    expr = expr.strip()
    parts = expr.split(".")

    # 标量特殊处理
    if parts[0] == "标量":
        if len(parts) < 2:
            return ""
        scalar_name = parts[1]
        # 标量存放在无行维度的结果里（单行），遍历所有 sheet 查找
        for sheet_name, df in pivot_data.items():
            if scalar_name in df.columns and len(df) == 1:
                v = df[scalar_name].iloc[0]
                return v if raw else _parse_value(v)
        return ""

    # 尝试 "Sheet名.区块名" 双段形式（至少3段：Sheet.区块.列）
    # 先用最长前缀匹配，找出 Sheet.区块 的边界
    block = None
    col = None
    row_val = None
    agg = None

    if len(parts) >= 3:
        # 尝试前两段作为 "Sheet名.区块名"
        candidate_two = f"{parts[0]}.{parts[1]}"
        if candidate_two in pivot_data:
            block = candidate_two
            col = parts[2]
            if len(parts) >= 4:
                fourth = parts[3]
                if fourth in ("sum", "avg", "mean", "max", "min", "count"):
                    agg = fourth
                    row_val = None
                else:
                    row_val = fourth
                    agg = parts[4] if len(parts) > 4 else None

    # 回退到旧逻辑：单段区块名
    if block is None:
        if len(parts) == 1:
            # 只有列名，依赖 default_block
            if not default_block:
                return ""
            col = parts[0]
            block = default_block
            row_val = None
            agg = None
        elif len(parts) == 2:
            # 区块名.列名 或 列名.行值/聚合
            if parts[0] in pivot_data:
                block, col = parts[0], parts[1]
                row_val = None
                agg = None
            elif default_block and parts[1] in ("sum", "avg", "mean", "max", "min", "count"):
                block = default_block
                col = parts[0]
                row_val = None
                agg = parts[1]
            elif default_block:
                block = default_block
                col = parts[0]
                row_val = parts[1]
                agg = None
            else:
                return ""
        elif len(parts) >= 3:
            block = parts[0]
            col = parts[1]
            third = parts[2]
            # 判断第三个是聚合还是行值
            if third in ("sum", "avg", "mean", "max", "min", "count"):
                agg = third
                row_val = None
            else:
                row_val = third
                agg = parts[3] if len(parts) > 3 else None
        else:
            return ""

    df = _lookup_block(pivot_data, block)
    if df is None:
        return ""

    # 特殊列名：行数
    if col == "行数":
        return len(df) if raw else str(len(df))

    if col not in df.columns:
        return ""

    # 按行值定位
    if row_val:
        # 第一列作为行维度
        row_dim_col = df.columns[0]
        matched = df[df[row_dim_col].astype(str) == row_val]
        if len(matched) == 0:
            return ""
        value = matched[col].iloc[0]
        return value if raw else _parse_value(value)

    # 聚合
    if agg:
        result = _aggregate_column(df, col, agg)
        if result is None:
            return ""
        if raw:
            return result
        if isinstance(result, float):
            if result == int(result):
                return str(int(result))
            return f"{result:.2f}"
        return str(result)

    # 默认：数值列求和，非数值列取首行
    numeric_series = pd.to_numeric(df[col], errors="coerce")
    if not numeric_series.isna().all():
        v = float(numeric_series.sum())
        return v if raw else _parse_value(v)
    v = df[col].iloc[0]
    return v if raw else _parse_value(v)


def _parse_slide_notes(notes_text: str) -> Dict[str, str]:
    """解析幻灯片备注区的配置"""
    config = {}
    for line in notes_text.split("\n"):
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        if "=" in line:
            key, value = line.split("=", 1)
            config[key.strip()] = value.strip()
    return config


def _update_notes_with_status(slide, status_map: Dict[str, str]):
    """v2.54.25+ 将形状替换状态回写到备注区对应行末尾。

    备注区原行格式：形状名=区块表达式
    回写后格式：形状名=区块表达式  # 成功(5行x3列)
    或：形状名=区块表达式  # 失败: 区块'XXX'未找到

    - 只更新通过 shape_block_map 声明的形状名行（即 status_map 中的 key）
    - 同一形状名多次出现时更新首个匹配行
    - 原行已有 # 状态标记时先清除旧标记再写新标记
    - 无 notes_slide 时自动创建
    """
    if not status_map:
        return
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return
    tf = notes_slide.notes_text_frame
    lines = tf.text.split("\n")
    new_lines = []
    updated_keys = set()
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.startswith("#") or "=" not in stripped:
            new_lines.append(line)
            continue
        key = stripped.split("=", 1)[0].strip()
        # 跳过已知前缀（区块/数据源/别名.）
        if key.startswith(("区块", "数据源", "别名.")):
            new_lines.append(line)
            continue
        if key in status_map and key not in updated_keys:
            # 去除原行末尾可能存在的旧状态标记
            base = re.sub(r"\s*#.*$", "", stripped)
            status = status_map[key]
            new_lines.append(f"{base}  # {status}")
            updated_keys.add(key)
        else:
            new_lines.append(line)
    # 未在备注区找到对应形状名行的状态，追加到末尾
    for key, status in status_map.items():
        if key not in updated_keys:
            new_lines.append(f"{key}=?  # {status}")
    try:
        tf.text = "\n".join(new_lines)
    except Exception:
        pass


def _replace_in_text_frame(text_frame, pivot_data: Dict[str, pd.DataFrame],
                           default_block: Optional[str],
                           image_collector: Optional[List[str]] = None,
                           alias_map: Optional[Dict[str, str]] = None,
                           mark_missing: bool = True) -> int:
    """替换文本框中的占位符，返回替换次数。

    保留各 run 的原始格式（颜色、字号、加粗等），仅替换占位符所在的 run 文本，
    不影响同段落中其他 run 的格式。

    如果 image_collector 不为 None，{{图片:...}} 会被收集到列表中并从文本中移除。
    alias_map: 备注区声明的别名映射（{别名: 表达式}），正文 {{别名}} 会被展开为对应表达式再解析。
    """
    replace_count = 0
    for para in text_frame.paragraphs:
        if not para.runs:
            continue

        runs = para.runs
        # 构建完整文本和各 run 的位置映射
        full_text = ""
        run_positions = []  # [(start, end) for each run]
        for run in runs:
            start = len(full_text)
            full_text += run.text
            end = len(full_text)
            run_positions.append((start, end))

        if not _PLACEHOLDER_RE.search(full_text):
            continue

        # 查找所有占位符匹配，逆序处理避免位置偏移
        matches = list(_PLACEHOLDER_RE.finditer(full_text))
        for m in reversed(matches):
            expr = m.group(1).strip()
            match_start = m.start()
            match_end = m.end()

            # 解析占位符，获取替换值
            replacement = None  # None = 保留原文（不替换）
            should_count = False
            is_missing = False  # 数据缺失（非图片正常移除）

            if expr.startswith("图表") or expr.startswith("图表:"):
                continue  # 保留原文
            elif (expr.startswith("图片") or expr.startswith("图片:")) and image_collector is not None:
                image_collector.append(expr[2:].lstrip(":：").strip())
                replacement = ""
                should_count = True
            elif expr.startswith("计算:") or expr.startswith("计算："):
                calc_expr = expr[3:].lstrip(":：").strip()
                value = _resolve_calc_expr(calc_expr, pivot_data, default_block)
                replacement = value or ""
                should_count = bool(value)
                is_missing = not bool(value)
            elif alias_map and expr in alias_map:
                alias_expr = alias_map[expr]
                # v2.54.26+ 别名表达式支持格式后缀（备注区声明时可能写 "区块.列|.2f"）
                # 注意：计算占位符（计算:...）的格式串由 _resolve_calc_expr 内部处理，这里不提前分离
                if alias_expr.startswith("计算:") or alias_expr.startswith("计算："):
                    calc_expr = alias_expr[3:].lstrip(":：").strip()
                    value = _resolve_calc_expr(calc_expr, pivot_data, default_block)
                else:
                    alias_expr_clean, alias_fmt = _split_expr_and_fmt(alias_expr)
                    if alias_fmt:
                        raw_value = _resolve_text_placeholder(alias_expr_clean, pivot_data, default_block, raw=True)
                        if raw_value == "" or raw_value is None:
                            value = ""
                        else:
                            try:
                                value = _format_calc_result(float(raw_value), alias_fmt)
                            except (ValueError, TypeError):
                                value = str(raw_value)
                    else:
                        value = _resolve_text_placeholder(alias_expr, pivot_data, default_block)
                replacement = value or ""
                should_count = bool(value)
                is_missing = not bool(value)
            else:
                # 文本占位符：支持可选格式后缀 {{区块.列|格式}}
                text_expr, text_fmt = _split_expr_and_fmt(expr)
                if text_fmt:
                    raw_value = _resolve_text_placeholder(text_expr, pivot_data, default_block, raw=True)
                    if raw_value == "" or raw_value is None:
                        replacement = ""
                        should_count = False
                        is_missing = True
                    else:
                        try:
                            value = _format_calc_result(float(raw_value), text_fmt)
                        except (ValueError, TypeError):
                            value = str(raw_value)
                        replacement = value
                        should_count = bool(value)
                        is_missing = not bool(value)
                else:
                    value = _resolve_text_placeholder(text_expr, pivot_data, default_block)
                    replacement = value or ""
                    should_count = bool(value)
                    is_missing = not bool(value)

            # 缺失标注：未替换的占位符标为 [缺失:expr] 红字加粗
            mark_this = False
            if is_missing and mark_missing:
                replacement = f"[缺失:{expr}]"
                mark_this = True

            if should_count:
                replace_count += 1

            # 定位占位符跨越的 run 范围
            first_run_idx = None
            last_run_idx = None
            for i, (rs, re_) in enumerate(run_positions):
                if first_run_idx is None and rs <= match_start < re_:
                    first_run_idx = i
                if rs < match_end <= re_:
                    last_run_idx = i
                    break

            if first_run_idx is None:
                continue
            if last_run_idx is None:
                last_run_idx = len(runs) - 1

            # 替换第一个 run 中的占位符部分为替换值
            rs, re_ = run_positions[first_run_idx]
            run = runs[first_run_idx]
            local_start = match_start - rs
            local_end = min(match_end, re_) - rs
            run.text = run.text[:local_start] + replacement + run.text[local_end:]

            # 缺失标注：给替换 run 设红字加粗
            if mark_this:
                try:
                    run.font.color.rgb = _MISSING_FONTCOLOR
                    run.font.bold = True
                except Exception:
                    pass

            # 清除后续 run 中的占位符残留部分
            for i in range(first_run_idx + 1, last_run_idx + 1):
                rs_i, re_i = run_positions[i]
                run_i = runs[i]
                local_end_i = min(match_end, re_i) - rs_i
                run_i.text = run_i.text[local_end_i:]

    return replace_count


def _replace_chart_data(slide, pivot_data: Dict[str, pd.DataFrame],
                        default_block: Optional[str],
                        alias_map: Optional[Dict[str, str]] = None,
                        shape_block_map: Optional[Dict[str, str]] = None,
                        mark_missing: bool = True,
                        status_map: Optional[Dict[str, str]] = None) -> int:
    """替换幻灯片中的图表数据，返回替换次数

    - 图表数据源占位符 {{图表:xxx}} 从形状名称/替代文字读取（不污染图表标题）
    - 形状名称/替代文字无占位符时，回退到 shape_block_map（备注区声明 形状名=区块名）
    - 图表标题支持文本占位符（{{区块.列}}、{{计算:...}}、别名、格式后缀等），
      与正文文本框共用替换逻辑
    - status_map: 若非 None，按形状名记录替换状态（成功/失败原因），供备注区回写
    """
    replace_count = 0
    for shape in slide.shapes:
        if not shape.has_chart:
            continue

        chart = shape.chart

        target_expr = None
        try:
            for attr in ("name", "alternative_text"):
                val = (getattr(shape, attr, None) or "")
                m = _CHART_PLACEHOLDER_RE.search(val)
                if m:
                    target_expr = m.group(1).strip()
                    break
        except Exception:
            pass

        # 形状名在备注区声明了映射（方案C：备注区写 "形状名=区块表达式"）
        if not target_expr and shape_block_map:
            try:
                mapped = shape_block_map.get(shape.name)
                if mapped:
                    target_expr = mapped
            except Exception:
                pass

        # 兼容：如果形状名称没有占位符，再回退到图表标题读取（旧模板兼容）
        if not target_expr:
            try:
                if chart.has_title and chart.chart_title.has_text_frame:
                    title_text = chart.chart_title.text_frame.text
                    m = _CHART_PLACEHOLDER_RE.search(title_text)
                    if m:
                        target_expr = m.group(1).strip()
            except Exception:
                pass

        # 替换图表标题中的文本占位符（与正文文本框共用逻辑）
        try:
            if chart.has_title and chart.chart_title.has_text_frame:
                title_count = _replace_in_text_frame(
                    chart.chart_title.text_frame, pivot_data, default_block,
                    image_collector=None, alias_map=alias_map, mark_missing=mark_missing
                )
                if title_count > 0:
                    replace_count += title_count
                    print(f"    [OK] 图表标题文本替换: {title_count} 处")
        except Exception:
            pass

        if not target_expr:
            continue

        # 支持多区块：用 ; 分隔多个区块表达式（仅散点图有意义，每区块形成一个系列）
        # 语法示例：SheetA.区块A|xy|X,Y ; SheetB.区块B|xy|X,Y ; 区块C
        sub_exprs = [s.strip() for s in target_expr.split(";") if s.strip()]

        if len(sub_exprs) <= 1:
            # 单区块：维持原逻辑
            sub_expr = sub_exprs[0] if sub_exprs else target_expr
            sub_expr, transpose = _extract_transpose_flag(sub_expr)
            sub_expr, xy_pair = _extract_xy_pair_flag(sub_expr)
            target_block, cols = _parse_block_and_cols(sub_expr)
            df = _lookup_block(pivot_data, target_block)
            if df is None:
                print(f"    [警告] 图表数据区块 '{target_block}' 未在透视结果中找到")
                if status_map is not None:
                    status_map[shape.name] = f"失败: 区块'{target_block}'未找到"
                if mark_missing:
                    try:
                        shape.name = f"{shape.name}[缺失:{target_block}]"
                    except Exception:
                        pass
                continue
            if df.empty:
                if status_map is not None:
                    status_map[shape.name] = f"失败: 区块'{target_block}'数据为空"
                if mark_missing:
                    try:
                        shape.name = f"{shape.name}[缺失:{target_block}]"
                    except Exception:
                        pass
                continue
            # v2.54.28+ 多级分类 X 轴：自动保留前 N 列分类层级列（N=模板层级数）
            # 用户 | 后只写数据列即可；若显式写了分类列则去重不重复添加
            template_level_count = _get_template_multi_level_count(chart)
            df = _filter_df_columns(df, cols, min_keep_cols=template_level_count if template_level_count > 0 else 1)

            try:
                _write_chart_data(chart, df, xy_pair=xy_pair, transpose=transpose)
                replace_count += 1
                cols_info = f", 仅列: {cols}" if cols else ""
                print(f"    [OK] 图表数据替换: {target_block} ({df.shape[0]}行 x {df.shape[1]}列{cols_info})")
                if status_map is not None:
                    status_map[shape.name] = f"成功({df.shape[0]}行x{df.shape[1]}列)"
                try:
                    if shape.name and _CHART_PLACEHOLDER_RE.search(shape.name):
                        shape.name = _CHART_PLACEHOLDER_RE.sub(target_block, shape.name)
                except Exception:
                    pass
            except Exception as e:
                print(f"    [警告] 图表数据替换失败 [{target_block}]: {e}")
                if status_map is not None:
                    status_map[shape.name] = f"失败: {e}"
        else:
            # 多区块：散点图专属，每区块解析后形成一个独立系列
            block_dfs = []  # [(df, xy_pair, cols, block_name), ...]
            for sub_expr in sub_exprs:
                sub_expr_clean, xy_pair = _extract_xy_pair_flag(sub_expr)
                target_block, cols = _parse_block_and_cols(sub_expr_clean)
                df = _lookup_block(pivot_data, target_block)
                if df is None:
                    print(f"    [警告] 图表数据区块 '{target_block}' 未在透视结果中找到")
                    continue
                if df.empty:
                    continue
                df = _filter_df_columns(df, cols)
                block_dfs.append((df, xy_pair, cols, target_block))

            if not block_dfs:
                print(f"    [警告] 多区块表达式无有效数据: {target_expr}")
                if status_map is not None:
                    status_map[shape.name] = f"失败: 多区块无有效数据({target_expr})"
                if mark_missing:
                    try:
                        shape.name = f"{shape.name}[缺失:{target_expr}]"
                    except Exception:
                        pass
                continue

            try:
                _write_chart_data_multi(chart, block_dfs)
                replace_count += 1
                blocks_info = ", ".join(f"{b[3]}({b[0].shape[0]}行)" for b in block_dfs)
                print(f"    [OK] 图表多区块替换: {blocks_info}")
                if status_map is not None:
                    status_map[shape.name] = f"成功({blocks_info})"
                # 清除形状名称中的占位符
                try:
                    if shape.name and _CHART_PLACEHOLDER_RE.search(shape.name):
                        shape.name = _CHART_PLACEHOLDER_RE.sub(
                            "+".join(b[3] for b in block_dfs), shape.name)
                except Exception:
                    pass
            except Exception as e:
                print(f"    [警告] 图表多区块替换失败: {e}")
                if status_map is not None:
                    status_map[shape.name] = f"失败: {e}"
    return replace_count


_PCT_KEYWORDS = ["占比", "pct", "百分比", "比例"]


def _fix_embedded_workbook_for_combo(chart, categories: list, series_data: dict):
    """v2.54.37+ 修复组合图（barChart+lineChart 等多 chart 类型共存）嵌入工作簿

    python-pptx 的 replace_data 在组合图上行为异常：
    - 嵌入工作簿表头行（第1行）写入的是数字索引（3/4）而非系列名
    - 分类列（A列）写入的是数字 0/1/2 而非文本类别
    - A1 单元格为空
    导致 PowerPoint 检测到 c:tx 引用（$B$1=3）与缓存（系列名）不匹配，报"链接不可用"或"nan/inf not supported"。

    本函数直接修改嵌入 xlsx 包，重写 sheet1.xml：
    - A1=空或"类别"，A2:A{n}=文本类别
    - B1/C1/...=系列名，B2:B{n}/C2:C{n}/...=数值

    Args:
        chart: pptx Chart 对象
        categories: X 轴类别列表（文本）
        series_data: {系列名: [数值列表]} 有序字典
    """
    import io
    import zipfile
    from pptx.oxml.ns import qn
    from lxml import etree

    try:
        chart_space = chart._chartSpace
        # 检测是否为组合图（plotArea 下有多个不同类型的 chart 子节点）
        plot_area = chart_space.find('.//' + qn('c:plotArea'))
        if plot_area is None:
            return
        chart_types = []
        for child in plot_area:
            tag = etree.QName(child).localname
            if tag.endswith('Chart') and tag != 'plotArea':
                chart_types.append(tag)
        # 只有1种图表类型时不需要修复
        if len(set(chart_types)) <= 1:
            return

        # 找到嵌入工作簿
        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        ext_data = chart_space.find('.//' + qn('c:externalData'))
        if ext_data is None:
            return
        rid = ext_data.get('{%s}id' % ns_r)
        if not rid:
            return

        # 通过 part 获取嵌入工作簿
        chart_part = chart.part
        embed_part = chart_part.related_part(rid)
        if embed_part is None:
            return

        embed_bytes = embed_part.blob
        ns_s = 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'

        with zipfile.ZipFile(io.BytesIO(embed_bytes), 'r') as zin:
            file_data = {n: zin.read(n) for n in zin.namelist()}

        # 读取 sharedStrings（如果有）
        shared_strings = []
        if 'xl/sharedStrings.xml' in file_data:
            ss_root = etree.fromstring(file_data['xl/sharedStrings.xml'])
            for si in ss_root.findall('{%s}si' % ns_s):
                # 合并所有 t 节点文本
                texts = [t.text or '' for t in si.findall('.//{%s}t' % ns_s)]
                shared_strings.append(''.join(texts))

        # 重写 sheet1.xml
        sheet_root = etree.fromstring(file_data['xl/worksheets/sheet1.xml'])
        # 删除原有 sheetData
        old_sd = sheet_root.find('{%s}sheetData' % ns_s)
        if old_sd is not None:
            sheet_root.remove(old_sd)

        # 构建新 sheetData
        new_sd = etree.SubElement(sheet_root, '{%s}sheetData' % ns_s)
        n_rows = len(categories)
        n_series = len(series_data)

        # 准备字符串写入 sharedStrings（追加）
        new_strings = []
        str_idx_map = {}
        def get_str_idx(s):
            if s in str_idx_map:
                return str_idx_map[s]
            idx = len(shared_strings) + len(new_strings)
            new_strings.append(s)
            str_idx_map[s] = idx
            return idx

        # 第1行：A1=空，B1/C1/...=系列名
        row1 = etree.SubElement(new_sd, '{%s}row' % ns_s, r='1')
        # A1 留空（或写"类别"）
        # B1, C1, ... 系列名
        for i, sname in enumerate(series_data.keys()):
            col_letter = _col_index_to_letter(i + 2)  # B, C, D...
            c = etree.SubElement(row1, '{%s}c' % ns_s, r='%s1' % col_letter, t='s')
            v = etree.SubElement(c, '{%s}v' % ns_s)
            v.text = str(get_str_idx(sname))

        # 数据行
        for row_idx in range(n_rows):
            r = row_idx + 2
            row = etree.SubElement(new_sd, '{%s}row' % ns_s, r=str(r))
            # A列=类别（文本）
            cat_val = str(categories[row_idx])
            c_a = etree.SubElement(row, '{%s}c' % ns_s, r='A%d' % r, t='s')
            v_a = etree.SubElement(c_a, '{%s}v' % ns_s)
            v_a.text = str(get_str_idx(cat_val))
            # B, C, ... 数值
            for i, (sname, vals) in enumerate(series_data.items()):
                col_letter = _col_index_to_letter(i + 2)
                c = etree.SubElement(row, '{%s}c' % ns_s, r='%s%d' % (col_letter, r))
                v = etree.SubElement(c, '{%s}v' % ns_s)
                val = vals[row_idx] if row_idx < len(vals) else 0
                v.text = str(val)

        # 更新 dimension
        dim = sheet_root.find('{%s}dimension' % ns_s)
        max_col = _col_index_to_letter(n_series + 1)
        dim_ref = 'A1:%s%d' % (max_col, n_rows + 1)
        if dim is not None:
            dim.set('ref', dim_ref)
        else:
            dim = etree.Element('{%s}dimension' % ns_s, ref=dim_ref)
            sheet_root.insert(0, dim)

        # 更新 sharedStrings
        if new_strings:
            if 'xl/sharedStrings.xml' in file_data:
                ss_root = etree.fromstring(file_data['xl/sharedStrings.xml'])
            else:
                ss_root = etree.Element('{%s}sst' % ns_s,
                                        xmlns=ns_s,
                                        count=str(len(new_strings)),
                                        uniqueCount=str(len(new_strings)))
            # 更新 count/uniqueCount
            total_count = len(shared_strings) + len(new_strings)
            ss_root.set('count', str(total_count))
            ss_root.set('uniqueCount', str(total_count))
            for s in new_strings:
                si = etree.SubElement(ss_root, '{%s}si' % ns_s)
                t = etree.SubElement(si, '{%s}t' % ns_s)
                t.text = s
            file_data['xl/sharedStrings.xml'] = etree.tostring(ss_root, xml_declaration=True,
                                                               encoding='UTF-8', standalone=True)

        # 序列化 sheet1.xml
        file_data['xl/worksheets/sheet1.xml'] = etree.tostring(sheet_root, xml_declaration=True,
                                                                encoding='UTF-8', standalone=True)

        # 重写嵌入工作簿 zip
        out_buf = io.BytesIO()
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for name, data in file_data.items():
                zout.writestr(name, data)

        # 写回 part
        embed_part._blob = out_buf.getvalue()
        # 同时更新 rels 缓存
        try:
            embed_part.blob = out_buf.getvalue()
        except Exception:
            pass

        print(f"    [OK] 组合图嵌入工作簿已修复: {n_series}系列, {n_rows}行类别")
    except Exception as e:
        print(f"    [警告] 组合图嵌入工作簿修复失败: {e}")


def _trim_extra_series(chart, expected_count: int):
    """replace_data 后剪除模板遗留的多余系列（模板6系列→数据4系列时剪掉2个空系列）"""
    from pptx.oxml.ns import qn
    try:
        plot = chart.plots[0]
        ser_elements = plot._element.findall(qn("c:ser"))
        if len(ser_elements) > expected_count and expected_count > 0:
            for extra in ser_elements[expected_count:]:
                plot._element.remove(extra)
    except Exception:
        pass


def _ensure_vary_colors(chart):
    """确保图表 varyColors=true，使不同系列自动分配不同颜色"""
    try:
        _ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        plot = chart.plots[0]
        root = plot._element
        vc = root.find("{%s}varyColors" % _ns)
        if vc is None:
            from lxml import etree
            vc = etree.Element("{%s}varyColors" % _ns, {"val": "true"})
            root.insert(0, vc)
        else:
            if vc.get("val") != "1" and vc.get("val") != "true":
                vc.set("val", "1")
    except Exception:
        pass


def _fix_plot_by(chart):
    """修正图表 plotBy 方向为 col，防止模板行列方向与新数据不一致

    若已有 plotBy 元素（可能为 "row"）→ 改为 "col"；
    若不存在 plotBy → 在 plotArea 下创建并设为 "col"。
    """
    try:
        _ns = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        plot = chart.plots[0]
        # plotBy 在 plotArea 下，barChart/pieChart 等具体类型之上
        plot_area = plot._element.getparent()
        if plot_area is None:
            return
        for elem in plot_area.iter():
            if elem.tag == f"{{{_ns}}}plotBy":
                if elem.get("val") != "col":
                    elem.set("val", "col")
                return
        # 不存在 plotBy 时创建，插入到 plotArea 的第一个子元素前
        from lxml import etree
        pb = etree.Element(f"{{{_ns}}}plotBy", {"val": "col"})
        plot_area.insert(0, pb)
    except Exception:
        pass


def _is_pct_column(col_name: str, values: list) -> bool:
    """判断列是否为百分比数据：列名含占比关键词且值域在 0~1 之间"""
    if not col_name:
        return False
    name = str(col_name)
    if not any(kw in name for kw in _PCT_KEYWORDS):
        return False
    nums = [float(v) for v in values if isinstance(v, (int, float))]
    if not nums:
        return False
    return all(0 <= v <= 1 for v in nums)


def _write_chart_data_multi(chart, block_dfs: list):
    """多区块写入图表（散点图专属，每区块形成一个独立系列）

    Args:
        block_dfs: [(df, xy_pair, cols, block_name), ...]
            - df: DataFrame
            - xy_pair: 是否为独立 xy 对模式
            - cols: 指定列（可选）
            - block_name: 区块名（用于系列名）
    """
    from pptx.chart.data import XyChartData, CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    try:
        is_scatter = chart.chart_type == XL_CHART_TYPE.XY_SCATTER
    except Exception:
        is_scatter = False

    if is_scatter:
        # 散点图：每区块形成一个或多个系列
        chart_data = XyChartData()
        multi_series_count = 0
        for df, xy_pair, cols, block_name in block_dfs:
            if xy_pair and len(df.columns) >= 2 and len(df.columns) % 2 == 0:
                # xy 对模式：列按 (X1,Y1,X2,Y2,...) 配对
                for i in range(0, len(df.columns), 2):
                    y_col = df.columns[i + 1]
                    x_vals = pd.to_numeric(df.iloc[:, i], errors="coerce").fillna(0).tolist()
                    y_vals = pd.to_numeric(df.iloc[:, i + 1], errors="coerce").fillna(0).tolist()
                    series = chart_data.add_series(f"{block_name}.{y_col}")
                    for x_val, y_val in zip(x_vals, y_vals):
                        series.add_data_point(float(x_val), float(y_val))
                    multi_series_count += 1
            else:
                # 共享 X 模式：第一列作 X，其余列各成一个系列
                x_values = pd.to_numeric(df.iloc[:, 0], errors="coerce").fillna(0).tolist()
                for col_idx in range(1, len(df.columns)):
                    col_name = df.columns[col_idx]
                    y_values = pd.to_numeric(df.iloc[:, col_idx], errors="coerce").fillna(0).tolist()
                    series = chart_data.add_series(f"{block_name}.{col_name}")
                    for x_val, y_val in zip(x_values, y_values):
                        series.add_data_point(float(x_val), float(y_val))
                    multi_series_count += 1

        chart.replace_data(chart_data)
        _trim_extra_series(chart, multi_series_count)
        try:
            plot = chart.plots[0]
            if not plot.has_data_labels:
                plot.has_data_labels = True
            plot.data_labels.number_format = '#,##0.##'
            plot.data_labels.number_format_is_linked = False
        except Exception:
            pass
    else:
        # 非散点图：取第一个区块走原逻辑
        df, xy_pair, cols, block_name = block_dfs[0]
        _write_chart_data(chart, df, xy_pair=xy_pair)


def _get_template_multi_level_count(chart) -> int:
    """v2.54.27+ 检测模板图表本身是否含多级分类 X 轴，返回层级数

    通过检查图表 XML 的 c:cat 节点下是否存在 c:multiLvlStrRef 判断：
    - 存在 multiLvlStrRef → 返回其下 c:lvl 子节点数（多级分类层级数）
    - 不存在 → 返回 0（单级分类）

    根据模板图表自身结构判断，而非根据数据列类型自动猜测。
    """
    from pptx.oxml.ns import qn

    try:
        chart_space = chart._chartSpace
        cat_elem = chart_space.find('.//' + qn('c:cat'))
        if cat_elem is None:
            return 0
        multi_lvl = cat_elem.find(qn('c:multiLvlStrRef'))
        if multi_lvl is None:
            return 0
        # 统计 c:lvl 子节点数（每个 lvl 是一个层级）
        lvls = multi_lvl.findall('.//' + qn('c:lvl'))
        return len(lvls)
    except Exception:
        return 0


def _restore_multi_level_categories(chart, level_data: list, series_values: list = None,
                                     level_headers: list = None, series_header: str = None):
    """v2.54.27+ replace_data 后恢复多级分类 X 轴

    python-pptx 的 CategoryChartData 只支持单级分类，chart.replace_data() 会用
    单级 strRef 覆盖掉模板原有的 multiLvlStrRef。本函数手动重建多级分类 XML，
    并同步把所有层级列写入嵌入工作簿（v2.54.32+ 修复"选择数据"只显示2列的问题）。

    Args:
        chart: pptx Chart 对象
        level_data: 层级数据列表，最深层（子分类）在前，最浅层（父分类）在后
                    例如 [[产品A,产品B,...], [华东,华东,...]] → level 0 是产品（子），level 1 是地区（父）
        series_values: 第一个数据系列的数值列表（用于补写嵌入工作簿的数据列），
                       若为 None 则不修改嵌入工作簿的数据列
        level_headers: 各层级列表头名称（顺序同 level_data，最深层在前），
                       若为 None 则用"分类1/分类2/..."占位
        series_header: 数据系列列表头名称，若为 None 则用"数据"占位
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    try:
        chart_space = chart._chartSpace
        cat_elem = chart_space.find('.//' + qn('c:cat'))
        if cat_elem is None:
            return

        # 清除 replace_data 写入的单级 strRef
        str_ref = cat_elem.find(qn('c:strRef'))
        if str_ref is not None:
            cat_elem.remove(str_ref)

        # 重建 multiLvlStrRef
        multi_lvl = etree.SubElement(cat_elem, qn('c:multiLvlStrRef'))

        # 层级数 = level_data 长度；嵌入工作簿需要 N 列分类 + 1 列数据
        n_levels = len(level_data)
        n_rows = len(level_data[0]) if level_data else 0
        # 分类引用范围：A1:{第N列}{n_rows+1}（包含表头行）
        last_col_letter = _col_index_to_letter(n_levels)  # N 列分类
        f_elem = etree.SubElement(multi_lvl, qn('c:f'))
        f_elem.text = f'Sheet1!$A$1:${last_col_letter}${n_rows + 1}'

        cache = etree.SubElement(multi_lvl, qn('c:multiLvlStrCache'))

        pt_count = etree.SubElement(cache, qn('c:ptCount'))
        pt_count.set('val', str(n_rows))

        # 添加各层级（最深层在前，与 PowerPoint XML 规范一致）
        for level_vals in level_data:
            lvl = etree.SubElement(cache, qn('c:lvl'))
            for i, v in enumerate(level_vals):
                pt = etree.SubElement(lvl, qn('c:pt'))
                pt.set('idx', str(i))
                v_elem = etree.SubElement(pt, qn('c:v'))
                v_elem.text = str(v)

        # v2.54.32+ 同步把所有层级列写入嵌入工作簿
        # replace_data 只写了单级分类（最深层），父分类列缺失，导致"选择数据"只看到2列
        _rebuild_embedded_workbook_for_multi_level(chart, level_data, series_values,
                                                    level_headers, series_header)

        # v2.54.32+ 更新 c:val 的 c:f 引用范围（数据列从 B 变为第 N+1 列）
        if series_values is not None:
            val_elem = chart_space.find('.//' + qn('c:val'))
            if val_elem is not None:
                num_ref = val_elem.find(qn('c:numRef'))
                if num_ref is not None:
                    data_col_letter = _col_index_to_letter(n_levels + 1)  # 数据列 = N+1
                    f_val = num_ref.find(qn('c:f'))
                    if f_val is not None:
                        f_val.text = f'Sheet1!${data_col_letter}$2:${data_col_letter}${n_rows + 1}'

        # v2.54.35+ 更新 c:tx 的 c:f 引用范围（系列名从 B1 变为第 N+1 列第1行）
        # replace_data 把系列名写入 B1，但补齐父分类列后系列名实际在第 N+1 列
        # 不更新会导致 PowerPoint 检查 c:tx 引用（B1=频段）与缓存（平均下行速率）不匹配，报"链接不可用"
        tx_elem = chart_space.find('.//' + qn('c:ser') + '/' + qn('c:tx'))
        if tx_elem is not None:
            tx_str_ref = tx_elem.find(qn('c:strRef'))
            if tx_str_ref is not None:
                tx_f = tx_str_ref.find(qn('c:f'))
                if tx_f is not None and series_header is not None:
                    data_col_letter = _col_index_to_letter(n_levels + 1)
                    tx_f.text = f'Sheet1!${data_col_letter}$1'

        print(f"    [OK] 多级分类 X 轴已恢复: {n_levels} 级, {n_rows} 个类别")
    except Exception as e:
        print(f"    [警告] 多级分类恢复失败: {e}")


def _col_index_to_letter(col_idx: int) -> str:
    """将列序号（1-based）转为 Excel 列字母（1→A, 2→B, 27→AA）"""
    result = ''
    while col_idx > 0:
        col_idx, rem = divmod(col_idx - 1, 26)
        result = chr(65 + rem) + result
    return result


def _rebuild_embedded_workbook_for_multi_level(chart, level_data: list, series_values: list = None,
                                                level_headers: list = None, series_header: str = None):
    """v2.54.32+ 重建图表嵌入工作簿，补齐多级分类的父层级列

    python-pptx replace_data 只写单级分类列（最深层子分类），父分类列缺失。
    本函数直接修改嵌入的 xlsx 包，把所有层级列 + 数据列重新写入。

    Args:
        chart: pptx Chart 对象
        level_data: 层级数据列表（最深层在前）
        series_values: 数据系列数值列表（单系列场景），None 则保留原数据列
        level_headers: 各层级列表头名称（顺序同 level_data，最深层在前），
                       None 则用"分类1/分类2/..."占位
        series_header: 数据系列列表头名称，None 则用"数据"占位
    """
    import io
    import zipfile
    from lxml import etree
    from pptx.oxml.ns import qn

    try:
        chart_part = chart.part
        # 找嵌入工作簿关系（package 类型指向 embeddings/xxx.xlsx）
        embed_rel = None
        for rel_id, rel in chart_part.rels.items():
            if 'package' in rel.reltype:
                embed_rel = rel
                break
        if embed_rel is None:
            return

        # 读取嵌入工作簿 blob
        embed_part = embed_rel.target_part
        embed_blob = embed_part.blob

        n_levels = len(level_data)
        n_rows = len(level_data[0]) if level_data else 0
        # level_data: 最深层在前 → 列写入顺序倒过来（父分类在前，子分类在后）
        # 即 level_data[-1] 是最浅层（父）写入 A 列，level_data[0] 是最深层（子）写入第 N 列
        # 数据列写入第 N+1 列
        total_cols = n_levels + (1 if series_values is not None else 0)

        # 解析嵌入 xlsx
        with zipfile.ZipFile(io.BytesIO(embed_blob), 'r') as ez:
            # 读 sharedStrings（保留旧字符串，追加新字符串）
            shared = []
            if 'xl/sharedStrings.xml' in ez.namelist():
                ss_xml = ez.read('xl/sharedStrings.xml')
                root = etree.fromstring(ss_xml)
                ns = {'s': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
                for si in root.findall('s:si', ns):
                    texts = si.findall('.//s:t', ns)
                    shared.append(''.join(t.text or '' for t in texts))

            # 收集需要写入的所有字符串值（层级数据 + 系列名若需要）
            # 字符串索引映射：值 → 在 sharedStrings 中的索引
            str_index_map = {s: i for i, s in enumerate(shared)}

            def get_or_add_str(s):
                if s in str_index_map:
                    return str_index_map[s]
                idx = len(shared)
                shared.append(s)
                str_index_map[s] = idx
                return idx

            # 构建新的 sheetData
            # 表头行：第1列=父分类名...第N列=子分类名，第N+1列=数据列名
            # 由于无法轻易获取原始列名，用通用占位符（PowerPoint 显示数据时不依赖表头名）
            # 实际从 chart XML 的 numCache/strCache 读取数据，表头只占位
            ns = {'s': 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'}
            S_NS = 'http://schemas.openxmlformats.org/spreadsheetml/2006/main'

            rows_xml = []
            # 表头行（row 1）：各列占位
            # 列写入顺序：父分类（level_data[-1]）在 A 列，子分类（level_data[0]）在第 N 列
            # level_headers 顺序同 level_data（最深层在前），需倒序对应列序号
            header_cells = []
            for col_idx in range(1, total_cols + 1):
                col_letter = _col_index_to_letter(col_idx)
                if col_idx <= n_levels:
                    # 分类列表头：列 col_idx 对应 level_data[n_levels - col_idx]
                    level_idx = n_levels - col_idx
                    if level_headers and level_idx < len(level_headers):
                        header_name = str(level_headers[level_idx])
                    else:
                        header_name = f'分类{col_idx}'
                    str_idx = get_or_add_str(header_name)
                else:
                    # 数据列表头
                    header_name = str(series_header) if series_header else '数据'
                    str_idx = get_or_add_str(header_name)
                header_cells.append(f'<c r="{col_letter}1" t="s"><v>{str_idx}</v></c>')
            rows_xml.append(f'<row r="1">{"".join(header_cells)}</row>')

            # 数据行（row 2 ~ n_rows+1）
            for row_idx in range(n_rows):
                row_num = row_idx + 2
                cells = []
                for col_idx in range(1, total_cols + 1):
                    col_letter = _col_index_to_letter(col_idx)
                    if col_idx <= n_levels:
                        # level_data[0] 是最深层（子分类），应写入最后一列（第 N 列）
                        # level_data[-1] 是最浅层（父分类），应写入第 1 列
                        # 所以列 col_idx 对应 level_data[n_levels - col_idx]
                        level_idx = n_levels - col_idx
                        val = str(level_data[level_idx][row_idx])
                        str_idx = get_or_add_str(val)
                        cells.append(f'<c r="{col_letter}{row_num}" t="s"><v>{str_idx}</v></c>')
                    else:
                        # 数据列
                        val = series_values[row_idx] if series_values is not None and row_idx < len(series_values) else 0
                        cells.append(f'<c r="{col_letter}{row_num}"><v>{val}</v></c>')
                rows_xml.append(f'<row r="{row_num}">{"".join(cells)}</row>')

            new_sheetdata_xml = '<sheetData>' + ''.join(rows_xml) + '</sheetData>'
            last_col_letter = _col_index_to_letter(total_cols)
            new_dimension = f'A1:{last_col_letter}{n_rows + 1}'

            # 重建 sheet1.xml
            # OOXML worksheet 子节点顺序规范：
            # dimension → sheetViews → sheetFormatPr → cols → sheetData → ... → pageMargins
            sheet_xml = ez.read('xl/worksheets/sheet1.xml')
            sheet_root = etree.fromstring(sheet_xml)
            # 移除旧 sheetData 和 dimension
            for old_sd in sheet_root.findall('s:sheetData', ns):
                sheet_root.remove(old_sd)
            for old_dim in sheet_root.findall('s:dimension', ns):
                sheet_root.remove(old_dim)

            # 按规范顺序插入：dimension 在最前面（index 0）
            new_dim_elem = etree.Element('{%s}dimension' % S_NS)
            new_dim_elem.set('ref', new_dimension)
            sheet_root.insert(0, new_dim_elem)

            # sheetData 应在 cols 之后、pageMargins 之前
            # 找到 pageMargins（或其后第一个节点）的位置，在其前面插入 sheetData
            new_sd_elem = etree.fromstring(new_sheetdata_xml)
            insert_idx = len(sheet_root)  # 默认末尾
            for i, child in enumerate(sheet_root):
                child_tag = etree.QName(child).localname
                if child_tag in ('pageMargins', 'pageSetup', 'headerFooter',
                                 'rowBreaks', 'colBreaks', 'customProperties',
                                 'cellWatches', 'ignoredErrors', 'smartTags',
                                 'drawing', 'legacyDrawing', 'extLst'):
                    insert_idx = i
                    break
            sheet_root.insert(insert_idx, new_sd_elem)

            # 重新序列化（带 XML 声明）
            new_sheet_xml = etree.tostring(sheet_root, xml_declaration=True, encoding='UTF-8', standalone=True)

            # 重建 sharedStrings.xml
            new_ss_xml = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="%d" uniqueCount="%d">' % (len(shared), len(shared))
            for s in shared:
                # 转义特殊字符
                s_esc = s.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                new_ss_xml += f'<si><t>{s_esc}</t></si>'
            new_ss_xml += '</sst>'

        # 重建 xlsx zip 包
        new_embed_blob = io.BytesIO()
        with zipfile.ZipFile(io.BytesIO(embed_blob), 'r') as ez_in:
            with zipfile.ZipFile(new_embed_blob, 'w', zipfile.ZIP_DEFLATED) as ez_out:
                for item in ez_in.infolist():
                    if item.filename == 'xl/worksheets/sheet1.xml':
                        ez_out.writestr(item, new_sheet_xml)
                    elif item.filename == 'xl/sharedStrings.xml':
                        ez_out.writestr(item, new_ss_xml)
                    else:
                        ez_out.writestr(item, ez_in.read(item.filename))

        # 写回 chart part
        embed_part._blob = new_embed_blob.getvalue()

    except Exception as e:
        print(f"    [警告] 重建嵌入工作簿失败: {e}")


def _write_chart_data(chart, df: pd.DataFrame, xy_pair: bool = False, transpose: bool = False):
    """将 DataFrame 写入图表的内嵌 WorkBook，并同步数据标签格式

    按图表类型分支：
    - 散点图（XY_SCATTER）：用 XyChartData
        - xy_pair=False（默认）：第一列作共享 X，其余列各成一个 Y 系列
        - xy_pair=True：列按 (X1,Y1,X2,Y2,...) 顺序配对，每对形成一个独立系列
    - 其他图表：用 CategoryChartData，第一列作类别，其余列作系列
    """
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # 散点图分支：X 轴必须为数值，不能用字符串类别
    try:
        is_scatter = chart.chart_type == XL_CHART_TYPE.XY_SCATTER
    except Exception:
        is_scatter = False

    if is_scatter:
        chart_data = XyChartData()
        if xy_pair:
            if len(df.columns) >= 2 and len(df.columns) % 2 == 0:
                for i in range(0, len(df.columns), 2):
                    x_col = df.columns[i]
                    y_col = df.columns[i + 1]
                    x_vals = pd.to_numeric(df.iloc[:, i], errors="coerce").fillna(0).tolist()
                    y_vals = pd.to_numeric(df.iloc[:, i + 1], errors="coerce").fillna(0).tolist()
                    series = chart_data.add_series(f"{y_col}")
                    for x_val, y_val in zip(x_vals, y_vals):
                        series.add_data_point(float(x_val), float(y_val))
            else:
                print(f"    [警告] xy 对模式要求列数为偶数(>=2)，当前 {len(df.columns)} 列，回退到共享 X 模式")
                xy_pair = False

        if not xy_pair:
            x_values = pd.to_numeric(df.iloc[:, 0], errors="coerce").fillna(0).tolist()
            for col_idx in range(1, len(df.columns)):
                col_name = df.columns[col_idx]
                y_values = pd.to_numeric(df.iloc[:, col_idx], errors="coerce").fillna(0).tolist()
                series = chart_data.add_series(col_name)
                for x_val, y_val in zip(x_values, y_values):
                    series.add_data_point(float(x_val), float(y_val))

        chart.replace_data(chart_data)
        # 散点图数据标签：默认显示数值
        try:
            plot = chart.plots[0]
            if not plot.has_data_labels:
                plot.has_data_labels = True
            plot.data_labels.number_format = '#,##0.##'
            plot.data_labels.number_format_is_linked = False
        except Exception:
            pass
        return

    # 第一列作为类别（X轴），其余列作为系列（Y轴）
    # v2.54.27+ 多级分类 X 轴支持：
    #   根据模板图表自身是否含 multiLvlStrRef 判断是否为多级分类
    #   若模板有多级分类，按模板层级数从 DataFrame 前 N 列文本维度读取层级数据
    #   最深层（子分类）用最后一列文本维度，前面列作为父分类层级
    template_level_count = _get_template_multi_level_count(chart)

    categories = df.iloc[:, 0].astype(str).tolist()
    series_data = {}
    pct_flags = {}  # 记录每个系列是否为百分比列
    # 优先使用 load_pivot_results 读取的单元格数字格式元信息（精确）
    # 回退到列名关键词 + 值域检测（兜底）
    attr_pct_cols = df.attrs.get("pct_columns") if hasattr(df, "attrs") else None
    # 多级分类时，数据列从第 template_level_count 列开始（跳过前面的文本维度列）
    data_start_col = template_level_count if template_level_count > 0 else 1
    for col_idx in range(data_start_col, len(df.columns)):
        col_name = df.columns[col_idx]
        values = pd.to_numeric(df.iloc[:, col_idx], errors="coerce").fillna(0).tolist()
        series_data[col_name] = values
        if attr_pct_cols is not None:
            # 有元信息时以元信息为准
            pct_flags[col_name] = col_name in attr_pct_cols
        else:
            # 无元信息时回退到关键词 + 值域检测
            pct_flags[col_name] = _is_pct_column(col_name, values)

    chart_data = CategoryChartData()

    if transpose:
        # 行列转置：列名作 X 轴类别，行维度值作系列名
        categories_t = [str(c) for c in df.columns[data_start_col:]]
        chart_data.categories = categories_t
        pct_flags_t = {}
        for row_idx in range(len(df)):
            series_name = str(df.iloc[row_idx, 0])
            values = pd.to_numeric(df.iloc[row_idx, data_start_col:], errors="coerce").fillna(0).tolist()
            chart_data.add_series(series_name, values)
            pct_flags_t[series_name] = False
        pct_flags = pct_flags_t
        actual_series_count = len(df)
    else:
        # 多级分类时，子分类（最深层）用最后一列文本维度
        if template_level_count > 0:
            categories = df.iloc[:, template_level_count - 1].astype(str).tolist()
        chart_data.categories = categories
        for name, values in series_data.items():
            chart_data.add_series(name, values)
        actual_series_count = len(series_data)

    # v2.54.27+ 记录多级分类层级数据，replace_data 后恢复
    # 仅当模板本身有多级分类时才恢复，且 DataFrame 提供足够的前置文本维度列
    multi_level_data = None
    multi_level_headers = None  # v2.54.32+ 各层级列表头名（最深层在前）
    if template_level_count > 0 and not transpose and len(df.columns) >= template_level_count:
        multi_level_data = []
        multi_level_headers = []
        # PowerPoint 多级分类层级：最深层（子分类）在前，最浅层（父分类）在后
        # 例如：[产品, 地区] → level 0 是产品（子），level 1 是地区（父）
        for i in range(template_level_count - 1, -1, -1):
            level_vals = df.iloc[:, i].astype(str).tolist()
            multi_level_data.append(level_vals)
            multi_level_headers.append(str(df.columns[i]))

    chart.replace_data(chart_data)
    _trim_extra_series(chart, actual_series_count)
    _ensure_vary_colors(chart)

    # v2.54.37+ 修复组合图（多 chart 类型共存）嵌入工作簿
    # replace_data 在组合图上会把表头写成数字索引、分类列写成数字，导致 PowerPoint 报错
    if not transpose and not multi_level_data:
        _fix_embedded_workbook_for_combo(chart, categories, series_data)

    # v2.54.27+ 恢复多级分类 X 轴（replace_data 会清空 multiLvlStrRef，需要手动重建）
    # v2.54.32+ 同步重建嵌入工作簿，补齐父分类列，修复"选择数据"只显示2列的问题
    if multi_level_data:
        # 取第一个数据系列的值和名称（用于补写嵌入工作簿的数据列和表头）
        first_series_values = None
        first_series_header = None
        if series_data:
            first_name = next(iter(series_data))
            first_series_values = series_data[first_name]
            first_series_header = first_name
        _restore_multi_level_categories(chart, multi_level_data, first_series_values,
                                         multi_level_headers, first_series_header)

    # 同步数据标签格式：百分比列显示 0.0%
    try:
        plot = chart.plots[0]
        if plot.has_data_labels:
            data_labels = plot.data_labels
        else:
            plot.has_data_labels = True
            data_labels = plot.data_labels

        # 所有系列都是百分比 → 用 0.0% 格式
        # 多数系列是百分比（超过一半）→ 也用 0.0% 格式
        # 否则用普通数值格式
        all_pct = all(pct_flags.values()) if pct_flags else False
        pct_count = sum(1 for v in pct_flags.values() if v)
        total_count = len(pct_flags)
        majority_pct = total_count > 0 and pct_count * 2 > total_count

        if all_pct or majority_pct:
            data_labels.number_format = '0.0%'
            data_labels.number_format_is_linked = False
        else:
            data_labels.number_format = '#,##0.##'
            data_labels.number_format_is_linked = False
    except Exception:
        pass


def _resolve_image_path(expr: str, pivot_data: Dict[str, pd.DataFrame],
                        default_block: Optional[str], template_dir: str,
                        output_dir: Optional[str] = None,
                        image_dir: Optional[str] = None) -> Optional[str]:
    """解析 {{图片:...}} 表达式，返回图片文件绝对路径或 None

    路径解析顺序：
    1. 绝对路径 → 直接用
    2. @output/... 前缀 → 拼接到输出目录（解决带时间戳的动态输出目录）
    3. 其他相对路径 → 优先在 image_dir 查找，找不到回退到 template_dir
    4. 透视数据解析 → 取到的值再按 1/2/3 查

    支持通配符 * 和 ?（用于文件名或目录名带时间戳的模糊匹配）：
        {{图片:@output/chart_*.png}}        匹配输出目录下唯一 chart_xxxxx.png
        {{图片:@output/report_*/a.png}}     匹配输出目录下 report_xxxxx 子目录里的 a.png
    多个匹配时取第一个（按名称排序），匹配前会提示。
    """
    import glob

    expr = expr.strip()
    if not expr:
        return None

    def _try_resolve(path_str: str) -> Optional[str]:
        """对单个路径字符串尝试解析（绝对 / @output / 相对目录），支持通配符"""
        path_str = path_str.strip()
        if not path_str:
            return None

        # @output/ 前缀：相对输出目录（动态时间戳目录）
        if path_str.startswith("@output/") or path_str.startswith("@output\\"):
            rel = path_str[len("@output/"):].lstrip("/\\")
            if not output_dir:
                return None
            candidate = os.path.join(output_dir, rel)
        elif os.path.isabs(path_str):
            # 绝对路径
            candidate = path_str
        else:
            # 相对路径：优先 image_dir，回退 template_dir
            candidate = os.path.join(image_dir or template_dir, path_str)

        # 含通配符：用 glob 模糊匹配
        if "*" in candidate or "?" in candidate:
            matches = sorted(glob.glob(candidate))
            if not matches:
                # 回退到 template_dir
                if image_dir:
                    fallback = os.path.join(template_dir, path_str)
                    matches = sorted(glob.glob(fallback))
                if not matches:
                    return None
            if len(matches) > 1:
                print(f"    [信息] 通配符匹配到 {len(matches)} 个文件，使用第一个: {os.path.basename(matches[0])}")
            return os.path.abspath(matches[0])

        # 无通配符：精确匹配
        if os.path.isfile(candidate):
            return os.path.abspath(candidate)
        # 相对路径回退：image_dir 找不到时尝试 template_dir
        if image_dir and not os.path.isabs(path_str):
            fallback = os.path.join(template_dir, path_str)
            if os.path.isfile(fallback):
                return os.path.abspath(fallback)
        return None

    # 1) 直接尝试解析表达式本身作为路径
    result = _try_resolve(expr)
    if result:
        return result

    # 2) 透视数据解析（区块.列 或 区块.列.行值）
    text_value = _resolve_text_placeholder(expr, pivot_data, default_block)
    if text_value:
        result = _try_resolve(text_value)
        if result:
            return result

    return None


def _replace_pictures(slide, pivot_data: Dict[str, pd.DataFrame],
                      default_block: Optional[str],
                      template_dir: str,
                      text_image_exprs: Optional[List[str]] = None,
                      output_dir: Optional[str] = None,
                      shape_block_map: Optional[Dict[str, str]] = None,
                      image_dir: Optional[str] = None,
                      status_map: Optional[Dict[str, str]] = None) -> int:
    """替换幻灯片中的图片，返回替换次数。

    匹配优先级：
    1. 图片形状的 name 或 alternative_text 中含 {{图片:...}}
    2. 备注区声明的 形状名=图片路径（shape_block_map）
    3. 文本框中收集到的 {{图片:...}} 表达式 → 匹配同页第一张未被其他方式匹配的图片

    路径解析时，相对路径优先在 image_dir 查找，找不到再回退到 template_dir。
    status_map: 若非 None，按形状名记录替换状态（成功/失败原因），供备注区回写。
    """
    if text_image_exprs is None:
        text_image_exprs = []
    if shape_block_map is None:
        shape_block_map = {}

    replaced = 0
    matched_shape_ids = set()

    for shape in list(slide.shapes):
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue

        expr = None
        # 优先级1：形状名称/替代文字中的 {{图片:...}}
        for attr in ("name", "alternative_text"):
            val = (getattr(shape, attr, None) or "")
            m = _IMAGE_PLACEHOLDER_RE.search(val)
            if m:
                expr = m.group(1).strip()
                break

        # 优先级2：备注区声明的 形状名=图片路径
        if not expr and shape_block_map:
            try:
                mapped = shape_block_map.get(shape.name)
                if mapped:
                    expr = mapped
            except Exception:
                pass

        if not expr:
            continue

        image_path = _resolve_image_path(expr, pivot_data, default_block, template_dir, output_dir, image_dir)
        if not image_path:
            print(f"    [警告] 图片路径无效 [{expr}]: 文件不存在")
            if status_map is not None:
                status_map[shape.name] = f"失败: 图片路径无效({expr})"
            continue

        _do_replace_picture(slide, shape, image_path)
        matched_shape_ids.add(shape.shape_id)
        replaced += 1
        print(f"    [OK] 图片替换: {os.path.basename(image_path)}")
        if status_map is not None:
            status_map[shape.name] = f"成功({os.path.basename(image_path)})"

    # 文本关联模式：{{图片:...}} 在文本框中 → 替换同页第一张未被匹配的图片
    for expr in text_image_exprs:
        image_path = _resolve_image_path(expr, pivot_data, default_block, template_dir, output_dir, image_dir)
        if not image_path:
            continue
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if shape.shape_id in matched_shape_ids:
                continue
            _do_replace_picture(slide, shape, image_path)
            matched_shape_ids.add(shape.shape_id)
            replaced += 1
            print(f"    [OK] 图片替换(文本关联): {os.path.basename(image_path)}")
            break

    return replaced


def _do_replace_picture(slide, old_shape, image_path: str):
    """删除旧图片，在原位置插入新图片，保留旋转角度"""
    left, top, width, height = old_shape.left, old_shape.top, old_shape.width, old_shape.height
    rotation = old_shape.rotation
    sp = old_shape._element
    sp.getparent().remove(sp)
    new_shape = slide.shapes.add_picture(image_path, left, top, width, height)
    if rotation:
        new_shape.rotation = rotation


def _expand_table_columns(table, target_col_count: int):
    """给表格添加列到目标列数，复制最后一列的样式"""
    current = len(table.columns)
    if target_col_count <= current:
        return
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl = table._tbl
    for row_idx, tr in enumerate(tbl.findall(f"{{{ns}}}tr")):
        # 取该行最后一格做模板
        last_tc = tr.findall(f"{{{ns}}}tc")[-1]
        for _ in range(target_col_count - current):
            new_tc = last_tc.__deepcopy__(True)
            # 清空新格的文本段落
            for p in new_tc.findall(f".//{{{ns}}}p"):
                for r in p.findall(f"{{{ns}}}r"):
                    r.text = ""
            tr.append(new_tc)
    # 刷新 table 的网格列定义
    grid = tbl.find(f"{{{ns}}}tblGrid")
    if grid is not None:
        last_gc = grid.findall(f"{{{ns}}}gridCol")[-1]
        for _ in range(target_col_count - current):
            grid.append(last_gc.__deepcopy__(True))


def _expand_table_rows(table, target_row_count: int):
    """给表格添加行到目标行数（含表头），复制最后一行的样式"""
    current = len(table.rows)
    if target_row_count <= current:
        return
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl = table._tbl
    tr_list = tbl.findall(f"{{{ns}}}tr")
    last_tr = tr_list[-1]
    for _ in range(target_row_count - current):
        new_tr = last_tr.__deepcopy__(True)
        # 清空新行所有单元格的文本
        for tc in new_tr.findall(f"{{{ns}}}tc"):
            for p in tc.findall(f".//{{{ns}}}p"):
                for r in p.findall(f"{{{ns}}}r"):
                    r.text = ""
        tbl.append(new_tr)


def _clear_table_row(table, row_idx: int):
    """清空表格某行的所有单元格文本"""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tbl = table._tbl
    tr_list = tbl.findall(f"{{{ns}}}tr")
    if row_idx >= len(tr_list):
        return
    tr = tr_list[row_idx]
    for tc in tr.findall(f"{{{ns}}}tc"):
        for p in tc.findall(f".//{{{ns}}}p"):
            for r in p.findall(f"{{{ns}}}r"):
                r.text = ""


def _set_cell_text_preserve_format(cell, text: str):
    """替换单元格文本但保留模板原有格式（字体/颜色/对齐/填充）。

    python-pptx 的 cell.text setter 会清空段落重建 run，导致格式丢失。
    本函数保留第一个段落的第一个 run，只改其文本；多余 run 清空文本。
    """
    tf = cell.text_frame
    paragraphs = tf.paragraphs
    if not paragraphs:
        cell.text = text
        return
    # 用第一个段落承载文本
    first_para = paragraphs[0]
    runs = first_para.runs
    if runs:
        # 第一个 run 写入新文本
        runs[0].text = text
        # 其余 run 清空（避免拼接出旧文本）
        for r in runs[1:]:
            r.text = ""
    else:
        # 段落没有 run（空段落），直接用 add_run（会继承段落默认格式）
        first_para.add_run().text = text
    # 多余段落清空文本（保留段落本身的格式属性）
    for p in paragraphs[1:]:
        for r in p.runs:
            r.text = ""


def _replace_table_data(slide, pivot_data: Dict[str, pd.DataFrame],
                        default_block: Optional[str],
                        shape_block_map: Optional[Dict[str, str]] = None,
                        mark_missing: bool = True,
                        status_map: Optional[Dict[str, str]] = None) -> int:
    """替换幻灯片中的表格数据（整表替换），返回替换次数。
    通过表格形状的 name 或 alternative_text 中的 {{表格:区块名}} 匹配。
    形状名称/替代文字无占位符时，回退到 shape_block_map（备注区声明 形状名=区块名）。
    自动扩展表格行列以容纳完整数据，数据少于模板时清空多余行。
    status_map: 若非 None，按形状名记录替换状态（成功/失败原因），供备注区回写。
    """
    replace_count = 0
    for shape in list(slide.shapes):
        if not shape.has_table:
            continue

        target_expr = None
        for attr in ("name", "alternative_text"):
            val = (getattr(shape, attr, None) or "")
            m = _TABLE_PLACEHOLDER_RE.search(val)
            if m:
                target_expr = m.group(1).strip()
                break

        # 形状名在备注区声明了映射（方案C：备注区写 "形状名=区块表达式"）
        if not target_expr and shape_block_map:
            try:
                mapped = shape_block_map.get(shape.name)
                if mapped:
                    target_expr = mapped
            except Exception:
                pass

        if not target_expr:
            continue

        # 解析 "区块名|列1,列2" 语法
        target_block, cols = _parse_block_and_cols(target_expr)
        df = _lookup_block(pivot_data, target_block)
        if df is None:
            print(f"    [警告] 表格数据区块 '{target_block}' 未在透视结果中找到")
            if status_map is not None:
                status_map[shape.name] = f"失败: 区块'{target_block}'未找到"
            if mark_missing:
                try:
                    _mark_missing_cell(shape.table.cell(0, 0), target_block)
                except Exception:
                    pass
            continue
        if df.empty:
            if status_map is not None:
                status_map[shape.name] = f"失败: 区块'{target_block}'数据为空"
            if mark_missing:
                try:
                    _mark_missing_cell(shape.table.cell(0, 0), target_block)
                except Exception:
                    pass
            continue
        # 按指定列筛选
        df = _filter_df_columns(df, cols)

        table = shape.table
        headers = list(df.columns)
        data_rows = df.values.tolist()

        need_rows = 1 + len(data_rows)  # 表头 + 数据行
        need_cols = len(headers)

        try:
            # 自动扩展
            if need_cols > len(table.columns):
                _expand_table_columns(table, need_cols)
            if need_rows > len(table.rows):
                _expand_table_rows(table, need_rows)

            # 填表头（第一行）——保留模板单元格格式（字体/颜色/对齐/填充）
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                _set_cell_text_preserve_format(cell, str(header))

            # 填数据行
            for row_idx, row_data in enumerate(data_rows):
                table_row = row_idx + 1
                for col_idx, value in enumerate(row_data):
                    cell = table.cell(table_row, col_idx)
                    _set_cell_text_preserve_format(cell, _parse_value(value))

            # 数据少于模板行数时，清空多余行
            for extra_row in range(need_rows, len(table.rows)):
                _clear_table_row(table, extra_row)

            replace_count += 1
            cols_info = f", 仅列: {cols}" if cols else ""
            print(f"    [OK] 表格数据替换: {target_block} ({len(headers)}列 x {len(data_rows)}行{cols_info})")
            if status_map is not None:
                status_map[shape.name] = f"成功({len(headers)}列x{len(data_rows)}行)"
            # 清除表格形状名称/替代文字中的占位符
            try:
                if shape.name and _TABLE_PLACEHOLDER_RE.search(shape.name):
                    shape.name = _TABLE_PLACEHOLDER_RE.sub(target_block, shape.name)
            except Exception:
                pass
        except Exception as e:
            print(f"    [警告] 表格数据替换失败 [{target_block}]: {e}")
            if status_map is not None:
                status_map[shape.name] = f"失败: {e}"
    return replace_count


def fill_template(template_path: str, pivot_data_path: str, output_path: str, image_dir: Optional[str] = None, mark_missing: bool = True) -> Dict:
    """填充 PPT 模板

    Args:
        template_path: PPT 模板文件路径
        pivot_data_path: 透视结果 xlsx 路径
        output_path: 输出 PPT 路径
        image_dir: 图片搜索目录（相对路径图片在此目录查找，为 None 时回退到模板所在目录）
        mark_missing: 未替换的占位符是否以黄底[缺失:...]标注（默认 True）

    Returns:
        dict: 替换统计 {slides, text_replacements, chart_replacements, picture_replacements, table_replacements}
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板文件不存在: {template_path}")

    template_dir = os.path.dirname(os.path.abspath(template_path))
    output_dir = os.path.dirname(os.path.abspath(output_path))

    print(f"[模板/1] 加载透视结果: {os.path.basename(pivot_data_path)}")
    pivot_data = load_pivot_results(pivot_data_path)
    print(f"    → 共 {len(pivot_data)} 个数据区块: {', '.join(pivot_data.keys())}")

    print(f"[模板/2] 加载模板: {os.path.basename(template_path)}")
    prs = Presentation(template_path)
    print(f"    → 共 {len(prs.slides)} 页")

    total_text = 0
    total_chart = 0
    total_picture = 0
    total_table = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        default_block = None
        alias_map: Dict[str, str] = {}
        shape_block_map: Dict[str, str] = {}
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                config = _parse_slide_notes(notes_text)
                default_block = config.get("区块")
                known_prefixes = ("区块", "数据源", "别名.")
                for key, value in config.items():
                    if key.startswith("别名."):
                        alias_name = key[3:].strip()
                        if alias_name:
                            alias_map[alias_name] = value
                    elif not any(key.startswith(p) for p in known_prefixes):
                        shape_block_map[key] = value
        except Exception:
            pass

        image_collector: List[str] = []
        text_count = 0
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_count += _replace_in_text_frame(shape.text_frame, pivot_data, default_block, image_collector, alias_map, mark_missing=mark_missing)

        # v2.54.25+ 收集本页图表/图片/表格替换状态，用于回写备注区
        page_status_map: Dict[str, str] = {}
        chart_count = _replace_chart_data(slide, pivot_data, default_block, alias_map, shape_block_map, mark_missing=mark_missing, status_map=page_status_map)

        picture_count = _replace_pictures(slide, pivot_data, default_block, template_dir, image_collector, output_dir, shape_block_map, image_dir, status_map=page_status_map)

        table_count = _replace_table_data(slide, pivot_data, default_block, shape_block_map, mark_missing=mark_missing, status_map=page_status_map)

        # v2.54.25+ 将替换状态回写到备注区对应行末尾
        _update_notes_with_status(slide, page_status_map)

        total_text += text_count
        total_chart += chart_count
        total_picture += picture_count
        total_table += table_count
        print(f"    [页{slide_idx}] 文本替换: {text_count}, 图表替换: {chart_count}, 图片替换: {picture_count}, 表格替换: {table_count}")

    print(f"[模板/3] 保存: {output_path}")
    prs.save(output_path)

    stats = {
        "slides": len(prs.slides),
        "text_replacements": total_text,
        "chart_replacements": total_chart,
        "picture_replacements": total_picture,
        "table_replacements": total_table,
    }
    print(f"\n[OK] 模板填充完成！共 {stats['slides']} 页, 文本替换 {total_text} 处, 图表替换 {total_chart} 处, 图片替换 {total_picture} 处, 表格替换 {total_table} 处")
    return stats
