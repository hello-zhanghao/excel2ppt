import os
import sys
import json
import shutil
import tempfile
import zipfile
import threading
import time
import io
import queue
from html import escape as html_escape
import re
from contextlib import redirect_stdout
from datetime import datetime
from flask import Flask, request, jsonify, send_file, render_template_string, Response

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from main import _run_pivot_mode, _run_ppt_mode, _run_html_from_pivot, __VERSION__

app = Flask(__name__, static_folder="static", static_url_path="/static")
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024

TEMP_ROOT = os.path.join(tempfile.gettempdir(), "excel2ppt_web")
if not os.path.exists(TEMP_ROOT):
    os.makedirs(TEMP_ROOT)

sessions = {}
sessions_lock = threading.Lock()
analysis_lock = threading.Lock()
_SID_RE = re.compile(r"^[A-Za-z0-9_-]{16,128}$")


def _get_sid(raw_sid):
    """Reject malformed identifiers instead of silently sharing a default job."""
    sid = str(raw_sid or "")
    if not _SID_RE.fullmatch(sid):
        return None
    return sid


def _get_or_create_session(sid):
    with sessions_lock:
        if sid not in sessions:
            sessions[sid] = {
                "work_dir": tempfile.mkdtemp(prefix="job_", dir=TEMP_ROOT),
                "files": {},
                "log_queue": queue.Queue(),
                "status": "idle",
                "output": {},
            }
        return sessions[sid]


def _push_log(sid, message):
    s = _get_or_create_session(sid)
    s["log_queue"].put(f"[{datetime.now().strftime('%H:%M:%S')}] {message}")


# ==================== 页面 ====================

@app.route("/")
def index():
    html_path = os.path.join(app.static_folder, "web", "index.html")
    if os.path.exists(html_path):
        with open(html_path, "r", encoding="utf-8") as f:
            return render_template_string(f.read(), version=__VERSION__)
    return f"<h1>Excel 统一分析工具 v{__VERSION__}</h1><p>Web 界面文件不存在</p>"


# ==================== API ====================

@app.route("/api/upload", methods=["POST"])
def api_upload():
    sid = _get_sid(request.form.get("sid"))
    if not sid:
        return jsonify({"ok": False, "error": "无效会话，请刷新页面后重试"}), 400
    s = _get_or_create_session(sid)
    uploaded = []

    for key in request.files:
        for file in request.files.getlist(key):
            if file.filename == "":
                continue
            safe_name = os.path.basename(file.filename)
            # 过滤掉无关文件（只保留 xlsx）
            if not safe_name.endswith(".xlsx"):
                continue
            # 跳过 Excel 临时文件
            if safe_name.startswith("~$"):
                continue
            # 同名冲突处理：用上级目录名作前缀
            original_base = safe_name
            if safe_name in s["files"]:
                parent = os.path.basename(os.path.dirname(file.filename.replace("\\", "/")))
                if parent:
                    safe_name = f"{parent}_{original_base}"
            dst = os.path.join(s["work_dir"], safe_name)
            file.save(dst)
            s["files"][safe_name] = dst
            uploaded.append(safe_name)

    _push_log(sid, f"上传完成: {', '.join(uploaded)}")

    # 检测配置和数据文件
    configs = [f for f in uploaded if "配置" in f or "config" in f.lower()]
    data_files = [f for f in uploaded if f not in configs]

    return jsonify({
        "ok": True,
        "work_dir": s["work_dir"],
        "configs": configs,
        "data_files": data_files,
        "uploaded": uploaded,
    })


@app.route("/api/analyze", methods=["POST"])
def api_analyze():
    payload = request.get_json(silent=True) or {}
    sid = _get_sid(payload.get("sid"))
    if not sid:
        return jsonify({"ok": False, "error": "无效会话，请刷新页面后重试"}), 400
    s = _get_or_create_session(sid)

    if s["status"] == "running":
        return jsonify({"ok": False, "error": "已有任务在运行"})

    config_name = payload.get("config", "")
    config_path = s["files"].get(config_name) or _find_config(s["work_dir"])
    if not config_path or not os.path.exists(config_path):
        return jsonify({"ok": False, "error": "未找到配置文件，请先上传"})

    s["status"] = "running"
    threading.Thread(target=_run_analysis_web, args=(sid, config_path), daemon=True).start()
    return jsonify({"ok": True, "message": "分析已启动"})


def _find_config(work_dir):
    """自动查找工作目录下的配置文件"""
    candidates = []
    for f in os.listdir(work_dir):
        if f.endswith(".xlsx") and not f.startswith("~$"):
            fp = os.path.join(work_dir, f)
            candidates.append((fp, f))
    # 优先选文件名含"配置"的
    for fp, name in candidates:
        if "配置" in name or "config" in name.lower():
            return fp
    return candidates[0][0] if candidates else None


def _detect_mode_web(config_path):
    """检测配置类型（Web 版本，内联避免 pyc 缓存问题）"""
    import openpyxl
    wb = openpyxl.load_workbook(config_path, read_only=True)
    ppt_keywords = {"页码", "页面类型", "页面标题", "图表类型"}
    pivot_keywords = {"数据源", "行维度", "列维度", "值字段", "聚合方式"}
    ppt_found = False
    pivot_found = False
    for name in wb.sheetnames:
        ws = wb[name]
        all_texts = set()
        for row in ws.iter_rows(min_row=1, max_row=5):
            for cell in row:
                if cell.value is not None:
                    all_texts.add(str(cell.value).strip())
        if len(ppt_keywords & all_texts) >= 2:
            ppt_found = True
        if len(pivot_keywords & all_texts) >= 2:
            pivot_found = True
    wb.close()
    if ppt_found and pivot_found:
        return "all"
    elif ppt_found:
        return "ppt"
    elif pivot_found:
        return "pivot"
    return "unknown"


def _run_analysis_web(sid, config_path):
    s = _get_or_create_session(sid)

    class WebLogRedirector:
        def write(self, msg):
            if msg and msg.strip():
                s["log_queue"].put(msg.rstrip())
        def flush(self):
            pass

    # The legacy generators use print(). stdout is process-global, so jobs are
    # serialised while it is redirected; without this, two users receive each
    # other's logs and an exception can leave stdout redirected forever.
    with analysis_lock, redirect_stdout(WebLogRedirector()):
        try:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            base = os.path.basename(config_path).rsplit(".", 1)[0]
            detected = _detect_mode_web(config_path)
            _push_log(sid, f"配置类型: {detected}")

            if detected in ("pivot", "all"):
                try:
                    pivot_out = os.path.join(s["work_dir"], f"{base}_分析_{ts}.xlsx")
                    _pivot_ret = _run_pivot_mode(config_path, pivot_out)
                    _pivot_tasks = _pivot_ret[1] if isinstance(_pivot_ret, tuple) else None
                    s["output"]["pivot"] = pivot_out
                    html_out = os.path.join(s["work_dir"], f"{base}_报告_{ts}.html")
                    _run_html_from_pivot(pivot_out, html_out, tasks=_pivot_tasks)
                    s["output"]["html"] = html_out
                except SystemExit:
                    pass
                except Exception as e:
                    _push_log(sid, f"透视分析失败: {e}")

            if detected in ("ppt", "all"):
                try:
                    ppt_out = os.path.join(s["work_dir"], f"{base}_报告_{ts}.pptx")
                    _run_ppt_mode(config_path, ppt_out, pivot_data_file=s["output"].get("pivot"))
                    s["output"]["ppt"] = ppt_out
                except SystemExit:
                    _push_log(sid, "⚠ PPT 生成未完成（数据可能不完整）")
                except Exception as e:
                    _push_log(sid, f"PPT 生成失败: {e}")

            _push_log(sid, "✅ 分析完成！" if s["output"] else "⚠ 无输出文件")
            s["status"] = "done"
        except Exception as e:
            _push_log(sid, f"❌ 分析失败: {e}")
            s["status"] = "error"


@app.route("/api/logs")
def api_logs():
    sid = _get_sid(request.args.get("sid"))
    if not sid:
        return "无效会话", 400
    s = _get_or_create_session(sid)

    def generate():
        q = s["log_queue"]
        while True:
            try:
                msg = q.get(timeout=0.3)
                yield f"data: {json.dumps({'msg': msg})}\n\n"
            except queue.Empty:
                yield f"data: {json.dumps({'heartbeat': True})}\n\n"
                if s["status"] in ("done", "error"):
                    yield f"data: {json.dumps({'status': s['status']})}\n\n"
                    break

    return Response(generate(), mimetype="text/event-stream")


@app.route("/api/status")
def api_status():
    sid = _get_sid(request.args.get("sid"))
    if not sid:
        return jsonify({"error": "无效会话"}), 400
    s = _get_or_create_session(sid)
    return jsonify({
        "status": s["status"],
        "output": {k: os.path.basename(v) for k, v in s["output"].items()},
    })


@app.route("/api/preview/excel")
def api_preview_excel():
    sid = _get_sid(request.args.get("sid"))
    if not sid:
        return "无效会话", 400
    s = _get_or_create_session(sid)
    file_key = request.args.get("file", "pivot")
    filepath = s["output"].get(file_key)

    if not filepath or not os.path.exists(filepath):
        return "<p>文件尚未生成</p>"

    return _excel_to_html(filepath)


def _excel_to_html(filepath):
    import openpyxl
    wb = openpyxl.load_workbook(filepath, data_only=True)

    sheets_html = []
    for sname in wb.sheetnames:
        ws = wb[sname]
        rows = list(ws.iter_rows())
        if not rows:
            continue

        max_col = max(len(row) for row in rows) if rows else 1

        html = f'<div class="sheet-name">{html_escape(sname)}</div>\n<table>'
        for ri, row in enumerate(rows):
            html += "<tr>"
            is_header = (ri == 0)
            tag = "th" if is_header else "td"
            for ci in range(max_col):
                cell = row[ci] if ci < len(row) else None
                val = cell.value if cell else ""
                if val is None:
                    val = ""
                # 百分比格式
                if isinstance(val, float):
                    if cell and cell.number_format and "%" in str(cell.number_format):
                        val = f"{val:.1%}"
                    elif abs(val) < 1 and val != 0:
                        val = round(val, 4)
                html += f"<{tag}>{html_escape(str(val))}</{tag}>"
            html += "</tr>\n"
        html += "</table>"
        sheets_html.append(html)

    wb.close()
    css = """
    <style>
    body { font-family: 'Microsoft YaHei',sans-serif; margin:0; padding:12px; background:#f7f8fa; }
    .sheet-name { font-size:14px; font-weight:bold; color:#182B49; margin:12px 0 6px; }
    table { border-collapse:collapse; width:100%; font-size:12px; background:#fff; border:1px solid #ddd; }
    th { background:#182B49; color:#fff; padding:6px 10px; text-align:left; }
    td { padding:4px 10px; border-bottom:1px solid #eee; }
    tr:hover td { background:#f0f4ff; }
    </style>
    """
    return css + "<br>".join(sheets_html)


@app.route("/api/preview/ppt/<int:page>")
def api_preview_ppt_page(page):
    sid = _get_sid(request.args.get("sid"))
    if not sid:
        return "无效会话", 400
    s = _get_or_create_session(sid)
    filepath = s["output"].get("ppt")
    if not filepath or not os.path.exists(filepath):
        return "<p>PPT 文件尚未生成</p>", 404

    from pptx import Presentation
    import io

    prs = Presentation(filepath)
    slides = list(prs.slides)
    if page < 1 or page > len(slides):
        return "<p>页码超出范围</p>", 404

    # 优先用 PowerPoint COM 导出（Windows 原生，完美渲染）
    buf = _export_slide_via_com(filepath, page)
    if buf is None:
        # COM 不可用时回退到 Pillow 绘制
        img = _render_slide_pillow(slides[page - 1], prs)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
    return send_file(buf, mimetype="image/png")


def _export_slide_via_com(pptx_path, page):
    try:
        import pythoncom
        pythoncom.CoInitialize()
    except Exception:
        pass
    try:
        import win32com.client
        import tempfile
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = False
        pres = app.Presentations.Open(pptx_path, WithWindow=False)
        tmp = tempfile.mktemp(suffix=".png")
        pres.Slides(page).Export(tmp, "PNG")
        pres.Close()
        app.Quit()
        with open(tmp, "rb") as f:
            data = f.read()
        try:
            os.unlink(tmp)
        except Exception:
            pass
        buf = io.BytesIO(data)
        buf.seek(0)
        return buf
    except Exception:
        return None
    finally:
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def _render_slide_pillow(slide, prs):
    from PIL import Image, ImageDraw, ImageFont
    from pptx.oxml.ns import qn
    from lxml import etree

    width = int(prs.slide_width / 12700)
    height = int(prs.slide_height / 12700)

    # 读取幻灯片背景色（从 XML）
    bg_color = "#FFFFFF"
    try:
        cSld = slide._element.find(qn("p:cSld"))
        if cSld is not None:
            bg_el = cSld.find(qn("p:bg"))
            if bg_el is not None:
                srgb = bg_el.find(".//" + qn("a:srgbClr"))
                if srgb is not None:
                    bg_color = f"#{srgb.get('val')}"
    except Exception:
        pass

    img = Image.new("RGB", (width, height), bg_color)
    draw = ImageDraw.Draw(img)

    try:
        font_large = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 22)
        font_mid = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 16)
        font_small = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 12)
        font_tiny = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 9)
    except Exception:
        font_large = font_mid = font_small = font_tiny = ImageFont.load_default()

    for shape in slide.shapes:
        l = int(shape.left / 12700)
        t = int(shape.top / 12700)
        w = int(shape.width / 12700)
        h = int(shape.height / 12700)

        # 绘制形状背景
        try:
            if hasattr(shape, "fill"):
                ft = shape.fill.type
                if ft is not None and ft != 0:
                    try:
                        rgb = str(shape.fill.fore_color.rgb)
                        if len(rgb) == 6:
                            r, g, b = int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16)
                            draw.rectangle([l, t, l + w, t + h], fill=(r, g, b))
                    except Exception:
                        pass
        except Exception:
            pass

        # 绘制文本
        if shape.has_text_frame:
            tf = shape.text_frame
            for pi, para in enumerate(tf.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                try:
                    sz = para.font.size
                    color = para.font.color
                    is_bold = para.font.bold

                    # 选择字体
                    if sz:
                        pt = int(sz / 12700)
                        if pt >= 20:
                            ft = font_large
                        elif pt >= 14:
                            ft = font_mid
                        elif pt >= 10:
                            ft = font_small
                        else:
                            ft = font_tiny
                    else:
                        ft = font_small

                    # 字体颜色
                    try:
                        if color and color.rgb:
                            col = str(color.rgb)
                            if len(col) == 6:
                                txt_color = f"#{col}"
                            else:
                                txt_color = "#333333"
                        else:
                            txt_color = "#333333"
                    except Exception:
                        txt_color = "#333333"

                    y_offset = t + pi * int(ft.size * 1.35)
                    draw.text((l + 6, y_offset + 4), text, fill=txt_color, font=ft)
                except Exception:
                    draw.text((l + 6, t + pi * 18 + 4), text, fill="#333333")

    return img


@app.route("/api/preview/ppt/info")
def api_preview_ppt_info():
    sid = _get_sid(request.args.get("sid"))
    if not sid:
        return jsonify({"error": "无效会话"}), 400
    s = _get_or_create_session(sid)
    filepath = s["output"].get("ppt")
    if not filepath or not os.path.exists(filepath):
        return jsonify({"pages": 0})

    from pptx import Presentation
    prs = Presentation(filepath)
    return jsonify({"pages": len(list(prs.slides))})


@app.route("/api/download/<filetype>")
def api_download(filetype):
    if filetype not in {"pivot", "ppt"}:
        return "不支持的文件类型", 404
    sid = _get_sid(request.args.get("sid"))
    if not sid:
        return "无效会话", 400
    s = _get_or_create_session(sid)
    filepath = s["output"].get(filetype)
    if not filepath or not os.path.exists(filepath):
        return "文件不存在", 404
    return send_file(filepath, as_attachment=True, download_name=os.path.basename(filepath))


# ==================== 启动 ====================

def start_server(port=8899, open_browser=True, host="127.0.0.1"):
    if open_browser:
        import webbrowser
        threading.Timer(1.0, lambda: webbrowser.open(f"http://localhost:{port}")).start()
    print(f"\n🌐 Excel 统一分析工具 Web 模式 v{__VERSION__}")
    print(f"   访问地址: http://localhost:{port}")
    print(f"   按 Ctrl+C 停止服务\n")
    app.run(host=host, port=port, debug=False, threaded=True)


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("port", nargs="?", type=int, default=8899)
    parser.add_argument("--no-browser", action="store_true", help="不自动打开浏览器")
    parser.add_argument("--host", default="127.0.0.1", help="监听地址，默认仅本机访问")
    args = parser.parse_args()
    start_server(args.port, open_browser=not args.no_browser, host=args.host)
